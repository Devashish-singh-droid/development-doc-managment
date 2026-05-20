from xml.parsers.expat import model
from fastapi.templating import Jinja2Templates
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
import os, shutil
from pathlib import Path
from fastapi import Request, APIRouter, Query, HTTPException, UploadFile, File, BackgroundTasks, Form
from fastapi.responses import FileResponse, RedirectResponse, HTMLResponse
from fastapi.responses import FileResponse
from fastapi.responses import StreamingResponse
import numpy as np
from utils.gem import gem_processing
from utils.logger import get_logger
logger = get_logger('api_routes')
from config import settings
from routers.auth import get_microsoft_access_token_for_session, get_session_data

# NEW: RBAC imports
from services.rbac import (
    require_roles, 
    check_permission, 
    can_create_role,
    validate_role_string,
    get_session_role,
    filter_qa_response_by_role,
    UserRole,
)

from services.models import (
    ActivityLogResponse,
    BulkDeleteDocumentsRequest,
    BulkDeleteResponse,
    BulkDeleteEmployeeProfilesRequest,
    BulkDeleteEmployeeProfilesResponse,
    CleanupTemporaryDocumentsResponse,
    DeleteDocumentResponse,
    DocumentDetailResponse,
    DocumentsListResponse,
    EmployeeDocumentProfileResolutionResponse,
    EmployeeProfileDeleteResponse,
    EmployeeProfileRequest,
    EmployeeProfileSaveResponse,
    EmployeeProfilesListResponse,
    HealthResponse,
    ImportFromLinkRequest,
    ImportFromLinkResponse,
    MultiUploadQueuedResponse,
    ProcessingStatusResponse,
    QARequest,
    QAResponse,
    QAHistoryDetailResponse,
    QAHistoryListResponse,
    QAPresentationRequest,
    RefreshManifestRequest,
    ResolveEmployeeProfileRequest,
    SimpleStatusResponse,
    SystemSettingsResponse,
    StatisticsResponse,
    UpdateSystemSettingsRequest,
    UploadQueuedResponse,
)
from utils.request_context import get_client_ip
from db.database import DatabaseService, clear_qa_retrieval_cache
from services.employee_utils import (
    normalize_choice,
    resolve_retention_hours,
    build_expiry_at,
    extract_employee_identity,
    build_profile_suggestion,
    DEFAULT_TEMP_RETENTION_HOURS,
    CLEANUP_INTERVAL_SECONDS,
)
from typing import Optional, List, Any, Dict
from pathlib import Path
import os
import shutil
import tempfile
from io import BytesIO
from PIL import Image
import fitz
from ocr.paddle_ocr_engine import paddle_ocr_extract, get_last_ocr_diagnostics
from openpyxl import load_workbook
import pandas as pd
from services.gemini_structurer import extract_metadata_with_gemini
from services.qa_hf_answerer import answer_question, answer_question_stream, clear_cache as clear_qa_answer_cache
from services.gemini_client import (
    generate_json_text as gemini_generate_json_text,
    is_configured as gemini_is_configured,
    model_name as gemini_model_name,
)
from services.json_utils import safe_json_parse
from services.video_to_text import is_video_file, transcribe_video
from utils.pdf_text_utils import (
    pdf_has_selectable_text,
    extract_pdf_text,
    extract_pdf_link_targets,
)
from utils.excel_structurer import (
    is_excel_file,
    parse_excel_to_structured_content,
)
from utils.word_text_utils import extract_word_text
from db.mongo_db import MongoDBManager
from utils.remote_link_importer import import_documents_from_link, RemoteImportError
from datetime import datetime, timezone
import html
import mimetypes
import re
import traceback
import json
import time
import hashlib
import threading
from uuid import uuid4
from collections import OrderedDict

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    _MATPLOTLIB_AVAILABLE = True
    _MATPLOTLIB_IMPORT_ERROR = ""
except Exception as _matplotlib_exc:
    plt = None
    _MATPLOTLIB_AVAILABLE = False
    _MATPLOTLIB_IMPORT_ERROR = str(_matplotlib_exc)

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
    _PPTX_AVAILABLE = True
    _PPTX_IMPORT_ERROR = ""
except Exception as _pptx_exc:
    Presentation = None
    RGBColor = None
    MSO_AUTO_SHAPE_TYPE = None
    MSO_ANCHOR = None
    PP_ALIGN = None
    Inches = None
    Pt = None
    _PPTX_AVAILABLE = False
    _PPTX_IMPORT_ERROR = str(_pptx_exc)


def _ensure_pptx_available() -> bool:
    global Presentation, RGBColor, MSO_AUTO_SHAPE_TYPE, MSO_ANCHOR, PP_ALIGN, Inches, Pt
    global _PPTX_AVAILABLE, _PPTX_IMPORT_ERROR
    if _PPTX_AVAILABLE:
        return True
    try:
        from pptx import Presentation as _Presentation
        from pptx.dml.color import RGBColor as _RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as _MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import MSO_ANCHOR as _MSO_ANCHOR, PP_ALIGN as _PP_ALIGN
        from pptx.util import Inches as _Inches, Pt as _Pt

        Presentation = _Presentation
        RGBColor = _RGBColor
        MSO_AUTO_SHAPE_TYPE = _MSO_AUTO_SHAPE_TYPE
        MSO_ANCHOR = _MSO_ANCHOR
        PP_ALIGN = _PP_ALIGN
        Inches = _Inches
        Pt = _Pt
        _PPTX_AVAILABLE = True
        _PPTX_IMPORT_ERROR = ""
        return True
    except Exception as exc:
        _PPTX_AVAILABLE = False
        _PPTX_IMPORT_ERROR = str(exc)
        return False

# NEW: RBAC Helper functions
def _get_session_email(request: Request) -> str:
    """Extract user email from session"""
    session_id = request.cookies.get("session_id")
    session_data = get_session_data(session_id)
    return session_data.get("email") if session_data else "unknown"

def _get_user_role(request: Request) -> Optional[UserRole]:
    """Extract user role from session"""
    session_id = request.cookies.get("session_id")
    session_data = get_session_data(session_id)
    if not session_data:
        return None
    return get_session_role(session_data)

def _get_session_data_safe(request: Request) -> Optional[dict]:
    """Safely get session data from request"""
    session_id = request.cookies.get("session_id")
    return get_session_data(session_id) if session_id else None


def _get_request_session(request: Request) -> Optional[dict]:
    session_data = getattr(request.state, "session_data", None)
    return session_data if isinstance(session_data, dict) else _get_session_data_safe(request)


def _get_request_role(request: Request) -> Optional[UserRole]:
    session_data = _get_request_session(request)
    return get_session_role(session_data) if session_data else None


def _get_request_actor(request: Request) -> str:
    return str(
        getattr(request.state, "email", None)
        or getattr(request.state, "username", None)
        or _get_session_email(request)
        or ""
    ).strip()


def _normalize_kbac_owner(value: Optional[str]) -> Optional[str]:
    owner = str(value or "").strip().lower()
    return owner or None


def _get_request_creator(request: Request) -> Optional[str]:
    session_data = _get_request_session(request)
    if not isinstance(session_data, dict):
        return None
    return _normalize_kbac_owner(session_data.get("created_by"))


def _build_request_kbac_context(request: Request) -> dict:
    user_role = _get_request_role(request)
    actor = _normalize_kbac_owner(_get_request_actor(request))
    creator = _get_request_creator(request)

    if user_role in {UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        team_owner = None
        scope_key = "all"
    elif user_role == UserRole.MANAGER:
        team_owner = actor
        scope_key = f"team:{team_owner}" if team_owner else "global"
    elif user_role == UserRole.EMPLOYEE:
        team_owner = creator
        scope_key = f"team:{team_owner}" if team_owner else "global"
    else:
        team_owner = None
        scope_key = "global"

    return {
        "role": user_role.value if user_role else None,
        "actor": actor,
        "created_by": creator,
        "team_owner": team_owner,
        "allow_all": bool(user_role in {UserRole.ADMIN, UserRole.SUPER_ADMIN}),
        "scope_key": scope_key,
    }


def _resolve_document_kbac_scope(assignment: Optional[dict], source_info: Optional[dict]) -> tuple[str, Optional[str]]:
    assignment = assignment if isinstance(assignment, dict) else {}
    source = source_info if isinstance(source_info, dict) else {}
    uploaded_by = _normalize_kbac_owner(source.get("uploaded_by"))
    profile_owner = _normalize_kbac_owner(assignment.get("profile_created_by"))

    uploaded_role = None
    raw_uploaded_role = str(source.get("uploaded_by_role") or "").strip()
    if raw_uploaded_role:
        try:
            uploaded_role = validate_role_string(raw_uploaded_role)
        except ValueError:
            uploaded_role = None

    if uploaded_role in {UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        return "global", None
    if profile_owner:
        return "team", profile_owner
    if uploaded_role == UserRole.MANAGER and uploaded_by:
        return "team", uploaded_by
    return "global", None


def _require_permission(request: Request, permission: str, detail: Optional[str] = None) -> UserRole:
    session_data = _get_request_session(request)
    user_role = get_session_role(session_data) if session_data else None
    if not user_role or not check_permission(session_data, permission):
        raise HTTPException(status_code=403, detail=detail or "Access denied")
    return user_role


def _get_employee_profile_for_access(emp_id: Optional[str] = None, profile_uuid: Optional[str] = None) -> Optional[dict]:
    if emp_id:
        return db_service.db.get_employee_profile(emp_id=(emp_id or "").strip())
    if profile_uuid:
        return db_service.db.get_employee_profile(profile_uuid=(profile_uuid or "").strip())
    return None


def _require_employee_profile_access(
    request: Request,
    emp_id: Optional[str] = None,
    profile_uuid: Optional[str] = None,
    write: bool = False,
) -> tuple[UserRole, Optional[dict]]:
    user_role = _get_request_role(request)
    actor = _get_request_actor(request)
    if not user_role:
        raise HTTPException(status_code=403, detail="Access denied")

    profile = _get_employee_profile_for_access(emp_id=emp_id, profile_uuid=profile_uuid)

    if user_role in {UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        return user_role, profile

    if user_role != UserRole.MANAGER:
        raise HTTPException(status_code=403, detail="You are not allowed to access employee profiles.")

    if profile:
        profile_owner = str(profile.get("created_by") or "").strip().lower()
        if not profile_owner or profile_owner != actor.lower():
            action = "modify" if write else "view"
            raise HTTPException(
                status_code=403,
                detail=f"Managers can only {action} employee profiles they created.",
            )
    return user_role, profile


def _document_read_access_context(request: Request) -> Optional[dict]:
    """Return the KBAC context for reading a document or raise if the role is not allowed."""
    session_data = _get_request_session(request)
    user_role = get_session_role(session_data) if session_data else None
    if not user_role:
        raise HTTPException(status_code=403, detail="Access denied")

    if user_role in {UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        if not check_permission(session_data, "view_documents"):
            raise HTTPException(
                status_code=403,
                detail="Admins and Super Admins only can access documents.",
            )
        return None

    if user_role == UserRole.MANAGER and check_permission(session_data, "view_own_employee_documents"):
        return _build_request_kbac_context(request)

    raise HTTPException(
        status_code=403,
        detail="You are not allowed to view this document.",
    )


def _get_document_for_read(request: Request, doc_id: str) -> dict:
    access_context = _document_read_access_context(request)
    doc = db_service.get_document_by_id(doc_id, access_context=access_context)
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    return doc

def _get_allowed_extensions():
    """Get allowed file extensions from environment"""
    return set(settings.allowed_file_extensions)

_ALLOWED_EXTENSIONS = _get_allowed_extensions()
_LINKED_PDF_MAX_DEPTH = max(0, settings.get_int("PDF_LINK_MAX_DEPTH", 2))
_LINKED_PDF_MAX_DOCUMENTS = max(0, settings.get_int("PDF_LINK_MAX_DOCUMENTS", 100))
_LINKED_PDF_MAX_DOCUMENTS_PER_LINK = max(
    1,
    settings.get_int("PDF_LINK_MAX_DOCUMENTS_PER_LINK", 10),
)
_LINKED_PDF_PREVIEW_CHARS = max(
    300,
    settings.get_int("PDF_LINK_PREVIEW_CHARS", 1200),
)
_LINKED_PDF_ALLOWED_EXTENSIONS = {".pdf"}
_ORIGINAL_DOCUMENTS_DIR = Path(__file__).resolve().parent.parent / "output" / "original_documents"

router = APIRouter(prefix="/api", tags=["documents"])
db_service = DatabaseService()

_qa_response_cache = OrderedDict()
_qa_response_cache_lock = threading.Lock()
_QA_RESPONSE_CACHE_TTL = max(1, settings.qa_response_cache_ttl_seconds)
_QA_RESPONSE_CACHE_MAX = max(1, settings.qa_response_cache_max_size)
_qa_conversation_store = OrderedDict()
_qa_conversation_lock = threading.Lock()
_QA_CONVERSATION_TTL = max(300, settings.get_int("QA_CONVERSATION_TTL_SECONDS", 7200))
_QA_CONVERSATION_MAX_SESSIONS = max(10, settings.get_int("QA_CONVERSATION_MAX_SESSIONS", 200))
_QA_CONVERSATION_MAX_TURNS = max(1, settings.get_int("QA_CONVERSATION_MAX_TURNS", 6))
_qa_history_store = OrderedDict()
_qa_history_lock = threading.Lock()
_QA_HISTORY_TTL = max(3600, settings.get_int("QA_HISTORY_TTL_HOURS", 24) * 3600)
_QA_HISTORY_MAX_USERS = max(10, settings.get_int("QA_HISTORY_MAX_USERS", 500))
_QA_HISTORY_MAX_CONVERSATIONS_PER_USER = max(1, settings.get_int("QA_HISTORY_MAX_CONVERSATIONS_PER_USER", 50))
_QA_HISTORY_MAX_TURNS_PER_CONVERSATION = max(
    _QA_CONVERSATION_MAX_TURNS,
    settings.get_int("QA_HISTORY_MAX_TURNS_PER_CONVERSATION", 40),
)
_qa_graph_store = OrderedDict()
_qa_graph_lock = threading.Lock()
_QA_GRAPH_TTL = max(300, settings.get_int("QA_GRAPH_TTL_SECONDS", _QA_HISTORY_TTL))
_QA_GRAPH_MAX_SESSIONS = max(10, settings.get_int("QA_GRAPH_MAX_SESSIONS", 200))
_QA_GRAPH_MAX_ITEMS_PER_SESSION = max(2, settings.get_int("QA_GRAPH_MAX_ITEMS_PER_SESSION", 12))
_QA_GRAPH_PURGE_AFTER_QUESTIONS = 2
_cleanup_scheduler_lock = threading.Lock()
_cleanup_scheduler_started = False


def _serialize_datetime(value):
    if isinstance(value, datetime):
        return _to_utc_iso(value)
    if value is None:
        return None
    return value


def _serialize_processing_status(status_data: dict) -> dict:
    serialized = {}
    for key, value in (status_data or {}).items():
        if key == "_id":
            continue
        serialized[key] = _serialize_datetime(value)
    return serialized


def _resolve_original_document_target(doc: dict) -> tuple[str, Optional[str]]:
    file_path = str(doc.get("file_path") or "").strip()
    if file_path:
        path_obj = Path(file_path)
        if path_obj.exists() and path_obj.is_file():
            return ("file", str(path_obj))

    source = doc.get("source") if isinstance(doc.get("source"), dict) else {}
    source_url = str(source.get("source_url") or "").strip()
    if source_url:
        if re.match(r"^https?://", source_url, flags=re.IGNORECASE):
            return ("url", source_url)
        path_obj = Path(source_url)
        if path_obj.exists() and path_obj.is_file():
            return ("file", str(path_obj))

    return ("missing", None)


def _preview_source_name(doc: dict, fallback: str = "document") -> str:
    source = doc.get("source") if isinstance(doc.get("source"), dict) else {}
    source_url = str(source.get("source_url") or "").strip()
    file_name = str(doc.get("file_name") or "").strip()
    return file_name or Path(source_url).name or fallback


def _preview_primary_text(doc: dict) -> str:
    content = doc.get("content") if isinstance(doc.get("content"), dict) else {}
    main_text = str(content.get("main_document_text") or "").strip()
    full_text = str(content.get("full_text") or "").strip()
    return main_text or full_text


def _render_text_preview_block(text: str) -> str:
    safe_text = html.escape(text or "Preview not available.")
    return (
        "<div class='preview-text-shell'>"
        f"<pre class='preview-text'>{safe_text}</pre>"
        "</div>"
    )


def _render_excel_preview_block(file_path: str, fallback_text: str) -> str:
    try:
        excel_file = pd.ExcelFile(file_path)
        sheets_html = []
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            safe_df = df.fillna("")
            table_html = safe_df.to_html(index=False, border=0, classes="preview-table", escape=True)
            sheets_html.append(
                "<section class='preview-sheet'>"
                f"<h3>{html.escape(sheet_name)}</h3>"
                f"{table_html}"
                "</section>"
            )
        if sheets_html:
            return "".join(sheets_html)
    except Exception as exc:
        logger.warning(f"Excel preview fallback used for {file_path}: {exc}")
    return _render_text_preview_block(fallback_text)


def _build_original_preview_html(doc_id: str, doc: dict, base_path: str = "") -> str:
    base_path = str(base_path or "").rstrip("/")
    target_type, target_value = _resolve_original_document_target(doc)
    source_name = _preview_source_name(doc)
    extension = Path(source_name).suffix.lower()
    source = doc.get("source") if isinstance(doc.get("source"), dict) else {}
    preview_text = _preview_primary_text(doc)
    original_url = f"{base_path}/api/documents/{doc_id}/original"

    body_html = ""
    note_html = ""

    if extension == ".pdf" and target_type in {"file", "url"}:
        body_html = (
            f"<iframe class='preview-frame' src='{html.escape(original_url)}#toolbar=1&navpanes=0'></iframe>"
        )
    elif extension in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tif", ".tiff"} and target_type in {"file", "url"}:
        body_html = (
            "<div class='preview-image-wrap'>"
            f"<img class='preview-image' src='{html.escape(original_url)}' alt='{html.escape(source_name)}' />"
            "</div>"
        )
    elif extension in {".mp4", ".avi", ".mov", ".mkv", ".flv", ".wmv", ".m4v", ".webm"} and target_type in {"file", "url"}:
        body_html = (
            "<div class='preview-image-wrap'>"
            f"<video class='preview-image' controls preload='metadata' src='{html.escape(original_url)}'>"
            "Your browser does not support inline video playback."
            "</video>"
            "</div>"
        )
        note_html = (
            "<div class='preview-note'>"
            "Video files are previewed inline when possible. Their audio transcript is also indexed for AI Buzz and document search."
            "</div>"
        )
    elif extension in {".docx", ".docm", ".doc"}:
        word_text = preview_text
        if target_type == "file" and target_value:
            try:
                word_text = extract_word_text(target_value) or word_text
            except Exception as exc:
                logger.warning(f"Word preview fallback used for {target_value}: {exc}")
        body_html = _render_text_preview_block(word_text)
        note_html = (
            "<div class='preview-note'>"
            "Word documents are shown as readable extracted text inside the page."
            "</div>"
        )
    elif extension in {".xlsx", ".xls", ".xlsm", ".xlsb"}:
        if target_type == "file" and target_value:
            body_html = _render_excel_preview_block(target_value, preview_text)
        else:
            body_html = _render_text_preview_block(preview_text)
        note_html = (
            "<div class='preview-note'>"
            "Spreadsheet preview is rendered as in-page sheet tables."
            "</div>"
        )
    elif extension in {".txt", ".csv", ".json", ".xml", ".html", ".htm"}:
        text_value = preview_text
        if target_type == "file" and target_value:
            try:
                text_value = Path(target_value).read_text(encoding="utf-8", errors="ignore") or text_value
            except Exception as exc:
                logger.warning(f"Text preview fallback used for {target_value}: {exc}")
        body_html = _render_text_preview_block(text_value)
    elif preview_text:
        body_html = _render_text_preview_block(preview_text)
        note_html = (
            "<div class='preview-note'>"
            "Exact browser rendering is not available for this file type, so a readable extracted preview is shown."
            "</div>"
        )
    elif target_type == "url":
        body_html = (
            "<div class='preview-note'>"
            "This file is available through a remote source link, but inline rendering was not possible."
            "</div>"
        )
    else:
        body_html = (
            "<div class='preview-note'>"
            "Preview not available for this document."
            "</div>"
        )

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1">
    <title>{html.escape(source_name)}</title>
    <style>
        body {{
            margin: 0;
            padding: 16px;
            font-family: "Segoe UI", Tahoma, sans-serif;
            background: #fffaf7;
            color: #2f2f2f;
        }}
        .preview-title {{
            margin: 0 0 12px 0;
            font-size: 18px;
            font-weight: 700;
            color: #9a3412;
        }}
        .preview-note {{
            margin: 0 0 12px 0;
            padding: 10px 12px;
            border: 1px solid #fdba74;
            border-radius: 8px;
            background: #fff7ed;
            color: #9a3412;
            line-height: 1.5;
            font-size: 13px;
        }}
        .preview-frame {{
            width: 100%;
            height: 88vh;
            border: none;
            border-radius: 8px;
            background: white;
        }}
        .preview-image-wrap {{
            display: flex;
            justify-content: center;
            align-items: flex-start;
            background: white;
            border: 1px solid #e5e7eb;
            border-radius: 8px;
            padding: 12px;
        }}
        .preview-image {{
            max-width: 100%;
            height: auto;
            border-radius: 6px;
        }}
        .preview-text-shell {{
            background: white;
            border: 1px solid #e5e7eb;
            border-radius: 8px;
            padding: 12px;
        }}
        .preview-text {{
            margin: 0;
            white-space: pre-wrap;
            word-break: break-word;
            font-family: Consolas, monospace;
            font-size: 13px;
            line-height: 1.6;
        }}
        .preview-sheet {{
            margin-bottom: 18px;
            background: white;
            border: 1px solid #e5e7eb;
            border-radius: 8px;
            padding: 12px;
            overflow-x: auto;
        }}
        .preview-sheet h3 {{
            margin: 0 0 10px 0;
            color: #7c2d12;
            font-size: 15px;
        }}
        .preview-table {{
            border-collapse: collapse;
            width: 100%;
            font-size: 13px;
        }}
        .preview-table th,
        .preview-table td {{
            border: 1px solid #e5e7eb;
            padding: 8px 10px;
            text-align: left;
            vertical-align: top;
        }}
        .preview-table th {{
            background: #fff7ed;
            color: #9a3412;
        }}
    </style>
</head>
<body>
    <div class="preview-title">{html.escape(source_name)}</div>
    {note_html}
    {body_html}
</body>
</html>"""


def _normalize_optional_text(value) -> str:
    return " ".join(str(value or "").strip().split())


_HIGHLIGHT_STOP_WORDS = {
    "the", "and", "for", "with", "from", "all", "show", "find", "get",
    "list", "in", "on", "of", "to", "a", "an", "is", "are", "by",
    "this", "that", "these", "those", "please", "documents", "document",
    "report", "reports", "me", "which", "number", "no", "id",
    "tell", "give", "want", "need", "may", "can", "could", "would",
    "know", "about", "kindly", "let", "name",
}


def _extract_highlight_terms(query: str) -> list[str]:
    normalized = " ".join(str(query or "").strip().split())
    if not normalized:
        return []

    terms = re.findall(r"[a-zA-Z0-9]{3,}", normalized)
    cleaned = []
    seen = set()
    for term in terms:
        key = term.lower()
        if key in _HIGHLIGHT_STOP_WORDS:
            continue
        if key in seen:
            continue
        seen.add(key)
        cleaned.append(term)

    cleaned.sort(key=len, reverse=True)
    return cleaned[:8]


def _highlight_text(text: str, terms: list[str]) -> str:
    raw = str(text or "")
    if not raw:
        return ""
    if not terms:
        return html.escape(raw)

    pattern = re.compile("|".join(re.escape(term) for term in terms), re.IGNORECASE)
    parts = []
    last_end = 0
    for match in pattern.finditer(raw):
        if match.start() < last_end:
            continue
        parts.append(html.escape(raw[last_end:match.start()]))
        parts.append("<mark>" + html.escape(raw[match.start():match.end()]) + "</mark>")
        last_end = match.end()
    parts.append(html.escape(raw[last_end:]))
    return "".join(parts)


def _resolve_storage_mode(choose: str) -> str:
    return normalize_choice(choose)


def _coerce_bool(value, default: bool = False) -> bool:
    if value is None:
        return bool(default)
    if isinstance(value, bool):
        return value
    raw = str(value).strip().lower()
    if raw in {"1", "true", "yes", "y", "on"}:
        return True
    if raw in {"0", "false", "no", "n", "off"}:
        return False
    return bool(default)


async def _persist_uploaded_file(upload_file: UploadFile) -> tuple[str, int, str]:
    suffix = Path(str(upload_file.filename or "")).suffix
    content = await upload_file.read()
    file_hash = hashlib.sha256(content).hexdigest()
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(content)
        temp_path = tmp.name
    return temp_path, len(content), file_hash


def _compute_file_sha256(path: str) -> str:
    digest = hashlib.sha256()
    with open(path, "rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


def _safe_delete_file(path: str) -> None:
    try:
        if path and os.path.exists(path):
            os.remove(path)
    except Exception as exc:
        logger.warning(f"Could not delete temp file '{path}': {exc}")


def _schedule_temp_file_cleanup(path: str, retention_hours: int) -> None:
    try:
        delay_seconds = max(0, int(retention_hours)) * 3600
    except Exception:
        delay_seconds = DEFAULT_TEMP_RETENTION_HOURS * 3600

    if delay_seconds <= 0:
        _safe_delete_file(path)
        return

    timer = threading.Timer(delay_seconds, _safe_delete_file, args=[path])
    timer.daemon = True
    timer.start()


def _qa_ppt_normalize_text(value: Any, max_length: int = 280) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    if max_length <= 0 or len(text) <= max_length:
        return text
    clipped = text[:max_length].rstrip(" ,;:-")
    boundary = max(clipped.rfind("."), clipped.rfind(";"), clipped.rfind(","), clipped.rfind(" "))
    if boundary >= max(40, int(max_length * 0.65)):
        clipped = clipped[:boundary].rstrip(" ,;:-")
    return clipped


def _qa_ppt_clean_text(value: Any, max_length: int = 280) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    if not text:
        return ""

    filler_patterns = [
        r"\bAs an AI(?: language model)?[, ]+",
        r"\bBased on (?:the )?(?:provided|available) (?:context|documents|information)[, ]+",
        r"\bAccording to (?:the )?(?:provided|available) (?:context|documents|information)[, ]+",
        r"\bIt is important to note that\s+",
        r"\bPlease note that\s+",
        r"\bI hope this helps\.?",
        r"\bLet me know if you need (?:anything else|more information)\.?",
    ]
    for pattern in filler_patterns:
        text = re.sub(pattern, "", text, flags=re.IGNORECASE)
    text = re.sub(r"\s+", " ", text).strip(" -")
    return _qa_ppt_normalize_text(text, max_length=max_length)


def _qa_ppt_chunk_text(value: Any, *, max_chars: int) -> List[str]:
    text = _qa_ppt_clean_text(value, max_length=0)
    if not text:
        return []

    pieces = [piece.strip() for piece in re.split(r"(?<=[.!?])\s+|\n+", text) if piece.strip()]
    if not pieces:
        pieces = [text]

    chunks: List[str] = []
    current = ""
    for piece in pieces:
        if len(piece) > max_chars:
            words = piece.split()
            for word in words:
                candidate = f"{current} {word}".strip()
                if current and len(candidate) > max_chars:
                    chunks.append(current)
                    current = word
                else:
                    current = candidate
            continue

        candidate = f"{current} {piece}".strip()
        if current and len(candidate) > max_chars:
            chunks.append(current)
            current = piece
        else:
            current = candidate

    if current:
        chunks.append(current)
    return chunks


def _qa_ppt_dedupe_items(items: List[str], *, max_items: int, max_length: int) -> List[str]:
    output: List[str] = []
    seen = set()
    for item in items or []:
        clean = _qa_ppt_clean_text(item, max_length=max_length)
        key = re.sub(r"[^a-z0-9]+", " ", clean.lower()).strip()
        if not clean or key in seen:
            continue
        seen.add(key)
        output.append(clean)
        if len(output) >= max_items:
            break
    return output


def _qa_ppt_content_key(value: Any) -> str:
    clean = _qa_ppt_clean_text(value, max_length=0).lower()
    clean = re.sub(r"\b(the|a|an|and|or|of|to|in|for|with|as|at|by|from)\b", " ", clean)
    return re.sub(r"[^a-z0-9]+", " ", clean).strip()


def _qa_ppt_unique_for_deck(items: List[str], seen: set[str], *, max_items: int, max_length: int) -> List[str]:
    output: List[str] = []
    for item in items or []:
        clean = _qa_ppt_clean_text(item, max_length=max_length)
        key = _qa_ppt_content_key(clean)
        if not clean or not key:
            continue
        is_repeat = key in seen or any(key in existing or existing in key for existing in seen if len(existing) > 22 and len(key) > 22)
        if is_repeat:
            continue
        seen.add(key)
        output.append(clean)
        if len(output) >= max_items:
            break
    return output


def _qa_ppt_extract_numbers(items: List[str], *, limit: int = 4) -> List[Dict[str, str]]:
    metrics: List[Dict[str, str]] = []
    seen = set()
    pattern = re.compile(
        r"(?P<prefix>₹|rs\.?|inr|usd|\$)?\s*(?P<number>\d[\d,]*(?:\.\d+)?)\s*(?P<suffix>%|percent|crore|lakh|k|m|million|billion)?",
        flags=re.IGNORECASE,
    )
    for item in items or []:
        text = _qa_ppt_clean_text(item, max_length=180)
        if not text:
            continue
        match = pattern.search(text)
        if not match:
            continue
        value = "".join(part for part in [match.group("prefix") or "", match.group("number") or "", match.group("suffix") or ""] if part)
        label = text[: match.start()].strip(" :,-") or text[match.end():].strip(" :,-") or "Metric"
        label = _qa_ppt_clean_text(label, max_length=58) or "Metric"
        key = (value.lower(), label.lower())
        if key in seen:
            continue
        seen.add(key)
        metrics.append({"value": value, "label": label})
        if len(metrics) >= limit:
            break
    return metrics


def _qa_ppt_graph_payload(payload: QAPresentationRequest) -> dict:
    graph_data = getattr(payload, "graph_data", {}) or {}
    if not isinstance(graph_data, dict):
        return {}
    normalized = _normalize_qa_graph_data(graph_data)
    if not normalized:
        return {}
    labels = [str(label or "").strip() for label in normalized.get("labels") or []]
    values = normalized.get("values") or []
    year_label_count = sum(1 for label in labels if re.fullmatch(r"(19|20)\d{2}", label))
    year_value_count = 0
    for value in values:
        try:
            numeric_value = float(value)
        except Exception:
            continue
        if 1900 <= numeric_value <= 2100 and float(numeric_value).is_integer():
            year_value_count += 1
    if labels and year_label_count == len(labels):
        return {}
    if values and year_value_count == len(values):
        return {}
    return normalized


def _qa_ppt_chart_options(payload: QAPresentationRequest, has_graph_data: bool) -> List[str]:
    return _normalize_qa_chart_options(getattr(payload, "chart_options", []), has_graph_data=has_graph_data)


def _qa_ppt_format_value(value: Any) -> str:
    try:
        number = float(value)
    except Exception:
        return _qa_ppt_clean_text(value, max_length=24)
    if abs(number) >= 10000000:
        return f"{number / 10000000:.1f}Cr"
    if abs(number) >= 100000:
        return f"{number / 100000:.1f}L"
    if abs(number) >= 1000:
        return f"{number / 1000:.1f}K"
    return f"{number:g}"


def _qa_ppt_slugify(value: str, fallback: str = "ai_buzz_presentation") -> str:
    slug = re.sub(r"[^a-z0-9]+", "_", str(value or "").strip().lower()).strip("_")
    return slug[:64] or fallback


def _qa_ppt_title_case_slug(value: str, fallback: str = "AI_Buzz_Presentation") -> str:
    parts = [part for part in re.split(r"[^A-Za-z0-9]+", str(value or "").strip()) if part]
    if not parts:
        return fallback
    normalized = "_".join(part[:1].upper() + part[1:] for part in parts[:8]).strip("_")
    return normalized or fallback


def _qa_ppt_sentence_points(answer: str, limit: int = 8) -> List[str]:
    sentences = re.split(r"(?<=[.!?])\s+|\n+", str(answer or "").strip())
    points: List[str] = []
    seen = set()
    for item in sentences:
        clean = _qa_ppt_clean_text(item, max_length=520)
        key = clean.lower()
        if not clean or key in seen:
            continue
        seen.add(key)
        points.append(clean)
        if len(points) >= limit:
            break
    return points


def _qa_ppt_collect_points(payload: QAPresentationRequest) -> List[str]:
    structured = payload.structured_answer or {}
    points: List[str] = []
    seen = set()

    def add_point(value: Any) -> None:
        clean = _qa_ppt_clean_text(value, max_length=520)
        key = clean.lower()
        if not clean or key in seen:
            return
        seen.add(key)
        points.append(clean)

    add_point(getattr(structured, "summary", ""))
    for item in getattr(structured, "highlights", []) or []:
        add_point(item)
    for section in getattr(structured, "sections", []) or []:
        heading = _qa_ppt_clean_text(getattr(section, "heading", ""), max_length=90)
        body = _qa_ppt_clean_text(getattr(section, "body", ""), max_length=720)
        if heading and body:
            add_point(f"{heading}: {body}")
        elif heading:
            add_point(heading)
        elif body:
            add_point(body)
        for bullet in getattr(section, "bullets", []) or []:
            if heading:
                add_point(f"{heading}: {bullet}")
            else:
                add_point(bullet)
    add_point(getattr(structured, "closing", ""))

    if not points:
        for item in _qa_ppt_sentence_points(payload.answer, limit=8):
            add_point(item)
    return points[:12]


def _qa_ppt_collect_sections(payload: QAPresentationRequest) -> List[Dict[str, Any]]:
    structured = payload.structured_answer or {}
    sections: List[Dict[str, Any]] = []
    for index, section in enumerate(getattr(structured, "sections", []) or [], start=1):
        title = _qa_ppt_clean_text(getattr(section, "heading", ""), max_length=90) or f"Insight {index}"
        content_items: List[str] = []
        body = _qa_ppt_clean_text(getattr(section, "body", ""), max_length=900)
        if body:
            content_items.append(body)
        for bullet in getattr(section, "bullets", []) or []:
            clean_bullet = _qa_ppt_clean_text(bullet, max_length=700)
            if clean_bullet:
                content_items.append(clean_bullet)
        if content_items:
            sections.append({"title": title, "items": content_items[:8]})

    if not sections:
        answer_chunks = _qa_ppt_chunk_text(payload.answer, max_chars=520)
        if answer_chunks:
            for chunk in answer_chunks:
                sections.append({
                    "title": f"Key Insight {len(sections) + 1}",
                    "items": [chunk],
                })
        else:
            points = _qa_ppt_collect_points(payload)
            if points:
                for index in range(0, len(points), 4):
                    chunk = points[index:index + 4]
                    sections.append({
                        "title": f"Key Insight {len(sections) + 1}",
                        "items": chunk,
                    })
    return sections


def _qa_ppt_collect_sources(payload: QAPresentationRequest) -> List[str]:
    lines: List[str] = []
    seen = set()
    cited_ids = {
        str(item.doc_id or "").strip()
        for item in (payload.citations or [])
        if str(item.doc_id or "").strip()
    }
    for source in payload.sources or []:
        source_id = str(source.doc_id or "").strip()
        if source_id and cited_ids and source_id not in cited_ids:
            continue
        document_type = _qa_ppt_clean_text(getattr(source, "document_type", "Document"), max_length=50)
        metadata = _qa_ppt_clean_text(getattr(source, "metadata", ""), max_length=140)
        snippet = _qa_ppt_clean_text(getattr(source, "snippet", ""), max_length=360)
        parts = [document_type]
        if metadata:
            parts.append(metadata)
        if snippet:
            parts.append(snippet)
        line = " | ".join(part for part in parts if part)
        key = line.lower()
        if not line or key in seen:
            continue
        seen.add(key)
        lines.append(line)
        if len(lines) >= 6:
            break
    return lines


def _qa_ppt_collect_next_steps(payload: QAPresentationRequest) -> List[str]:
    steps: List[str] = []
    seen = set()
    for item in [payload.suggestion, *(payload.follow_up_suggestions or [])]:
        clean = _qa_ppt_clean_text(item, max_length=160)
        key = clean.lower()
        if not clean or key in seen:
            continue
        seen.add(key)
        steps.append(clean)
    return steps[:4]


def _qa_ppt_collect_summary_bullets(payload: QAPresentationRequest, points: List[str]) -> List[str]:
    structured = payload.structured_answer or {}
    raw_items: List[str] = []
    for item in getattr(structured, "highlights", []) or []:
        raw_items.append(item)
    if not raw_items:
        raw_items.extend(points or [])
    return _qa_ppt_dedupe_items(raw_items, max_items=6, max_length=260)


def _qa_ppt_build_download_name(payload: QAPresentationRequest) -> str:
    structured = payload.structured_answer or {}
    base_title = (
        _qa_ppt_normalize_text(getattr(structured, "title", ""), max_length=70)
        or _qa_ppt_normalize_text(getattr(structured, "summary", ""), max_length=70)
        or "AI Buzz Presentation"
    )
    stamp = datetime.now().strftime("%Y%m%d_%H%M")
    safe_title = _qa_ppt_title_case_slug(base_title, fallback="AI_Buzz_Presentation")
    return f"{safe_title}_{stamp}.pptx"


def _qa_ppt_prepare_summary_insights(summary: str, points: List[str]) -> List[str]:
    normalized_summary = re.sub(r"\s+", " ", str(summary or "")).strip().lower()
    insights: List[str] = []
    seen = set()
    for item in points or []:
        clean = _qa_ppt_clean_text(item, max_length=110)
        key = clean.lower()
        if not clean or key in seen or key == normalized_summary:
            continue
        seen.add(key)
        insights.append(clean)
        if len(insights) >= 3:
            break
    if not insights:
        fallback = _qa_ppt_clean_text(summary, max_length=110)
        if fallback:
            insights.append(fallback)
    return insights[:3]


def _qa_ppt_split_items_by_capacity(
    items: List[str],
    *,
    max_items: int,
    max_chars: int,
) -> List[List[str]]:
    safe_items: List[str] = []
    for item in items or []:
        chunks = _qa_ppt_chunk_text(item, max_chars=max(140, max_chars))
        safe_items.extend(chunks or [_qa_ppt_clean_text(item, max_length=max_chars)])
    safe_items = [item for item in safe_items if item]
    if not safe_items:
        return []

    chunks: List[List[str]] = []
    current: List[str] = []
    current_chars = 0
    for item in safe_items:
        item_len = len(item)
        if current and (len(current) >= max_items or current_chars + item_len > max_chars):
            chunks.append(current)
            current = []
            current_chars = 0
        current.append(item)
        current_chars += item_len
    if current:
        chunks.append(current)
    return chunks


def _qa_ppt_split_sections_for_slides(sections: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    slide_sections: List[Dict[str, Any]] = []
    safe_sections = sections if isinstance(sections, list) else []
    for index, section in enumerate(safe_sections, start=1):
        raw_title = _qa_ppt_normalize_text(section.get("title", ""), max_length=72) or f"Key Insight {index}"
        raw_items = section.get("items", []) if isinstance(section, dict) else []
        chunks = _qa_ppt_split_items_by_capacity(raw_items, max_items=6, max_chars=720)
        if not chunks:
            continue
        for chunk_index, chunk in enumerate(chunks, start=1):
            title = raw_title if chunk_index == 1 else f"{raw_title} Cont."
            slide_sections.append(
                {
                    "title": title,
                    "items": chunk,
                    "is_continuation": chunk_index > 1,
                }
            )
    return slide_sections


def _qa_ppt_calculate_bullet_font_size(
    items: List[str],
    *,
    base_size: int,
    min_size: int,
) -> int:
    safe_items = [str(item or "").strip() for item in (items or []) if str(item or "").strip()]
    if not safe_items:
        return base_size
    longest = max(len(item) for item in safe_items)
    total_chars = sum(len(item) for item in safe_items)
    size = base_size
    if len(safe_items) >= 4:
        size -= 1
    if total_chars > 260:
        size -= 1
    if total_chars > 340:
        size -= 1
    if longest > 95:
        size -= 1
    if longest > 125:
        size -= 1
    return max(min_size, size)


def _qa_ppt_requires_appendix(text: str, *, threshold: int) -> bool:
    normalized = re.sub(r"\s+", " ", str(text or "")).strip()
    return len(normalized) > threshold


def _qa_ppt_build_appendix_entries(
    question: str,
    summary: str,
    answer: str,
) -> List[Dict[str, str]]:
    entries: List[Dict[str, str]] = []
    if _qa_ppt_requires_appendix(question, threshold=80):
        entries.append({"title": "Full Question", "body": _qa_ppt_normalize_text(question, max_length=340)})
    if _qa_ppt_requires_appendix(summary, threshold=170):
        entries.append({"title": "Full Summary", "body": _qa_ppt_normalize_text(summary, max_length=380)})
    elif _qa_ppt_requires_appendix(answer, threshold=220):
        entries.append({"title": "Extended Answer", "body": _qa_ppt_normalize_text(answer, max_length=420)})
    return entries[:2]


def _qa_ppt_add_appendix_slide(
    prs,
    entries: List[Dict[str, str]],
    *,
    navy: tuple[int, int, int],
    ivory: tuple[int, int, int],
    orange: tuple[int, int, int],
    slate: tuple[int, int, int],
    muted: tuple[int, int, int],
) -> None:
    if not entries:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*ivory)
    _qa_ppt_add_textbox(slide, 0.75, 0.55, 4.2, 0.4, "Appendix", font_size=24, bold=True, color=navy)
    _qa_ppt_add_textbox(slide, 0.75, 0.98, 5.4, 0.24, "Extended text moved here to keep the main slides clean.", font_size=11, color=muted)
    for idx, entry in enumerate(entries):
        top = 1.55 + (idx * 2.55)
        _qa_ppt_add_rounded_card(slide, 0.8, top, 11.7, 2.15, (255, 255, 255), (226, 232, 240))
        _qa_ppt_add_textbox(slide, 1.08, top + 0.18, 3.4, 0.22, entry.get("title", "Details"), font_size=13, bold=True, color=orange)
        _qa_ppt_add_textbox(slide, 1.08, top + 0.5, 11.0, 1.3, entry.get("body", ""), font_size=14, color=slate)
    _qa_ppt_add_watermark(slide)


def _qa_ppt_add_watermark(
    slide,
    *,
    text: str = "Created by AI Buzz",
    color: tuple[int, int, int] = (148, 163, 184),
) -> None:
    watermark = slide.shapes.add_textbox(Inches(7.9), Inches(6.72), Inches(4.75), Inches(0.38))
    frame = watermark.text_frame
    frame.word_wrap = False
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.RIGHT
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.italic = True
    run.font.name = "Aptos"
    run.font.color.rgb = RGBColor(*color)


def _qa_ppt_add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: tuple[int, int, int] = (15, 23, 42),
    align: str = "left",
    fit: bool = True,
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(2)
    text_frame.margin_right = Pt(2)
    text_frame.margin_top = Pt(1)
    text_frame.margin_bottom = Pt(1)
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.text = str(text or "")
    paragraph.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "Aptos"
    run.font.color.rgb = RGBColor(*color)
    if fit:
        try:
            text_frame.fit_text(font_family="Aptos", max_size=font_size)
        except Exception:
            pass
    return box


def _qa_ppt_add_bullet_list(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: List[str],
    *,
    font_size: int = 18,
    color: tuple[int, int, int] = (31, 41, 55),
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(8)
    frame.margin_right = Pt(8)
    frame.margin_top = Pt(4)
    frame.margin_bottom = Pt(4)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(item or "")
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(6)
        paragraph.space_before = Pt(0)
        paragraph.alignment = PP_ALIGN.LEFT
        if paragraph.runs:
            run = paragraph.runs[0]
        else:
            run = paragraph.add_run()
            run.text = str(item or "")
        run.font.size = Pt(font_size)
        run.font.name = "Aptos"
        run.font.color.rgb = RGBColor(*color)
    try:
        frame.fit_text(font_family="Aptos", max_size=font_size)
    except Exception:
        pass
    return box


def _qa_ppt_add_rounded_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: tuple[int, int, int],
    line_color: tuple[int, int, int],
) -> Any:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    shape.line.color.rgb = RGBColor(*line_color)
    shape.line.width = Pt(1)
    return shape


def _qa_ppt_apply_designed_background(
    slide,
    *,
    base: tuple[int, int, int] = (255, 247, 237),
    accent: tuple[int, int, int] = (234, 88, 12),
    navy: tuple[int, int, int] = (15, 23, 42),
    header: bool = False,
) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*base)
    top_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.16 if not header else 0.72),
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = RGBColor(*(navy if header else accent))
    top_band.line.fill.background()
    side_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.16),
        Inches(7.5),
    )
    side_band.fill.solid()
    side_band.fill.fore_color.rgb = RGBColor(*accent)
    side_band.line.fill.background()
    for left, top, size, color, transparency in [
        (11.45, 0.45, 0.86, navy, 0.82),
        (12.1, 5.95, 0.62, accent, 0.72),
        (0.55, 6.1, 0.5, navy, 0.86),
    ]:
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left),
            Inches(top),
            Inches(size),
            Inches(size),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(*color)
        dot.fill.transparency = transparency
        dot.line.fill.background()


def _qa_ppt_add_icon_badge(
    slide,
    left: float,
    top: float,
    text: str,
    *,
    fill_color: tuple[int, int, int],
    text_color: tuple[int, int, int] = (255, 255, 255),
) -> None:
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(0.42),
        Inches(0.42),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(*fill_color)
    badge.line.fill.background()
    _qa_ppt_add_textbox(slide, left + 0.08, top + 0.08, 0.26, 0.18, text, font_size=10, bold=True, color=text_color, align="center")


def _qa_ppt_add_section_header(
    slide,
    title: str,
    subtitle: str,
    *,
    navy: tuple[int, int, int],
    muted: tuple[int, int, int],
    accent: tuple[int, int, int],
) -> None:
    _qa_ppt_add_textbox(slide, 0.72, 0.42, 8.7, 0.42, title, font_size=24, bold=True, color=navy)
    if subtitle:
        _qa_ppt_add_textbox(slide, 0.74, 0.88, 8.8, 0.26, subtitle, font_size=11, color=muted)
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.74), Inches(1.22), Inches(1.1), Inches(0.05))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(*accent)
    rule.line.fill.background()


def _qa_ppt_add_metric_card(
    slide,
    left: float,
    top: float,
    width: float,
    metric: Dict[str, str],
    *,
    navy: tuple[int, int, int],
    orange: tuple[int, int, int],
    muted: tuple[int, int, int],
) -> None:
    _qa_ppt_add_rounded_card(slide, left, top, width, 1.12, (255, 255, 255), (226, 232, 240))
    _qa_ppt_add_textbox(slide, left + 0.22, top + 0.18, width - 0.44, 0.28, metric.get("value", ""), font_size=22, bold=True, color=orange)
    _qa_ppt_add_textbox(slide, left + 0.22, top + 0.62, width - 0.44, 0.28, metric.get("label", "Metric"), font_size=10, color=muted)
    _qa_ppt_add_icon_badge(slide, left + width - 0.68, top + 0.16, "+", fill_color=navy)


def _qa_ppt_add_visual_analytics_slide(
    prs,
    payload: QAPresentationRequest,
    graph_data: dict,
    metrics: List[Dict[str, str]],
    *,
    navy: tuple[int, int, int],
    orange: tuple[int, int, int],
    slate: tuple[int, int, int],
    muted: tuple[int, int, int],
    ivory: tuple[int, int, int],
) -> None:
    if not graph_data:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _qa_ppt_apply_designed_background(slide, base=(255, 251, 245), accent=orange, navy=navy)
    _qa_ppt_add_section_header(
        slide,
        "Visual Analytics",
        "Quantified signals and comparison points from the answer.",
        navy=navy,
        muted=muted,
        accent=orange,
    )

    if graph_data and _MATPLOTLIB_AVAILABLE:
        chart_options = _qa_ppt_chart_options(payload, has_graph_data=True)
        chart_type = "line" if "line" in chart_options else "bar"
        try:
            chart_bytes = _qa_render_chart_png(graph_data, chart_type)
            chart_stream = BytesIO(chart_bytes)
            slide.shapes.add_picture(chart_stream, Inches(0.76), Inches(1.55), width=Inches(7.15), height=Inches(4.45))
        except Exception as exc:
            logger.warning(f"[QA PPT] Chart could not be embedded: {exc}")
            graph_data = {}

    if graph_data:
        labels = list(graph_data.get("labels") or [])[:4]
        values = list(graph_data.get("values") or [])[:4]
        for idx, (label, value) in enumerate(zip(labels, values)):
            metric = {"value": _qa_ppt_format_value(value), "label": _qa_ppt_clean_text(label, max_length=48)}
            _qa_ppt_add_metric_card(slide, 8.35, 1.45 + (idx * 1.22), 4.15, metric, navy=navy, orange=orange, muted=muted)
    else:
        for idx, metric in enumerate(metrics[:4]):
            left = 1.0 + ((idx % 2) * 5.55)
            top = 1.62 + ((idx // 2) * 1.55)
            _qa_ppt_add_metric_card(slide, left, top, 4.95, metric, navy=navy, orange=orange, muted=muted)

    structured = payload.structured_answer or {}
    insight = _qa_ppt_clean_text(getattr(structured, "summary", "") or payload.answer, max_length=360)
    _qa_ppt_add_rounded_card(slide, 0.92, 5.58, 11.55, 0.88, (255, 255, 255), (226, 232, 240))
    _qa_ppt_add_textbox(slide, 1.18, 5.76, 11.0, 0.42, insight, font_size=11, color=slate)
    _qa_ppt_add_watermark(slide)


def _qa_ppt_add_timeline_slide(
    prs,
    sections: List[Dict[str, Any]],
    points: List[str],
    *,
    navy: tuple[int, int, int],
    orange: tuple[int, int, int],
    slate: tuple[int, int, int],
    muted: tuple[int, int, int],
    ivory: tuple[int, int, int],
) -> None:
    timeline_items = _qa_ppt_dedupe_items(
        [section.get("title", "") for section in sections if isinstance(section, dict)] or points,
        max_items=4,
        max_length=56,
    )
    if len(timeline_items) < 3:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _qa_ppt_apply_designed_background(slide, base=(248, 250, 252), accent=orange, navy=navy)
    _qa_ppt_add_section_header(
        slide,
        "Storyline And Flow",
        "A concise executive narrative from issue to implication to action.",
        navy=navy,
        muted=muted,
        accent=orange,
    )
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.25), Inches(3.1), Inches(10.7), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(203, 213, 225)
    line.line.fill.background()

    step_width = 10.4 / max(1, len(timeline_items) - 1)
    for idx, item in enumerate(timeline_items):
        left = 1.05 + (idx * step_width)
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(2.82), Inches(0.58), Inches(0.58))
        marker.fill.solid()
        marker.fill.fore_color.rgb = RGBColor(*(orange if idx == 0 else navy))
        marker.line.fill.background()
        _qa_ppt_add_textbox(slide, left + 0.14, 2.96, 0.3, 0.14, str(idx + 1), font_size=10, bold=True, color=(255, 255, 255), align="center")
        card_top = 1.62 if idx % 2 == 0 else 3.72
        fill = (255, 244, 230) if idx % 2 == 0 else (239, 246, 255)
        line_color = (253, 186, 116) if idx % 2 == 0 else (147, 197, 253)
        _qa_ppt_add_rounded_card(slide, left - 0.42, card_top, 2.25, 1.0, fill, line_color)
        _qa_ppt_add_textbox(slide, left - 0.18, card_top + 0.2, 1.75, 0.45, item, font_size=12, bold=True, color=slate, align="center")
    _qa_ppt_add_watermark(slide)


def _qa_ppt_add_text_continuation_slides(
    prs,
    title: str,
    chunks: List[str],
    *,
    navy: tuple[int, int, int],
    orange: tuple[int, int, int],
    slate: tuple[int, int, int],
    muted: tuple[int, int, int],
    ivory: tuple[int, int, int],
) -> None:
    for index, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _qa_ppt_apply_designed_background(slide, base=(255, 251, 245), accent=orange, navy=navy)
        slide_title = title if index == 1 else f"{title} Cont. {index}"
        _qa_ppt_add_section_header(
            slide,
            slide_title,
            "Additional detail is placed on separate slides so the text remains readable.",
            navy=navy,
            muted=muted,
            accent=orange,
        )
        _qa_ppt_add_rounded_card(slide, 0.9, 1.55, 11.55, 4.95, (255, 255, 255), (226, 232, 240))
        font_size = 16 if len(chunk) < 520 else 14
        _qa_ppt_add_textbox(slide, 1.22, 1.9, 10.9, 4.15, chunk, font_size=font_size, color=slate)
        _qa_ppt_add_watermark(slide)


def _qa_build_presentation_file(payload: QAPresentationRequest) -> str:
    if not _ensure_pptx_available():
        raise RuntimeError(f"python-pptx is unavailable: {_PPTX_IMPORT_ERROR}")

    question = _qa_ppt_clean_text(payload.question, max_length=140) or "AI Buzz Answer"
    full_question = re.sub(r"\s+", " ", str(payload.question or "")).strip()
    answer = _qa_ppt_clean_text(payload.answer, max_length=0)
    full_answer = re.sub(r"\s+", " ", str(payload.answer or "")).strip()
    structured = payload.structured_answer or {}
    title = _qa_ppt_clean_text(getattr(structured, "title", ""), max_length=110) or question
    summary = _qa_ppt_clean_text(getattr(structured, "summary", ""), max_length=0) or answer
    full_summary = re.sub(r"\s+", " ", str(getattr(structured, "summary", "") or "")).strip()
    points = _qa_ppt_collect_points(payload)
    sections = _qa_ppt_split_sections_for_slides(_qa_ppt_collect_sections(payload))
    sources = _qa_ppt_collect_sources(payload)
    next_steps = _qa_ppt_collect_next_steps(payload)
    summary_insights = _qa_ppt_prepare_summary_insights(summary, points)
    appendix_entries = _qa_ppt_build_appendix_entries(full_question, full_summary or summary, full_answer or answer)
    graph_data = _qa_ppt_graph_payload(payload)
    metrics = _qa_ppt_extract_numbers([summary, answer, *points, *sources], limit=4)
    deck_seen: set[str] = set()
    summary_bullets = _qa_ppt_collect_summary_bullets(payload, points)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    navy = (15, 23, 42)
    orange = (234, 88, 12)
    orange_soft = (255, 237, 213)
    slate = (51, 65, 85)
    muted = (100, 116, 139)
    ivory = (248, 250, 252)

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = RGBColor(*navy)
    for left, top, size, color in [
        (9.4, 0.5, 2.1, (30, 41, 59)),
        (10.5, 4.9, 1.5, (59, 130, 246)),
        (0.5, 5.6, 1.3, (234, 88, 12)),
    ]:
        circle = cover.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left),
            Inches(top),
            Inches(size),
            Inches(size),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*color)
        circle.fill.transparency = 0.7
        circle.line.fill.background()
    accent = cover.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(0.8),
        Inches(0.22),
        Inches(5.9),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*orange)
    accent.line.fill.background()
    _qa_ppt_add_textbox(cover, 1.2, 1.0, 9.4, 0.5, "AI Buzz Presentation", font_size=18, bold=True, color=(255, 244, 230))
    _qa_ppt_add_textbox(cover, 1.2, 1.7, 10.3, 1.5, title, font_size=30, bold=True, color=(255, 255, 255))
    _qa_ppt_add_textbox(cover, 1.2, 3.4, 10.1, 1.2, summary or answer or "Generated from the latest AI Buzz answer.", font_size=18, color=(226, 232, 240))
    _qa_ppt_add_textbox(
        cover,
        1.2,
        6.45,
        4.8,
        0.3,
        f"Generated on {datetime.now().strftime('%d %b %Y %I:%M %p')}",
        font_size=11,
        color=(203, 213, 225),
    )
    _qa_ppt_add_watermark(cover, color=(148, 163, 184))

    summary_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _qa_ppt_apply_designed_background(summary_slide, base=(255, 251, 245), accent=orange, navy=navy, header=True)
    _qa_ppt_add_textbox(summary_slide, 0.6, 0.18, 5.2, 0.34, "Executive Summary", font_size=22, bold=True, color=(255, 255, 255))
    summary_chunks = _qa_ppt_chunk_text(summary or answer or "No summary available.", max_chars=420)
    summary_text = summary_chunks[0] if summary_chunks else "No summary available."
    summary_items = _qa_ppt_unique_for_deck(summary_bullets or [summary_text], deck_seen, max_items=6, max_length=260)
    if len(summary_items) <= 3:
        for idx, item in enumerate(summary_items):
            top = 1.2 + (idx * 1.45)
            fill = (255, 244, 230) if idx % 2 == 0 else (239, 246, 255)
            line_color = (253, 186, 116) if idx % 2 == 0 else (147, 197, 253)
            _qa_ppt_add_rounded_card(summary_slide, 0.9, top, 11.5, 1.08, fill, line_color)
            _qa_ppt_add_icon_badge(summary_slide, 1.18, top + 0.32, str(idx + 1), fill_color=orange)
            _qa_ppt_add_textbox(summary_slide, 1.82, top + 0.28, 10.15, 0.42, item, font_size=15, color=navy)
    else:
        for idx, item in enumerate(summary_items[:6]):
            card_left = 0.85 + ((idx % 2) * 5.95)
            card_top = 1.05 + ((idx // 2) * 1.62)
            fill = orange_soft if idx % 2 == 0 else (239, 246, 255)
            line_color = (253, 186, 116) if idx % 2 == 0 else (147, 197, 253)
            _qa_ppt_add_rounded_card(summary_slide, card_left, card_top, 5.45, 1.25, fill, line_color)
            _qa_ppt_add_icon_badge(summary_slide, card_left + 0.22, card_top + 0.38, str(idx + 1), fill_color=orange)
            _qa_ppt_add_textbox(summary_slide, card_left + 0.82, card_top + 0.25, 4.2, 0.58, item, font_size=12, color=navy)
    _qa_ppt_add_watermark(summary_slide)

    if not summary_bullets and len(summary_chunks) > 1:
        _qa_ppt_add_text_continuation_slides(
            prs,
            "Executive Summary Details",
            summary_chunks[1:],
            navy=navy,
            orange=orange,
            slate=slate,
            muted=muted,
            ivory=ivory,
        )

    _qa_ppt_add_visual_analytics_slide(
        prs,
        payload,
        graph_data,
        metrics,
        navy=navy,
        orange=orange,
        slate=slate,
        muted=muted,
        ivory=ivory,
    )

    _qa_ppt_add_timeline_slide(
        prs,
        sections,
        points,
        navy=navy,
        orange=orange,
        slate=slate,
        muted=muted,
        ivory=ivory,
    )

    for idx, section in enumerate(sections, start=1):
        section_title = section["title"] if isinstance(section, dict) else f"Key Insight {idx}"
        section_items = section.get("items", []) if isinstance(section, dict) else []
        section_items = _qa_ppt_unique_for_deck(section_items, deck_seen, max_items=6, max_length=520)
        if not section_items:
            continue
        section_font = _qa_ppt_calculate_bullet_font_size(section_items, base_size=15, min_size=11)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _qa_ppt_apply_designed_background(slide, base=(255, 251, 245), accent=orange, navy=navy)
        section_palette = [
            ((255, 244, 230), (253, 186, 116)),
            ((239, 246, 255), (147, 197, 253)),
            ((240, 253, 250), (94, 234, 212)),
            ((245, 243, 255), (196, 181, 253)),
        ]
        primary_fill, primary_line = section_palette[(idx - 1) % len(section_palette)]
        header_band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.16),
            Inches(0),
            Inches(13.173),
            Inches(0.82),
        )
        header_band.fill.solid()
        header_band.fill.fore_color.rgb = RGBColor(*primary_fill)
        header_band.line.fill.background()

        _qa_ppt_add_textbox(slide, 0.7, 0.18, 8.8, 0.42, section_title, font_size=22, bold=True, color=navy)
        _qa_ppt_add_textbox(slide, 10.35, 0.23, 2.1, 0.22, f"Slide {idx + 2}", font_size=11, bold=True, color=muted, align="center")
        if len(section_items) <= 3:
            for item_index, item in enumerate(section_items):
                top = 1.35 + (item_index * 1.55)
                fill, line_color = section_palette[item_index % len(section_palette)]
                _qa_ppt_add_rounded_card(slide, 0.9, top, 11.55, 1.18, fill, line_color)
                _qa_ppt_add_icon_badge(slide, 1.18, top + 0.38, str(item_index + 1), fill_color=orange)
                _qa_ppt_add_textbox(slide, 1.82, top + 0.25, 10.15, 0.56, item, font_size=section_font, color=slate)
        else:
            for item_index, item in enumerate(section_items[:6]):
                left = 0.88 + ((item_index % 2) * 5.85)
                top = 1.25 + ((item_index // 2) * 1.62)
                fill, line_color = section_palette[item_index % len(section_palette)]
                _qa_ppt_add_rounded_card(slide, left, top, 5.35, 1.26, fill, line_color)
                _qa_ppt_add_icon_badge(slide, left + 0.22, top + 0.4, str(item_index + 1), fill_color=orange)
                _qa_ppt_add_textbox(slide, left + 0.82, top + 0.25, 4.08, 0.58, item, font_size=max(11, section_font - 1), color=slate)
        _qa_ppt_add_watermark(slide)

    if sources:
        source_chunks = _qa_ppt_split_items_by_capacity(sources, max_items=5, max_chars=560) or [sources[:5]]
        for chunk_index, source_chunk in enumerate(source_chunks, start=1):
            sources_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _qa_ppt_apply_designed_background(sources_slide, base=(248, 250, 252), accent=orange, navy=navy)
            source_font = _qa_ppt_calculate_bullet_font_size(source_chunk, base_size=13, min_size=10)
            title_text = "Source Highlights" if chunk_index == 1 else f"Source Highlights Cont. {chunk_index}"
            _qa_ppt_add_textbox(sources_slide, 0.7, 0.55, 6.6, 0.4, title_text, font_size=24, bold=True, color=navy)
            _qa_ppt_add_textbox(sources_slide, 0.7, 0.98, 7.8, 0.28, "Relevant references used by AI Buzz for this answer", font_size=12, color=muted)
            for item_index, item in enumerate(source_chunk):
                top = 1.5 + (item_index * 0.96)
                fill = (255, 244, 230) if item_index % 2 == 0 else (239, 246, 255)
                line_color = (253, 186, 116) if item_index % 2 == 0 else (147, 197, 253)
                _qa_ppt_add_rounded_card(sources_slide, 0.8, top, 11.8, 0.78, fill, line_color)
                _qa_ppt_add_textbox(sources_slide, 1.06, top + 0.14, 11.15, 0.46, item, font_size=source_font, color=slate)
            _qa_ppt_add_watermark(sources_slide)

    closing_slide = prs.slides.add_slide(prs.slide_layouts[6])
    closing_slide.background.fill.solid()
    closing_slide.background.fill.fore_color.rgb = RGBColor(*navy)
    ribbon = closing_slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(0.75),
        Inches(2.5),
        Inches(0.45),
    )
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = RGBColor(*orange)
    ribbon.line.fill.background()
    _qa_ppt_add_textbox(closing_slide, 0.92, 0.86, 2.02, 0.2, "Final Summary", font_size=14, bold=True, color=(255, 255, 255), align="center")
    closing_heading = "Executive takeaways and recommended decisions"
    closing_heading_font = 24 if len(closing_heading) < 44 else 22
    _qa_ppt_add_textbox(closing_slide, 0.9, 1.5, 6.9, 0.78, closing_heading, font_size=closing_heading_font, bold=True, color=(255, 255, 255))
    structured_closing = _qa_ppt_clean_text(getattr(structured, "closing", ""), max_length=260)
    closing_items = next_steps or ([structured_closing] if structured_closing else []) or [
        "Review the profile details with stakeholders.",
        "Use the summarized evidence to confirm the next decision.",
    ]
    closing_main_items = _qa_ppt_split_items_by_capacity(closing_items, max_items=3, max_chars=240)
    closing_primary_items = closing_main_items[0] if closing_main_items else closing_items[:3]
    closing_font = _qa_ppt_calculate_bullet_font_size(closing_primary_items, base_size=14, min_size=11)
    _qa_ppt_add_bullet_list(closing_slide, 1.0, 2.5, 6.35, 3.45, closing_primary_items, font_size=closing_font, color=(226, 232, 240))
    recap_card = _qa_ppt_add_rounded_card(closing_slide, 8.15, 1.7, 4.15, 3.55, (255, 244, 230), (253, 186, 116))
    recap_card.shadow.inherit = False
    _qa_ppt_add_textbox(closing_slide, 8.5, 2.05, 3.4, 0.25, "Conclusion", font_size=13, bold=True, color=orange)
    conclusion_text = _qa_ppt_clean_text(summary or question, max_length=260)
    conclusion_font = 12 if len(conclusion_text) > 130 else 14
    _qa_ppt_add_textbox(closing_slide, 8.5, 2.38, 3.45, 1.48, conclusion_text, font_size=conclusion_font, bold=True, color=navy)
    _qa_ppt_add_textbox(closing_slide, 8.5, 4.08, 3.45, 0.72, "Client-ready narrative built from the current answer, evidence, and data signals.", font_size=11, color=slate)
    _qa_ppt_add_watermark(closing_slide, color=(148, 163, 184))

    if len(closing_main_items) > 1:
        for chunk_index, chunk in enumerate(closing_main_items[1:], start=2):
            continuation_slide = prs.slides.add_slide(prs.slide_layouts[6])
            continuation_slide.background.fill.solid()
            continuation_slide.background.fill.fore_color.rgb = RGBColor(*navy)
            continuation_heading = f"Next Moves Cont. {chunk_index}"
            _qa_ppt_add_textbox(continuation_slide, 0.95, 0.95, 5.0, 0.36, continuation_heading, font_size=20, bold=True, color=(255, 255, 255))
            chunk_font = _qa_ppt_calculate_bullet_font_size(chunk, base_size=15, min_size=11)
            _qa_ppt_add_bullet_list(continuation_slide, 1.0, 1.75, 7.2, 4.8, chunk, font_size=chunk_font, color=(226, 232, 240))
            _qa_ppt_add_textbox(
                continuation_slide,
                8.55,
                1.8,
                3.1,
                2.0,
                _qa_ppt_normalize_text(summary or answer, max_length=170),
                font_size=14,
                color=(226, 232, 240),
            )
            _qa_ppt_add_watermark(continuation_slide, color=(148, 163, 184))

    _qa_ppt_add_appendix_slide(
        prs,
        appendix_entries,
        navy=navy,
        ivory=ivory,
        orange=orange,
        slate=slate,
        muted=muted,
    )

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        output_path = tmp.name
    prs.save(output_path)
    return output_path


def _sanitize_filename_for_storage(filename: str) -> str:
    raw_name = str(filename or "").strip() or "document"
    suffix = Path(raw_name).suffix
    stem = Path(raw_name).stem
    safe_stem = re.sub(r"[^A-Za-z0-9._-]+", "_", stem).strip("._-") or "document"
    safe_suffix = re.sub(r"[^A-Za-z0-9.]+", "", suffix) or ""
    return f"{safe_stem}{safe_suffix}"


def _move_file_to_original_store(source_path: str, original_filename: str) -> str:
    source = Path(str(source_path or "")).resolve()
    _ORIGINAL_DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)
    safe_name = _sanitize_filename_for_storage(original_filename)
    stored_name = f"{datetime.utcnow().strftime('%Y%m%d%H%M%S')}_{uuid4().hex[:10]}_{safe_name}"
    destination = (_ORIGINAL_DOCUMENTS_DIR / stored_name).resolve()
    shutil.move(str(source), str(destination))
    return str(destination)


def _prepare_processing_file(source_path: str, original_filename: str) -> tuple[str, Optional[str], bool]:
    resolved_source = str(Path(str(source_path or "")).resolve())
    if settings.save_original_documents:
        stored_path = _move_file_to_original_store(resolved_source, original_filename)
        return stored_path, stored_path, True
    return resolved_source, None, False


def _build_duplicate_skip_message(original_filename: str, duplicate_match: Optional[dict] = None, reason: Optional[str] = None) -> str:
    safe_name = str(original_filename or "file").strip() or "file"
    existing_name = str((duplicate_match or {}).get("file_name") or "").strip()
    doc_id = str((duplicate_match or {}).get("document_id") or "").strip()
    reason_text = str(reason or (duplicate_match or {}).get("reason") or "").strip()

    if existing_name and doc_id:
        return f"Duplicate detected for {safe_name}. Matched existing document {existing_name} ({doc_id}), so processing was skipped."
    if existing_name:
        return f"Duplicate detected for {safe_name}. Matched existing document {existing_name}, so processing was skipped."
    if reason_text == "batch_duplicate":
        return f"Duplicate detected for {safe_name} within the same upload batch, so processing was skipped."
    return f"Duplicate detected for {safe_name}, so processing was skipped."


def _register_duplicate_skip_status(
    display_filename: str,
    message: str,
    duplicate_match: Optional[dict] = None,
) -> dict:
    payload = {
        "status": "duplicate_skipped",
        "progress": 100,
        "message": message,
        "duplicate_of": duplicate_match,
        "duplicate_reason": str((duplicate_match or {}).get("reason") or "").strip() or "duplicate_detected",
    }
    if duplicate_match and duplicate_match.get("document_id"):
        payload["document_id"] = duplicate_match.get("document_id")
        payload["document_uuid"] = duplicate_match.get("document_uuid")
    processing_status[display_filename] = payload
    save_status_to_db(display_filename, payload)
    return payload


def _resolve_employee_assignment(
    db: MongoDBManager,
    requested_emp_id: str,
    requested_emp_name: str,
    requested_choice: str,
    retention_hours: int,
    structured_metadata: dict,
    owner_email: Optional[str] = None,
    owner_role: Optional[UserRole] = None,
) -> dict:
    extracted_emp_id, extracted_emp_name = extract_employee_identity(structured_metadata)
    requested_emp_id = _normalize_optional_text(requested_emp_id)
    requested_emp_name = _normalize_optional_text(requested_emp_name)
    suggested_emp_id = requested_emp_id or extracted_emp_id
    suggested_emp_name = requested_emp_name or extracted_emp_name
    normalized_choice = normalize_choice(requested_choice)
    owner_email = str(owner_email or "").strip() or None
    manager_scope = owner_role == UserRole.MANAGER and owner_email

    profile = None
    if requested_emp_id:
        blocked_existing_profile = False
        profile = db.get_employee_profile(emp_id=requested_emp_id)
        if manager_scope and profile:
            profile_owner = str(profile.get("created_by") or "").strip().lower()
            if not profile_owner or profile_owner != owner_email.lower():
                profile = None
                blocked_existing_profile = True
        if not profile and not blocked_existing_profile and (requested_emp_name or extracted_emp_name):
            profile = db.create_or_update_employee_profile(
                requested_emp_id,
                requested_emp_name or extracted_emp_name,
                extra={"created_by": owner_email} if manager_scope else {},
            )
        elif profile and (requested_emp_name or extracted_emp_name):
            db.create_or_update_employee_profile(
                profile.get("empID"),
                requested_emp_name or extracted_emp_name or profile.get("empName"),
                extra={"created_by": profile.get("created_by") or owner_email} if manager_scope else {},
            )
            profile = db.get_employee_profile(emp_id=requested_emp_id)
    # AUTO-LINK: If no explicitly requested empID but extracted_emp_name exists, try to auto-link
    elif extracted_emp_name:
        # Check if profile already exists with this name
        candidates = db.search_employee_profiles(
            extracted_emp_name,
            limit=1,
            created_by=owner_email if manager_scope else None,
        )
        if candidates:
            profile = candidates[0]
        else:
            # Create new profile with extracted name (no ID, system will generate one)
            generated_emp_id = f"EMP-{extracted_emp_name.upper().replace(' ', '-')[:20]}"
            profile = db.create_or_update_employee_profile(
                generated_emp_id,
                extracted_emp_name,
                extra={"created_by": owner_email} if manager_scope else {},
            )
            if profile:
                logger.info(f"[AUTO-LINK] Created new employee profile for: {extracted_emp_name} (ID: {generated_emp_id})")

    if profile:
        return {
            "empID": profile.get("empID"),
            "empName": profile.get("empName"),
            "employee_uuid": profile.get("uuid"),
            "profile_created_by": profile.get("created_by"),
            "storage_mode": normalized_choice,
            "isTemporary": normalized_choice == "temporary",
            "expiry_at": build_expiry_at(retention_hours) if normalized_choice == "temporary" else None,
            "assignment_status": "linked",
            "employee_action_required": False,
            "suggested_profile": build_profile_suggestion(profile.get("empID"), profile.get("empName")),
        }

    return {
        "empID": None,
        "empName": None,
        "employee_uuid": None,
        "profile_created_by": None,
        "storage_mode": normalized_choice,
        "isTemporary": normalized_choice == "temporary",
        "expiry_at": build_expiry_at(retention_hours) if normalized_choice == "temporary" else None,
        "assignment_status": "pending_profile",
        "employee_action_required": True,
        "suggested_profile": build_profile_suggestion(suggested_emp_id, suggested_emp_name),
    }


def _run_cleanup_cycle():
    cleanup_service = DatabaseService()
    try:
        result = cleanup_service.cleanup_expired_temporary_documents()
        if result.get("deleted_documents"):
            logger.info(
                f"[TEMP CLEANUP] Removed {result.get('deleted_documents', 0)} expired document(s) "
                f"and {result.get('deleted_files', 0)} file(s)"
            )
    except Exception as exc:
        logger.warning(f"[TEMP CLEANUP] Cleanup cycle failed: {exc}")
    finally:
        cleanup_service.close()


def start_cleanup_scheduler():
    global _cleanup_scheduler_started
    with _cleanup_scheduler_lock:
        if _cleanup_scheduler_started:
            return

        def _cleanup_loop():
            while True:
                _run_cleanup_cycle()
                time.sleep(CLEANUP_INTERVAL_SECONDS)

        worker = threading.Thread(target=_cleanup_loop, name="temp-document-cleanup", daemon=True)
        worker.start()
        _cleanup_scheduler_started = True


def _qa_cache_key(question: str, limit: int, chunk_limit: int, use_chunks: bool, scope_key: str = "all") -> str:
    normalized = " ".join(str(question or "").strip().lower().split())
    payload = (
        f"{normalized}|{int(limit or 0)}|chunks:{1 if use_chunks else 0}|"
        f"{int(chunk_limit or 0)}|scope:{scope_key or 'all'}"
    )
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _qa_cache_get(key: str):
    now = time.time()
    with _qa_response_cache_lock:
        item = _qa_response_cache.get(key)
        if not item:
            return None
        expires_at, value = item
        if expires_at <= now:
            _qa_response_cache.pop(key, None)
            return None
        _qa_response_cache.move_to_end(key)
        return value


def _qa_cache_set(key: str, value):
    expires_at = time.time() + _QA_RESPONSE_CACHE_TTL
    with _qa_response_cache_lock:
        _qa_response_cache[key] = (expires_at, value)
        _qa_response_cache.move_to_end(key)
        while len(_qa_response_cache) > _QA_RESPONSE_CACHE_MAX:
            _qa_response_cache.popitem(last=False)


def clear_qa_response_cache():
    with _qa_response_cache_lock:
        _qa_response_cache.clear()
    logger.info("[QA API] Response cache cleared")


def _get_request_session_id(request: Request) -> str:
    return str(request.cookies.get("session_id") or "").strip()


_QA_SOFT_NOT_FOUND_ANSWER = "I wasn't able to find matching information in my knowledge sources for that request."
_QA_INCOMPLETE_INPUT_ANSWER = "It looks like you were still typing. Please complete the sentence so I can help properly."


def _qa_incomplete_input_payload(conversation_id: str) -> dict:
    return {
        "answer": _QA_INCOMPLETE_INPUT_ANSWER,
        "suggestion": "",
        "conversation_id": str(conversation_id or "").strip(),
        "turn_id": "",
        "revision_count": 1,
        "structured_answer": {
            "style": "paragraph",
            "title": "",
            "summary": _QA_INCOMPLETE_INPUT_ANSWER,
            "highlights": [],
            "sections": [],
            "closing": "",
        },
        "citations": [],
        "sources": [],
        "chunk_sources": [],
        "graph_id": "",
        "chart_options": [],
        "follow_up_suggestions": [],
    }


def _qa_requested_chart_options(question: str, has_graph_data: bool = False) -> list[str]:
    if not has_graph_data:
        return []
    normalized = " ".join(str(question or "").strip().lower().split())
    preferred = []
    if any(token in normalized for token in ["doughnut", "donut", "donut chart"]):
        preferred.append("doughnut")
    elif any(token in normalized for token in ["pie chart", "pie graph", "pie"]):
        preferred.append("pie")
    elif any(token in normalized for token in ["line chart", "line chat", "line graph", "line plot", "trend line", "timeline", "over time", "trend"]):
        preferred.append("line")
    elif any(token in normalized for token in ["bar chart", "bar graph", "bar plot", "column chart", "column graph", "histogram"]):
        preferred.append("bar")

    ordered = []
    seen = set()
    for option in preferred + ["bar", "line", "pie", "doughnut"]:
        if option in seen:
            continue
        ordered.append(option)
        seen.add(option)
    return ordered


def _qa_merge_chart_options_for_question(question: str, options: Any, has_graph_data: bool = False) -> list[str]:
    normalized = _normalize_qa_chart_options(options, has_graph_data=has_graph_data)
    preferred = _qa_requested_chart_options(question, has_graph_data=has_graph_data)
    if not preferred:
        return normalized
    merged = []
    seen = set()
    for option in preferred + normalized:
        if option in seen:
            continue
        merged.append(option)
        seen.add(option)
    return merged


def _normalize_qa_chart_options(options: Any, has_graph_data: bool = False) -> list[str]:
    allowed = ("bar", "pie", "line", "doughnut")
    values = options if isinstance(options, list) else []
    normalized = []
    seen = set()
    for option in values:
        chart_type = str(option or "").strip().lower()
        if chart_type not in allowed or chart_type in seen:
            continue
        normalized.append(chart_type)
        seen.add(chart_type)
    if normalized:
        return normalized
    return ["bar", "pie"] if has_graph_data else []


def _normalize_qa_graph_data(graph_data: Any) -> dict:
    payload = graph_data if isinstance(graph_data, dict) else {}
    labels = payload.get("labels")
    values = payload.get("values")

    if not isinstance(labels, list) or not isinstance(values, list):
        points = payload.get("points")
        if isinstance(points, list):
            labels = []
            values = []
            for point in points:
                if not isinstance(point, dict):
                    continue
                label = str(point.get("label") or "").strip()
                value = point.get("value")
                if not label:
                    continue
                try:
                    numeric = float(value)
                except Exception:
                    continue
                labels.append(label[:60])
                values.append(round(numeric, 4))

    if not isinstance(labels, list) or not isinstance(values, list):
        return {}

    normalized_labels = []
    normalized_values = []
    for label, value in zip(labels, values):
        safe_label = str(label or "").strip()
        if not safe_label:
            continue
        try:
            numeric_value = float(value)
        except Exception:
            continue
        normalized_labels.append(safe_label[:60])
        normalized_values.append(round(numeric_value, 4))
        if len(normalized_labels) >= 24:
            break

    if len(normalized_labels) < 2 or len(normalized_labels) != len(normalized_values):
        return {}

    return {
        "title": str(payload.get("title") or "").strip()[:120],
        "x_label": str(payload.get("x_label") or "").strip()[:80],
        "y_label": str(payload.get("y_label") or "").strip()[:80],
        "labels": normalized_labels,
        "values": normalized_values,
    }


def _qa_graph_cleanup_locked(now: float) -> None:
    for session_id, payload in list(_qa_graph_store.items()):
        if not isinstance(payload, dict):
            _qa_graph_store.pop(session_id, None)
            continue
        if float(payload.get("expires_at") or 0) <= now:
            _qa_graph_store.pop(session_id, None)
            continue
        graphs = payload.get("graphs")
        if not isinstance(graphs, OrderedDict):
            graphs = OrderedDict(
                (str(graph_id), dict(graph_payload))
                for graph_id, graph_payload in (graphs or {}).items()
                if isinstance(graph_payload, dict)
            )
            payload["graphs"] = graphs
        for graph_id, graph_payload in list(graphs.items()):
            if float(graph_payload.get("expires_at") or 0) <= now:
                graphs.pop(graph_id, None)
        if not graphs:
            _qa_graph_store.pop(session_id, None)


def _qa_graph_rollover_for_question(session_id: str) -> None:
    safe_session_id = str(session_id or "").strip()
    if not safe_session_id:
        return

    now = time.time()
    with _qa_graph_lock:
        _qa_graph_cleanup_locked(now)
        payload = _qa_graph_store.get(safe_session_id)
        if not isinstance(payload, dict):
            return
        graphs = payload.get("graphs")
        if not isinstance(graphs, OrderedDict):
            graphs = OrderedDict()
        for graph_id, graph_payload in list(graphs.items()):
            age = int(graph_payload.get("question_age") or 0) + 1
            if age >= _QA_GRAPH_PURGE_AFTER_QUESTIONS:
                graphs.pop(graph_id, None)
                continue
            graph_payload["question_age"] = age
            graph_payload["expires_at"] = now + _QA_GRAPH_TTL

        if graphs:
            payload["graphs"] = graphs
            payload["expires_at"] = now + _QA_GRAPH_TTL
            _qa_graph_store[safe_session_id] = payload
            _qa_graph_store.move_to_end(safe_session_id)
        else:
            _qa_graph_store.pop(safe_session_id, None)


def _qa_graph_attach_to_response(session_id: str, response_payload: dict) -> dict:
    payload = dict(response_payload or {})
    raw_chart_options = payload.get("chart_options")
    payload.pop("graph_id", None)
    payload.pop("chart_options", None)

    safe_session_id = str(session_id or "").strip()
    graph_data = _normalize_qa_graph_data(payload.get("graph_data"))
    chart_options = _normalize_qa_chart_options(raw_chart_options, has_graph_data=bool(graph_data))
    if not safe_session_id or not graph_data:
        payload["graph_id"] = ""
        payload["chart_options"] = []
        return payload

    now = time.time()
    graph_id = str(uuid4())
    with _qa_graph_lock:
        _qa_graph_cleanup_locked(now)
        session_payload = _qa_graph_store.get(safe_session_id)
        if not isinstance(session_payload, dict):
            session_payload = {"graphs": OrderedDict()}
        graphs = session_payload.get("graphs")
        if not isinstance(graphs, OrderedDict):
            graphs = OrderedDict()

        graphs[graph_id] = {
            "graph_id": graph_id,
            "graph_data": graph_data,
            "chart_options": chart_options,
            "question_age": 0,
            "created_at": now,
            "expires_at": now + _QA_GRAPH_TTL,
        }
        graphs.move_to_end(graph_id)
        while len(graphs) > _QA_GRAPH_MAX_ITEMS_PER_SESSION:
            graphs.popitem(last=False)

        session_payload["graphs"] = graphs
        session_payload["expires_at"] = now + _QA_GRAPH_TTL
        _qa_graph_store[safe_session_id] = session_payload
        _qa_graph_store.move_to_end(safe_session_id)
        while len(_qa_graph_store) > _QA_GRAPH_MAX_SESSIONS:
            _qa_graph_store.popitem(last=False)

    payload["graph_id"] = graph_id
    payload["chart_options"] = chart_options
    return payload


def _qa_graph_get(session_id: str, graph_id: str) -> Optional[dict]:
    safe_session_id = str(session_id or "").strip()
    safe_graph_id = str(graph_id or "").strip()
    if not safe_session_id or not safe_graph_id:
        return None

    now = time.time()
    with _qa_graph_lock:
        _qa_graph_cleanup_locked(now)
        payload = _qa_graph_store.get(safe_session_id)
        if not isinstance(payload, dict):
            return None
        graphs = payload.get("graphs")
        if not isinstance(graphs, OrderedDict):
            return None
        graph_payload = graphs.get(safe_graph_id)
        if not isinstance(graph_payload, dict):
            return None
        graph_payload["expires_at"] = now + _QA_GRAPH_TTL
        payload["expires_at"] = now + _QA_GRAPH_TTL
        graphs.move_to_end(safe_graph_id)
        _qa_graph_store[safe_session_id] = payload
        _qa_graph_store.move_to_end(safe_session_id)
        return json.loads(json.dumps(graph_payload))


def _qa_render_chart_png(graph_data: dict, chart_type: str) -> bytes:
    labels = list(graph_data.get("labels") or [])
    values = list(graph_data.get("values") or [])
    if len(labels) < 2 or len(labels) != len(values):
        raise ValueError("Invalid graph data")

    numeric_values = [float(value) for value in values]
    safe_chart_type = str(chart_type or "").strip().lower()
    palette = [
        "#f97316",
        "#0f766e",
        "#2563eb",
        "#dc2626",
        "#7c3aed",
        "#ca8a04",
        "#0891b2",
        "#be185d",
        "#65a30d",
        "#ea580c",
    ]
    segment_colors = [palette[idx % len(palette)] for idx in range(len(labels))]

    figure, axis = plt.subplots(figsize=(8, 5), dpi=120)
    title = str(graph_data.get("title") or "").strip()
    x_label = str(graph_data.get("x_label") or "").strip()
    y_label = str(graph_data.get("y_label") or "").strip()

    if safe_chart_type == "pie":
        axis.pie(
            numeric_values,
            labels=labels,
            autopct="%1.1f%%",
            startangle=90,
            colors=segment_colors,
            wedgeprops={"linewidth": 1, "edgecolor": "#ffffff"},
        )
        axis.axis("equal")
    elif safe_chart_type == "doughnut":
        axis.pie(
            numeric_values,
            labels=labels,
            autopct="%1.1f%%",
            startangle=90,
            colors=segment_colors,
            wedgeprops={"width": 0.45, "linewidth": 1, "edgecolor": "#ffffff"},
        )
        axis.axis("equal")
    elif safe_chart_type == "line":
        axis.plot(labels, numeric_values, marker="o", color="#ea580c", linewidth=2)
        axis.grid(axis="y", alpha=0.25, linestyle="--")
        for idx, value in enumerate(numeric_values):
            axis.text(idx, value, f"{value:g}", ha="center", va="bottom", fontsize=8)
    else:
        bars = axis.bar(labels, numeric_values, color=segment_colors, edgecolor="#9a3412")
        axis.bar_label(bars, labels=[f"{value:g}" for value in numeric_values], padding=3, fontsize=8)
        axis.grid(axis="y", alpha=0.2, linestyle="--")

    if title:
        axis.set_title(title)
    if x_label and safe_chart_type not in {"pie", "doughnut"}:
        axis.set_xlabel(x_label)
    if y_label and safe_chart_type not in {"pie", "doughnut"}:
        axis.set_ylabel(y_label)
    if safe_chart_type not in {"pie", "doughnut"}:
        axis.tick_params(axis="x", rotation=25, labelsize=8)

    figure.tight_layout()
    image_buffer = BytesIO()
    figure.savefig(image_buffer, format="png")
    plt.close(figure)
    image_buffer.seek(0)
    return image_buffer.read()


def _qa_conversation_get(session_id: str) -> list[dict]:
    safe_session_id = str(session_id or "").strip()
    if not safe_session_id:
        return []

    now = time.time()
    with _qa_conversation_lock:
        payload = _qa_conversation_store.get(safe_session_id)
        if not payload:
            return []
        if float(payload.get("expires_at") or 0) <= now:
            _qa_conversation_store.pop(safe_session_id, None)
            return []
        _qa_conversation_store.move_to_end(safe_session_id)
        turns = payload.get("turns") if isinstance(payload.get("turns"), list) else []
        return [dict(turn) for turn in turns if isinstance(turn, dict)]


def _qa_conversation_set(session_id: str, turns: list[dict]) -> None:
    safe_session_id = str(session_id or "").strip()
    if not safe_session_id:
        return
    cleaned_turns = [dict(turn) for turn in (turns or []) if isinstance(turn, dict)][-_QA_CONVERSATION_MAX_TURNS:]
    with _qa_conversation_lock:
        _qa_conversation_store[safe_session_id] = {
            "expires_at": time.time() + _QA_CONVERSATION_TTL,
            "turns": cleaned_turns,
        }
        _qa_conversation_store.move_to_end(safe_session_id)
        while len(_qa_conversation_store) > _QA_CONVERSATION_MAX_SESSIONS:
            _qa_conversation_store.popitem(last=False)


def _qa_conversation_append(session_id: str, turn: dict) -> None:
    history = _qa_conversation_get(session_id)
    history.append(dict(turn or {}))
    _qa_conversation_set(session_id, history)


def _qa_history_user_key(request: Request) -> str:
    actor = _normalize_kbac_owner(_get_request_actor(request))
    if actor:
        return actor
    session_id = _get_request_session_id(request)
    return f"session:{session_id}" if session_id else ""


def _qa_history_make_title(question: str) -> str:
    normalized = " ".join(str(question or "").strip().split())
    if not normalized:
        return "New chat"
    return normalized[:80].rstrip()


def _qa_history_make_preview(turn: dict) -> str:
    question = " ".join(str(turn.get("question") or "").strip().split())
    answer = " ".join(str(turn.get("answer") or "").strip().split())
    preview = question or answer
    return preview[:120].rstrip()


def _qa_history_build_turn_payload(
    question: str,
    response_payload: dict,
    topic: str = "",
    *,
    turn_id: str = "",
    created_at: Optional[str] = None,
) -> dict:
    safe_timestamp = str(created_at or "").strip() or _to_utc_iso(datetime.now(timezone.utc))
    revision_count = _qa_normalize_revision_count(response_payload.get("revision_count"), default=1)
    return {
        "turn_id": str(turn_id or "").strip() or str(uuid4()),
        "question": str(question or "").strip(),
        "answer": str(response_payload.get("answer") or "").strip(),
        "suggestion": str(response_payload.get("suggestion") or "").strip(),
        "revision_count": revision_count,
        "structured_answer": dict(response_payload.get("structured_answer") or {}),
        "citations": list(response_payload.get("citations") or []),
        "sources": list(response_payload.get("sources") or []),
        "chunk_sources": list(response_payload.get("chunk_sources") or []),
        "graph_id": str(response_payload.get("graph_id") or "").strip(),
        "chart_options": list(response_payload.get("chart_options") or []),
        "follow_up_suggestions": list(response_payload.get("follow_up_suggestions") or []),
        "topic": str(topic or "").strip(),
        "created_at": safe_timestamp,
    }


def _qa_history_ensure_turn_ids_locked(conversation: dict) -> bool:
    if not isinstance(conversation, dict):
        return False
    turns = conversation.get("turns")
    if not isinstance(turns, list):
        return False
    changed = False
    normalized_turns = []
    for turn in turns:
        if not isinstance(turn, dict):
            continue
        payload = dict(turn)
        if not str(payload.get("turn_id") or "").strip():
            payload["turn_id"] = str(uuid4())
            changed = True
        revision_count = _qa_normalize_revision_count(payload.get("revision_count"), default=1)
        if payload.get("revision_count") != revision_count:
            payload["revision_count"] = revision_count
            changed = True
        normalized_turns.append(payload)
    if changed:
        conversation["turns"] = normalized_turns
    return changed


def _qa_history_cleanup_locked(now: float) -> None:
    expired_users = []
    for user_key, payload in list(_qa_history_store.items()):
        conversations = payload.get("conversations")
        if not isinstance(conversations, OrderedDict):
            conversations = OrderedDict(
                (str(conv_id), dict(conv_payload))
                for conv_id, conv_payload in (conversations or {}).items()
                if isinstance(conv_payload, dict)
            )
            payload["conversations"] = conversations
        for conversation_id, conversation in list(conversations.items()):
            if float(conversation.get("expires_at") or 0) <= now:
                conversations.pop(conversation_id, None)
        if not conversations:
            expired_users.append(user_key)
    for user_key in expired_users:
        _qa_history_store.pop(user_key, None)


def _qa_history_get_conversation(user_key: str, conversation_id: str) -> Optional[dict]:
    safe_user_key = str(user_key or "").strip().lower()
    safe_conversation_id = str(conversation_id or "").strip()
    if not safe_user_key or not safe_conversation_id:
        return None

    now = time.time()
    with _qa_history_lock:
        _qa_history_cleanup_locked(now)
        payload = _qa_history_store.get(safe_user_key)
        if not isinstance(payload, dict):
            return None
        conversations = payload.get("conversations")
        if not isinstance(conversations, OrderedDict):
            return None
        conversation = conversations.get(safe_conversation_id)
        if not isinstance(conversation, dict):
            return None
        if _qa_history_ensure_turn_ids_locked(conversation):
            conversations[safe_conversation_id] = conversation
        conversations.move_to_end(safe_conversation_id)
        _qa_history_store.move_to_end(safe_user_key)
        return json.loads(json.dumps(conversation))


def _qa_history_get_turns(user_key: str, conversation_id: str) -> list[dict]:
    conversation = _qa_history_get_conversation(user_key, conversation_id)
    turns = conversation.get("turns") if isinstance(conversation, dict) else []
    return [dict(turn) for turn in turns if isinstance(turn, dict)]


def _qa_history_list_conversations(user_key: str) -> list[dict]:
    safe_user_key = str(user_key or "").strip().lower()
    if not safe_user_key:
        return []

    now = time.time()
    with _qa_history_lock:
        _qa_history_cleanup_locked(now)
        payload = _qa_history_store.get(safe_user_key)
        if not isinstance(payload, dict):
            return []
        conversations = payload.get("conversations")
        if not isinstance(conversations, OrderedDict):
            return []
        _qa_history_store.move_to_end(safe_user_key)
        items = []
        for conversation_id, conversation in reversed(list(conversations.items())):
            if not isinstance(conversation, dict):
                continue
            turns = conversation.get("turns") if isinstance(conversation.get("turns"), list) else []
            last_turn = turns[-1] if turns else {}
            items.append({
                "conversation_id": conversation_id,
                "title": str(conversation.get("title") or "").strip() or "New chat",
                "preview": _qa_history_make_preview(last_turn),
                "updated_at": conversation.get("updated_at"),
                "turn_count": len(turns),
            })
        return items


def _qa_history_store_turn(
    user_key: str,
    conversation_id: str,
    question: str,
    response_payload: dict,
    topic: str = "",
) -> str:
    stored_conversation_id, _ = _qa_history_store_turn_result(
        user_key,
        conversation_id,
        question,
        response_payload,
        topic=topic,
    )
    return stored_conversation_id


def _qa_history_store_turn_result(
    user_key: str,
    conversation_id: str,
    question: str,
    response_payload: dict,
    topic: str = "",
) -> tuple[str, str]:
    safe_user_key = str(user_key or "").strip().lower()
    safe_conversation_id = str(conversation_id or "").strip() or str(uuid4())
    if not safe_user_key:
        turn_id = str(response_payload.get("turn_id") or "").strip() or str(uuid4())
        return safe_conversation_id, turn_id

    now_dt = datetime.now(timezone.utc)
    now_ts = time.time()
    timestamp = _to_utc_iso(now_dt)
    response_payload["revision_count"] = 1
    turn_payload = _qa_history_build_turn_payload(
        question,
        response_payload,
        topic=topic,
        turn_id=str(response_payload.get("turn_id") or "").strip(),
        created_at=timestamp,
    )

    with _qa_history_lock:
        _qa_history_cleanup_locked(now_ts)
        user_payload = _qa_history_store.get(safe_user_key)
        if not isinstance(user_payload, dict):
            user_payload = {"conversations": OrderedDict()}
        conversations = user_payload.get("conversations")
        if not isinstance(conversations, OrderedDict):
            conversations = OrderedDict()
            user_payload["conversations"] = conversations

        conversation = conversations.get(safe_conversation_id)
        if not isinstance(conversation, dict):
            conversation = {
                "conversation_id": safe_conversation_id,
                "title": _qa_history_make_title(question),
                "created_at": timestamp,
                "updated_at": timestamp,
                "expires_at": now_ts + _QA_HISTORY_TTL,
                "turns": [],
            }

        turns = conversation.get("turns") if isinstance(conversation.get("turns"), list) else []
        turns.append(turn_payload)
        conversation["turns"] = turns[-_QA_HISTORY_MAX_TURNS_PER_CONVERSATION:]
        conversation["title"] = str(conversation.get("title") or "").strip() or _qa_history_make_title(question)
        conversation["updated_at"] = timestamp
        conversation["expires_at"] = now_ts + _QA_HISTORY_TTL

        conversations[safe_conversation_id] = conversation
        conversations.move_to_end(safe_conversation_id)
        while len(conversations) > _QA_HISTORY_MAX_CONVERSATIONS_PER_USER:
            conversations.popitem(last=False)

        _qa_history_store[safe_user_key] = user_payload
        _qa_history_store.move_to_end(safe_user_key)
        while len(_qa_history_store) > _QA_HISTORY_MAX_USERS:
            _qa_history_store.popitem(last=False)

    return safe_conversation_id, str(turn_payload.get("turn_id") or "")


def _qa_history_get_turn_slice(
    user_key: str,
    conversation_id: str,
    turn_id: str,
) -> tuple[Optional[dict], int, list[dict]]:
    conversation = _qa_history_get_conversation(user_key, conversation_id)
    if not isinstance(conversation, dict):
        return None, -1, []
    turns = [dict(turn) for turn in conversation.get("turns") or [] if isinstance(turn, dict)]
    safe_turn_id = str(turn_id or "").strip()
    target_index = next(
        (
            index for index, turn in enumerate(turns)
            if str(turn.get("turn_id") or "").strip() == safe_turn_id
        ),
        -1,
    )
    return conversation, target_index, turns


def _qa_history_find_conversation_id_by_turn_id(user_key: str, turn_id: str) -> str:
    safe_user_key = str(user_key or "").strip().lower()
    safe_turn_id = str(turn_id or "").strip()
    if not safe_user_key or not safe_turn_id:
        return ""

    now = time.time()
    with _qa_history_lock:
        _qa_history_cleanup_locked(now)
        payload = _qa_history_store.get(safe_user_key)
        if not isinstance(payload, dict):
            return ""
        conversations = payload.get("conversations")
        if not isinstance(conversations, OrderedDict):
            return ""
        for conversation_id, conversation in conversations.items():
            if not isinstance(conversation, dict):
                continue
            _qa_history_ensure_turn_ids_locked(conversation)
            turns = conversation.get("turns") if isinstance(conversation.get("turns"), list) else []
            if any(str(turn.get("turn_id") or "").strip() == safe_turn_id for turn in turns if isinstance(turn, dict)):
                return str(conversation_id)
    return ""


def _qa_history_replace_turn_result(
    user_key: str,
    conversation_id: str,
    target_turn_id: str,
    question: str,
    response_payload: dict,
    topic: str = "",
) -> tuple[str, str]:
    safe_user_key = str(user_key or "").strip().lower()
    safe_conversation_id = str(conversation_id or "").strip()
    safe_turn_id = str(target_turn_id or "").strip()
    if not safe_user_key or not safe_conversation_id or not safe_turn_id:
        raise ValueError("Conversation and turn identifiers are required")

    now_dt = datetime.now(timezone.utc)
    now_ts = time.time()
    timestamp = _to_utc_iso(now_dt)

    with _qa_history_lock:
        _qa_history_cleanup_locked(now_ts)
        user_payload = _qa_history_store.get(safe_user_key)
        if not isinstance(user_payload, dict):
            raise KeyError("Conversation not found")
        conversations = user_payload.get("conversations")
        if not isinstance(conversations, OrderedDict):
            raise KeyError("Conversation not found")
        conversation = conversations.get(safe_conversation_id)
        if not isinstance(conversation, dict):
            raise KeyError("Conversation not found")

        _qa_history_ensure_turn_ids_locked(conversation)
        turns = [dict(turn) for turn in conversation.get("turns") or [] if isinstance(turn, dict)]
        target_index = next(
            (
                index for index, turn in enumerate(turns)
                if str(turn.get("turn_id") or "").strip() == safe_turn_id
            ),
            -1,
        )
        if target_index < 0:
            raise KeyError("Turn not found")

        existing_turn = turns[target_index]
        try:
            next_revision_count = max(1, int(existing_turn.get("revision_count") or 1)) + 1
        except Exception:
            next_revision_count = 2
        response_payload["revision_count"] = next_revision_count
        replacement_turn = _qa_history_build_turn_payload(
            question,
            response_payload,
            topic=topic,
            turn_id=safe_turn_id,
            created_at=existing_turn.get("created_at") or timestamp,
        )
        turns = turns[:target_index] + [replacement_turn]

        conversation["turns"] = turns[-_QA_HISTORY_MAX_TURNS_PER_CONVERSATION:]
        first_question = str(conversation["turns"][0].get("question") or "").strip() if conversation["turns"] else ""
        conversation["title"] = _qa_history_make_title(first_question)
        conversation["updated_at"] = timestamp
        conversation["expires_at"] = now_ts + _QA_HISTORY_TTL

        conversations[safe_conversation_id] = conversation
        conversations.move_to_end(safe_conversation_id)
        _qa_history_store[safe_user_key] = user_payload
        _qa_history_store.move_to_end(safe_user_key)

    return safe_conversation_id, safe_turn_id


def _qa_history_filter_turn_for_role(user_role: UserRole, turn: dict) -> dict:
    response_payload = {
        "turn_id": str(turn.get("turn_id") or "").strip(),
        "answer": str(turn.get("answer") or "").strip(),
        "suggestion": str(turn.get("suggestion") or "").strip(),
        "revision_count": _qa_normalize_revision_count(turn.get("revision_count"), default=1),
        "structured_answer": dict(turn.get("structured_answer") or {}),
        "citations": list(turn.get("citations") or []),
        "sources": list(turn.get("sources") or []),
        "chunk_sources": list(turn.get("chunk_sources") or []),
        "graph_id": str(turn.get("graph_id") or "").strip(),
        "chart_options": list(turn.get("chart_options") or []),
        "follow_up_suggestions": list(turn.get("follow_up_suggestions") or []),
    }
    filtered = filter_qa_response_by_role(user_role, dict(response_payload))
    return {
        "turn_id": str(turn.get("turn_id") or "").strip(),
        "question": str(turn.get("question") or "").strip(),
        "answer": str(filtered.get("answer") or "").strip(),
        "suggestion": str(filtered.get("suggestion") or "").strip(),
        "revision_count": _qa_normalize_revision_count(
            filtered.get("revision_count") or turn.get("revision_count"),
            default=1,
        ),
        "structured_answer": dict(filtered.get("structured_answer") or {}),
        "citations": list(filtered.get("citations") or []),
        "sources": list(filtered.get("sources") or []),
        "chunk_sources": list(filtered.get("chunk_sources") or []),
        "graph_id": str(filtered.get("graph_id") or "").strip(),
        "chart_options": list(filtered.get("chart_options") or []),
        "follow_up_suggestions": list(filtered.get("follow_up_suggestions") or []),
        "created_at": turn.get("created_at"),
    }


def _qa_normalize_revision_count(value: Any, default: int = 1) -> int:
    try:
        return max(1, int(value or default))
    except Exception:
        return max(1, int(default or 1))


def _qa_history_filter_live_graphs(session_id: str, turns: list[dict]) -> list[dict]:
    safe_session_id = str(session_id or "").strip()
    if not safe_session_id:
        return turns
    sanitized_turns = []
    for turn in turns:
        if not isinstance(turn, dict):
            continue
        payload = dict(turn)
        graph_id = str(payload.get("graph_id") or "").strip()
        if graph_id and not _qa_graph_get(safe_session_id, graph_id):
            payload["graph_id"] = ""
            payload["chart_options"] = []
        sanitized_turns.append(payload)
    return sanitized_turns


def _qa_is_contextual_followup(question: str) -> bool:
    normalized = " ".join(str(question or "").strip().lower().split())
    if not normalized:
        return False
    followup_markers = [
        "his ", "her ", "their ", "them ", "they ", "he ", "she ",
        "that ", "those ", "these ", "this ", "same ",
        "contact details", "phone number", "mobile number", "email address",
        "share me", "what about", "can you share", "tell me more",
        "based on that", "based on this", "from that", "from this",
    ]
    return any(marker in normalized for marker in followup_markers)


def _qa_is_situation_question(question: str) -> bool:
    normalized = " ".join(str(question or "").strip().lower().split())
    if not normalized:
        return False
    situation_markers = [
        "should i hire", "should we hire", "should hire",
        "fit for", "good fit", "suitable for", "right for",
        "recommend", "recommended", "eligible for", "capable of",
        "good for an", "good for a", "hire for",
    ]
    return any(marker in normalized for marker in situation_markers)


def _qa_build_context_signature(turns: list[dict]) -> str:
    if not turns:
        return "fresh"
    payload = []
    for turn in turns[-2:]:
        payload.append({
            "question": str(turn.get("question") or "")[:160],
            "topic": str(turn.get("topic") or "")[:120],
            "answer": str(turn.get("answer") or "")[:220],
        })
    encoded = json.dumps(payload, sort_keys=True, ensure_ascii=True).encode("utf-8")
    return hashlib.sha256(encoded).hexdigest()[:16]


def _qa_build_contextual_question(question: str, history: list[dict]) -> str:
    current_question = str(question or "").strip()
    if not current_question:
        return ""

    last_turn = history[-1] if history else {}
    retrieval_question = current_question

    if history and _qa_is_contextual_followup(current_question):
        previous_question = str(last_turn.get("question") or "").strip()
        previous_topic = str(last_turn.get("topic") or "").strip()
        previous_answer = str(last_turn.get("answer") or "").strip()
        contextual_parts = [current_question]
        if previous_question:
            contextual_parts.append(f"Previous question: {previous_question}")
        if previous_topic:
            contextual_parts.append(f"Previous topic: {previous_topic}")
        if previous_answer:
            contextual_parts.append(f"Previous answer summary: {previous_answer[:220]}")
        retrieval_question = "\n".join(contextual_parts)

    if _qa_is_situation_question(current_question):
        retrieval_question = (
            f"{retrieval_question}\n"
            "Relevant evidence to consider: resume skills experience projects education certifications achievements role fit ai ml machine learning."
        )

    return retrieval_question


def _qa_merge_sources(primary_sources: list[dict], secondary_sources: list[dict]) -> list[dict]:
    merged = []
    seen = set()
    for source in list(primary_sources or []) + list(secondary_sources or []):
        if not isinstance(source, dict):
            continue
        doc_id = str(source.get("doc_id") or "").strip()
        snippet = str(source.get("snippet") or "").strip()
        key = (doc_id, snippet[:160], str(source.get("document_type") or "").strip())
        if key in seen:
            continue
        seen.add(key)
        merged.append(dict(source))
    return merged


def _qa_collect_context_sources(question: str, history: list[dict]) -> list[dict]:
    if not history:
        return []

    collected = []
    for turn in reversed(history[-2:]):
        turn_sources = turn.get("sources") if isinstance(turn.get("sources"), list) else []
        if not turn_sources:
            continue
        collected = _qa_merge_sources(collected, turn_sources)
        if _qa_is_contextual_followup(question):
            break
    return collected[:8]


def _qa_build_conversation_context(history: list[dict]) -> str:
    if not history:
        return ""
    lines = []
    for idx, turn in enumerate(history[-2:], start=1):
        question = str(turn.get("question") or "").strip()
        answer = str(turn.get("answer") or "").strip()
        topic = str(turn.get("topic") or "").strip()
        if question:
            lines.append(f"Turn {idx} user question: {question}")
        if topic:
            lines.append(f"Turn {idx} topic: {topic}")
        if answer:
            lines.append(f"Turn {idx} assistant answer: {answer[:260]}")
    return "\n".join(lines)


def _qa_question_requests_graph(question: str) -> bool:
    normalized = " ".join(str(question or "").strip().lower().split())
    if not normalized:
        return False
    phrase_triggers = [
        "graph", "chart", "plot", "visual", "visualize", "visualise",
        "bar chart", "pie chart", "line chart", "line chat", "doughnut", "donut chart",
        "histogram", "trend", "distribution", "breakdown", "analytics view",
        "represent in graph", "show in graph", "show as chart",
        "make a chart", "generate chart", "draw a graph", "graphical view",
        "side by side chart", "comparison chart",
    ]
    regex_triggers = [
        r"\bplot\b.*\b(vs|against|between|comparison)\b",
        r"\bshow\b.*\b(chart|graph|visual)\b",
        r"\bcompare\b.*\b(chart|graph|visual)\b",
    ]
    if any(token in normalized for token in phrase_triggers):
        return True
    return any(re.search(pattern, normalized, flags=re.IGNORECASE) for pattern in regex_triggers)


def _qa_requests_broad_document_listing(question: str) -> bool:
    normalized = " ".join(str(question or "").strip().lower().split())
    normalized = re.sub(r"\bpo['’]s\b", "po", normalized, flags=re.IGNORECASE)
    if not normalized:
        return False

    document_markers = [
        "po", "purchase order", "invoice", "document", "vendor", "bill",
    ]
    broad_markers = [
        "all", "all the", "every", "list", "show", "what are all", "tell me about all",
    ]
    broad_patterns = [
        r"\b(?:all|every)\s+(?:the\s+)?(?:po|purchase order|purchase orders|invoice|invoices|documents)\b",
        r"\b(?:list|show|tell me about)\s+(?:me\s+)?(?:all\s+)?(?:the\s+)?(?:po|purchase order|purchase orders|invoice|invoices|documents)\b",
    ]

    if not any(marker in normalized for marker in document_markers):
        return False
    if any(re.search(pattern, normalized, flags=re.IGNORECASE) for pattern in broad_patterns):
        return True
    return any(marker in normalized for marker in broad_markers) and any(
        marker in normalized for marker in document_markers
    )


def _qa_requests_all_purchase_orders(question: str) -> bool:
    normalized = " ".join(str(question or "").strip().lower().split())
    normalized = re.sub(r"\bpo['’]s\b", "po", normalized, flags=re.IGNORECASE)
    if not normalized:
        return False
    if not any(marker in normalized for marker in ["po", "purchase order", "purchase orders"]):
        return False
    return _qa_requests_broad_document_listing(normalized)


def _qa_extract_label_candidates(text: str) -> list[str]:
    raw = str(text or "")
    if not raw:
        return []
    patterns = [
        r"\b(?:INV|INVOICE)[-\s]?\d{1,5}\b",
        r"\b(?:PO|PURCHASE ORDER)[-\s]?\d{1,5}\b",
    ]
    labels = []
    seen = set()
    for pattern in patterns:
        for match in re.findall(pattern, raw, flags=re.IGNORECASE):
            normalized = " ".join(str(match).upper().replace("PURCHASE ORDER", "PO").replace("INVOICE", "INV").split())
            normalized = normalized.replace(" ", "-")
            if normalized in seen:
                continue
            seen.add(normalized)
            labels.append(normalized)
    return labels


def _qa_extract_amount_candidates(text: str) -> list[float]:
    raw = str(text or "")
    if not raw:
        return []

    compact = raw.replace("â‚¹", "₹")
    matches = re.findall(
        r"(?:₹|rs\.?|inr)?\s*([0-9]{1,3}(?:,[0-9]{2,3})*(?:\.[0-9]+)?|[0-9]+(?:\.[0-9]+)?)",
        compact,
        flags=re.IGNORECASE,
    )
    values = []
    for token in matches:
        cleaned = str(token or "").replace(",", "").strip()
        if not cleaned:
            continue
        try:
            value = float(cleaned)
        except Exception:
            continue
        if value <= 0:
            continue
        values.append(round(value, 4))
    return values


def _qa_normalize_doc_ref(prefix: str, number_token: str) -> str:
    prefix_clean = str(prefix or "").strip().upper()
    digits = re.sub(r"[^0-9]", "", str(number_token or ""))
    if not prefix_clean or not digits:
        return ""
    return f"{prefix_clean}-{digits.zfill(3)}"


def _qa_extract_requested_doc_refs(question: str) -> set[str]:
    raw = str(question or "")
    refs = set()
    for match in re.finditer(r"\b(?:PO|PURCHASE\s*ORDER)\s*-?\s*(\d{1,6})\b", raw, flags=re.IGNORECASE):
        normalized = _qa_normalize_doc_ref("PO", match.group(1))
        if normalized:
            refs.add(normalized)
    for match in re.finditer(r"\b(?:INV|INVOICE)\s*-?\s*(\d{1,6})\b", raw, flags=re.IGNORECASE):
        normalized = _qa_normalize_doc_ref("INV", match.group(1))
        if normalized:
            refs.add(normalized)
    return refs


def _qa_source_doc_refs(source: dict) -> set[str]:
    if not isinstance(source, dict):
        return set()
    refs = set()
    corpus = " ".join(
        [
            str(source.get("doc_id") or ""),
            str(source.get("document_type") or ""),
            str(source.get("metadata") or ""),
            str(source.get("snippet") or ""),
            str(source.get("full_text") or "")[:3000],
        ]
    )
    for label in _qa_extract_label_candidates(corpus):
        parts = label.split("-", 1)
        if len(parts) == 2:
            normalized = _qa_normalize_doc_ref(parts[0], parts[1])
            if normalized:
                refs.add(normalized)
    for match in re.finditer(r"\b(?:PO|PURCHASE\s*ORDER)\s*-?\s*(\d{1,6})\b", corpus, flags=re.IGNORECASE):
        normalized = _qa_normalize_doc_ref("PO", match.group(1))
        if normalized:
            refs.add(normalized)
    for match in re.finditer(r"\b(?:INV|INVOICE)\s*-?\s*(\d{1,6})\b", corpus, flags=re.IGNORECASE):
        normalized = _qa_normalize_doc_ref("INV", match.group(1))
        if normalized:
            refs.add(normalized)
    return refs


def _qa_is_devices_vs_amount_request(question: str) -> bool:
    normalized = " ".join(str(question or "").strip().lower().split())
    if not normalized:
        return False
    amount_terms = [
        "amount", "value", "price", "total", "cost", "spend", "expense",
        "bill value", "invoice value", "net amount", "grand total", "payable",
        "money", "monetary", "revenue", "charges",
    ]
    quantity_terms = [
        "device", "devices", "laptop", "laptops", "quantity", "qty",
        "unit", "units", "count", "pieces", "item", "items", "products",
        "materials", "equipment", "shipment count",
    ]
    comparison_terms = [
        "vs", "versus", "against", "compare", "comparison",
        "side by side", "between", "across", "each", "all", "wise",
        "item wise", "device wise", "quantity wise",
    ]
    has_amount = any(term in normalized for term in amount_terms)
    has_quantity = any(term in normalized for term in quantity_terms)
    has_comparison = any(term in normalized for term in comparison_terms)
    return has_amount and has_quantity and has_comparison


def _qa_best_source_quantity(source: dict) -> Optional[float]:
    if not isinstance(source, dict):
        return None
    corpus = " ".join(
        [
            str(source.get("metadata") or ""),
            str(source.get("snippet") or ""),
            str(source.get("full_text") or "")[:2500],
        ]
    )
    qty_matches = re.findall(r"quantity[^0-9]{0,16}([0-9]{1,5}(?:\.[0-9]+)?)", corpus, flags=re.IGNORECASE)
    if not qty_matches:
        qty_matches = re.findall(r"\b([0-9]{1,5}(?:\.[0-9]+)?)\s*(?:laptops?|devices?|units?)\b", corpus, flags=re.IGNORECASE)
    values = []
    for token in qty_matches:
        try:
            value = float(str(token).replace(",", ""))
        except Exception:
            continue
        if value <= 0:
            continue
        values.append(value)
    return values[0] if values else None


def _qa_best_source_label(source: dict) -> str:
    if not isinstance(source, dict):
        return ""
    label_candidates = []
    for field_name in ("metadata", "snippet", "full_text"):
        value = source.get(field_name)
        if value:
            label_candidates.extend(_qa_extract_label_candidates(str(value)))
        if label_candidates:
            break
    if label_candidates:
        return label_candidates[0]
    doc_type = str(source.get("document_type") or "Document").strip()
    doc_id = str(source.get("doc_id") or "").strip()[:8]
    return f"{doc_type}-{doc_id}" if doc_id else doc_type


def _qa_best_source_amount(source: dict) -> Optional[float]:
    if not isinstance(source, dict):
        return None
    combined_text = "\n".join(
        [
            str(source.get("metadata") or ""),
            str(source.get("snippet") or ""),
            str(source.get("full_text") or "")[:2000],
        ]
    )
    amounts = _qa_extract_amount_candidates(combined_text)
    if not amounts:
        return None
    return max(amounts)


def _qa_build_graph_data_from_sources(question: str, sources: list[dict]) -> dict:
    if not _qa_question_requests_graph(question):
        return {}
    if not isinstance(sources, list) or not sources:
        return {}

    requested_refs = _qa_extract_requested_doc_refs(question)
    filtered_sources = [src for src in sources if isinstance(src, dict)]
    if requested_refs:
        strict_matches = []
        for source in filtered_sources:
            source_refs = _qa_source_doc_refs(source)
            if source_refs.intersection(requested_refs):
                strict_matches.append(source)
        if len(strict_matches) >= 2:
            filtered_sources = strict_matches
        else:
            return {}

    points = []
    seen_labels = set()
    wants_devices_vs_amount = _qa_is_devices_vs_amount_request(question)
    for source in filtered_sources:
        label = _qa_best_source_label(source)
        amount = _qa_best_source_amount(source)
        if not label or amount is None:
            continue
        quantity = _qa_best_source_quantity(source)
        if wants_devices_vs_amount and quantity is not None:
            safe_label = f"{label} ({int(quantity) if float(quantity).is_integer() else round(quantity, 2)} devices)"[:60]
        else:
            safe_label = label[:60]
        if safe_label in seen_labels:
            continue
        seen_labels.add(safe_label)
        points.append((safe_label, amount))
        if len(points) >= 12:
            break

    if len(points) >= 2:
        points.sort(key=lambda item: item[1], reverse=True)
        labels = [label for label, _ in points]
        values = [value for _, value in points]
        chart_title = "Comparison of Retrieved Document Amounts"
        if wants_devices_vs_amount:
            chart_title = "Devices vs Amount Comparison"
        elif requested_refs:
            chart_title = "Requested Document Amount Comparison"
        return {
            "title": chart_title,
            "x_label": "Document",
            "y_label": "Amount",
            "labels": labels,
            "values": values,
        }

    source_types = [str(src.get("document_type") or "").strip().lower() for src in filtered_sources if isinstance(src, dict)]
    po_count = sum(1 for item in source_types if "purchase" in item or item == "po")
    invoice_count = sum(1 for item in source_types if "invoice" in item or item == "inv")
    if po_count + invoice_count >= 2:
        labels = []
        values = []
        if po_count:
            labels.append("Purchase Orders")
            values.append(float(po_count))
        if invoice_count:
            labels.append("Invoices")
            values.append(float(invoice_count))
        return {
            "title": "Document Distribution for This Comparison",
            "x_label": "Document Type",
            "y_label": "Count",
            "labels": labels,
            "values": values,
        }

    return {}


def _qa_finalize_graph_payload(
    question: str,
    graph_data: Any,
    chart_options: Any,
    sources: list[dict] | None = None,
) -> tuple[dict, list[str], bool]:
    if not _qa_question_requests_graph(question):
        return {}, [], False

    normalized_graph = _normalize_qa_graph_data(graph_data)
    if not normalized_graph and sources:
        normalized_graph = _qa_build_graph_data_from_sources(question, sources)

    normalized_options = _qa_merge_chart_options_for_question(
        question,
        chart_options,
        has_graph_data=bool(normalized_graph),
    )
    return normalized_graph, normalized_options, bool(normalized_graph)


def _qa_detect_followup_domain(question: str, answer: str, sources: list[dict]) -> str:
    corpus = " ".join(
        [
            str(question or "").lower(),
            str(answer or "").lower(),
            " ".join(str(src.get("document_type") or "").lower() for src in (sources or []) if isinstance(src, dict)),
            " ".join(str(src.get("metadata") or "").lower()[:400] for src in (sources or []) if isinstance(src, dict)),
        ]
    )

    procurement_markers = [
        "purchase order", "po-", "invoice", "gst", "vendor", "payment terms",
        "delivery", "amount", "total value", "bill",
    ]
    manual_markers = [
        "manual", "user guide", "documentation", "api doc", "sop", "procedure",
        "troubleshoot", "troubleshooting", "installation", "configuration",
        "workflow", "steps", "how to", "setup", "onboarding",
    ]
    resume_markers = [
        "resume", "cv", "candidate", "profile", "skills", "experience",
        "education", "internship", "project", "certification",
    ]

    procurement_score = sum(1 for marker in procurement_markers if marker in corpus)
    manual_score = sum(1 for marker in manual_markers if marker in corpus)
    resume_score = sum(1 for marker in resume_markers if marker in corpus)

    if procurement_score >= max(2, resume_score + 1, manual_score + 1):
        return "procurement"
    if resume_score >= max(2, procurement_score + 1, manual_score + 1):
        return "resume"
    if manual_score >= max(2, procurement_score + 1, resume_score + 1):
        return "manual"
    return "general"


def _qa_build_follow_up_suggestions(
    question: str,
    answer: str,
    sources: list[dict],
    has_graph_data: bool,
) -> list[str]:
    domain = _qa_detect_followup_domain(question, answer, sources)
    question_norm = " ".join(str(question or "").strip().lower().split())

    if domain == "procurement":
        candidates = [
            "Plot a bar chart comparing these PO and invoice amounts.",
            "Compare these documents by amount, date, and payment terms.",
            "Show invoice-to-PO mapping and flag any mismatch.",
            "Create a timeline of PO and invoice dates.",
            "Group totals vendor-wise and highlight the highest spend.",
            "Compare quantity vs amount to detect pricing anomalies.",
        ]
        if has_graph_data:
            candidates.insert(0, "Show this comparison as a pie chart by value.")
    elif domain == "resume":
        candidates = [
            "Summarize the candidate's strengths and skill gaps for a target role.",
            "Compare this profile against a job description with a fit summary.",
            "Generate interview questions based on this candidate's projects.",
            "Create an experience timeline with key achievements.",
            "Map skills to role requirements and show missing capabilities.",
            "Score this profile across technical depth, communication, and ownership.",
        ]
    elif domain == "manual":
        candidates = [
            "Convert this into a short step-by-step execution checklist.",
            "List prerequisites, dependencies, and environment setup points.",
            "Generate a troubleshooting table with symptom, cause, and fix.",
            "Extract key commands/API endpoints from this manual.",
            "Create a quick-start summary for first-time users.",
            "Compare this process with a safer fallback workflow.",
        ]
    else:
        candidates = [
            "Summarize this in 5 key bullet points with important numbers and dates.",
            "Compare this with related documents using amount and date differences.",
            "If numeric values exist, plot a quick bar chart for easier review.",
            "Extract action items and decisions from this response.",
            "Generate a concise checklist from this information.",
        ]
        if has_graph_data:
            candidates.insert(0, "Show this data as a pie chart for quick distribution view.")

    suggestions = []
    seen = set()
    for candidate in candidates:
        text = str(candidate or "").strip()
        normalized = " ".join(text.lower().split())
        if not text or normalized in seen:
            continue
        if normalized == question_norm:
            continue
        seen.add(normalized)
        suggestions.append(text)
        if len(suggestions) >= 2:
            break
    return suggestions


def _invalidate_qa_caches():
    clear_qa_response_cache()
    clear_qa_retrieval_cache()
    clear_qa_answer_cache()


def _build_answer_sources(doc_sources: list, chunk_sources: list) -> list:
    if chunk_sources:
        merged = list(chunk_sources)
        existing_ids = {
            str(src.get("doc_id") or "").strip()
            for src in chunk_sources
            if isinstance(src, dict)
        }
        for doc_src in doc_sources or []:
            doc_id = str(doc_src.get("doc_id") or "").strip()
            if not doc_id or doc_id in existing_ids:
                continue
            merged.append(doc_src)
        return merged
    return list(doc_sources or [])


def _build_qa_ui_sources(doc_sources: list, chunk_sources: list) -> list:
    sources_for_ui = list(doc_sources or [])
    if chunk_sources:
        existing_ids = {src.get("doc_id") for src in sources_for_ui if isinstance(src, dict)}
        for chunk_src in chunk_sources:
            doc_id = chunk_src.get("doc_id") if isinstance(chunk_src, dict) else None
            if not doc_id or doc_id in existing_ids:
                continue
            sources_for_ui.append(
                {
                    "doc_id": doc_id,
                    "document_type": chunk_src.get("document_type", "unknown"),
                    "metadata": chunk_src.get("metadata", ""),
                    "snippet": chunk_src.get("snippet", ""),
                    "items": chunk_src.get("items", []),
                    "full_text": "",
                    "confidence": chunk_src.get("confidence", {}),
                    "score": chunk_src.get("score", 0),
                    "reasons": ["chunk_match"],
                }
            )
            existing_ids.add(doc_id)
    return sources_for_ui


def _qa_stream_event(payload: dict) -> bytes:
    return (json.dumps(payload, ensure_ascii=False) + "\n").encode("utf-8")


def _to_utc_iso(dt: datetime) -> str:
    """Serialize datetime as explicit UTC ISO-8601 (with Z suffix)."""
    if dt.tzinfo is None:
        dt = dt.replace(tzinfo=timezone.utc)
    else:
        dt = dt.astimezone(timezone.utc)
    return dt.isoformat().replace("+00:00", "Z")


def _build_tracking_filename(filename: str) -> str:
    """Keep status keys unique so parallel imports do not overwrite each other."""
    normalized_name = str(filename or "").strip() or "remote_document"
    if normalized_name not in processing_status:
        return normalized_name

    stem = Path(normalized_name).stem or "remote_document"
    suffix = Path(normalized_name).suffix
    counter = 2
    while True:
        candidate = f"{stem} ({counter}){suffix}"
        if candidate not in processing_status:
            return candidate
        counter += 1


def _queue_processing_task(
    background_tasks: BackgroundTasks,
    temp_path: str,
    display_filename: str,
    message: str,
    source_info: Optional[dict] = None,
):
    processing_status[display_filename] = {
        "status": "processing",
        "progress": 10,
        "message": message,
    }
    save_status_to_db(display_filename, processing_status[display_filename])
    background_tasks.add_task(process_ocr_file, temp_path, display_filename, source_info)

# ============= AUDIT LOG ENDPOINTS =============
@router.post("/activity/log", response_model=SimpleStatusResponse)
async def log_user_activity(request: Request, action: str, claim_id: str = None, details: str = None):
    """Log a user activity (audit log)."""
    try:
        db = MongoDBManager()
        username = getattr(request.state, "username", None) or "unknown"
        db.log_activity(
            username=username,
            action=action,
            claim_id=claim_id,
            details=details,
            client_ip=getattr(request.state, "client_ip", None) or get_client_ip(request),
        )
        db.close()
        return {"status": "success"}
    except Exception as e:
        logger.warning(f"[ACTIVITY] Failed to log {action}: {e}")
        raise HTTPException(status_code=500, detail="Failed to log activity")

@router.get("/activity-log", response_model=ActivityLogResponse)
def get_activity_log(
    request: Request,  # NEW: RBAC
    days: int = Query(7, ge=1, le=90),
    limit: int = Query(200, ge=1, le=1000),
    offset: int = Query(0, ge=0),
):
    """Retrieve all activity logs from the last N days (Super Admin only)."""
    _require_permission(request, "view_activity_log", "Only the Super Admin can access activity logs.")
    
    try:
        db = MongoDBManager()
        raw_activities = db.get_activity_log(days=days, limit=limit + 1, offset=offset)
        has_more = len(raw_activities) > limit
        activities = raw_activities[:limit]
        # Convert timestamps to ISO format for JSON serialization
        for act in activities:
            if "timestamp" in act:
                act["timestamp"] = _to_utc_iso(act["timestamp"])
        db.close()
        return {
            "status": "success",
            "days": days,
            "limit": limit,
            "offset": offset,
            "count": len(activities),
            "has_more": has_more,
            "activities": activities,
        }
    except Exception as e:
        logger.error(f"[ERROR] Error retrieving activity log: {e}")
        raise HTTPException(status_code=500, detail="Error retrieving activity log")

@router.get("/activity-log/user/{username}", response_model=ActivityLogResponse)
def get_user_activity_log(
    request: Request,
    username: str,
    days: int = Query(7, ge=1, le=90),
    limit: int = Query(200, ge=1, le=1000),
    offset: int = Query(0, ge=0),
):
    """Retrieve activity logs for a specific user from the last N days."""
    _require_permission(request, "view_activity_log", "Only the Super Admin can access activity logs.")
    try:
        db = MongoDBManager()
        raw_activities = db.get_activity_log(
            days=days,
            username=username,
            limit=limit + 1,
            offset=offset,
        )
        has_more = len(raw_activities) > limit
        activities = raw_activities[:limit]
        # Convert timestamps to ISO format for JSON serialization
        for act in activities:
            if "timestamp" in act:
                act["timestamp"] = _to_utc_iso(act["timestamp"])
        db.close()
        return {
            "status": "success",
            "username": username,
            "days": days,
            "limit": limit,
            "offset": offset,
            "count": len(activities),
            "has_more": has_more,
            "activities": activities,
        }
    except Exception as e:
        logger.error(f"[ERROR] Error retrieving activity log for user '{username}': {e}")
        raise HTTPException(status_code=500, detail="Error retrieving user activity log")

# Store processing status in memory
# Note: In production, consider using Redis for better reliability across instances
processing_status = {}

# Helper function to persist status to MongoDB
def save_status_to_db(filename: str, status_data: dict):
    """Save processing status to MongoDB for persistence"""
    try:
        from db.mongo_db import MongoDBManager
        db = MongoDBManager()
        payload = _serialize_processing_status(status_data)
        db.db['processing_status'].update_one(
            {"filename": filename},
            {"$set": {**payload, "filename": filename, "updated_at": datetime.utcnow()}},
            upsert=True
        )
        db.close()
    except Exception as e:
        logger.warning(f"Could not persist status to DB: {e}")


def _update_processing_progress(filename: str, progress: int, message: str):
    processing_status[filename] = {
        "status": "processing",
        "progress": progress,
        "message": message,
    }
    save_status_to_db(filename, processing_status[filename])


def _sanitize_source_info(source_info: Optional[dict]) -> dict:
    if not isinstance(source_info, dict):
        return {}
    return {
        key: value
        for key, value in source_info.items()
        if not str(key).startswith("_") and key not in {"file_path", "preserve_file"}
    }


def _extract_text_from_document(
    file_path: str,
    display_name: str,
    status_filename: Optional[str] = None,
    generate_video_text: Optional[bool] = None,
):
    all_text = ""
    confidences = []
    linked_targets = []
    structured_content = None
    file_ext = Path(file_path).suffix.lower()

    if file_ext == ".pdf":
        if pdf_has_selectable_text(file_path):
            logger.info(f"Selectable text detected in PDF: {display_name}")
            all_text = extract_pdf_text(file_path)
            confidences = [100]
        else:
            doc = fitz.open(file_path)
            logger.info(f"PDF opened with {len(doc)} pages: {display_name}")
            total_pages = len(doc)
            try:
                for page_num, page in enumerate(doc, start=1):
                    logger.info(f"→ {display_name}: page {page_num}/{total_pages}")
                    try:
                        pix = page.get_pixmap(dpi=200)
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        result = paddle_ocr_extract(img)
                        if result:
                            all_text += f"\n\n--- PAGE {page_num} ---\n\n"
                            all_text += result["raw_text"]
                            confidences.append(result["confidence"])
                            logger.info(f"✓ {display_name} ({result['confidence']}%)")
                        else:
                            ocr_diag = get_last_ocr_diagnostics()
                            if ocr_diag:
                                logger.warning(
                                    f"OCR empty result on page {page_num} for {display_name} | "
                                    f"reason={ocr_diag.get('status')} detail={ocr_diag.get('detail')} "
                                    f"raw_candidates={ocr_diag.get('raw_candidates')} "
                                    f"kept_lines={ocr_diag.get('kept_lines')} "
                                    f"min_confidence={ocr_diag.get('min_confidence')}"
                                )
                            else:
                                logger.warning(f"No text on page {page_num} for {display_name}")
                    except Exception as page_err:
                        logger.error(f"Page extraction error for {display_name}: {page_err}")
                        continue

                    if status_filename:
                        _update_processing_progress(
                            status_filename,
                            min(45, 20 + int(25 * (page_num / max(total_pages, 1)))),
                            f"OCR in progress... ({page_num}/{total_pages})",
                        )
            finally:
                doc.close()

        linked_targets = extract_pdf_link_targets(file_path, allowed_extensions=_ALLOWED_EXTENSIONS)
    elif file_ext in [".docx", ".docm"]:
        logger.info(f"Word document detected: {display_name}")
        all_text = extract_word_text(file_path)
        confidences = [100]
    elif file_ext == ".doc":
        logger.info(f"Legacy Word document detected: {display_name}")
        try:
            all_text = extract_word_text(file_path)
            confidences = [100]
        except Exception as doc_err:
            logger.warning(f"Could not extract .doc with python-docx for {display_name}: {doc_err}")
            all_text = f"[Could not fully extract from .doc file: {display_name}]"
            confidences = [50]
    elif file_ext in [".xlsx", ".xlsm", ".xls", ".xlsb"]:
        logger.info(f"Excel file detected: {display_name}")
        try:
            structured_content = parse_excel_to_structured_content(file_path, display_name=display_name)
            all_text = str(structured_content.get("retrieval_text") or "").strip()
            confidences = [100]
        except Exception as excel_err:
            logger.error(f"Excel extraction error for {display_name}: {excel_err}")
            all_text = f"[Could not extract from Excel file: {display_name}]"
            confidences = [50]
    elif is_video_file(file_path):
        logger.info(f"Video file detected: {display_name}")
        should_generate_video_text = (
            settings.generate_video_transcripts
            if generate_video_text is None
            else _coerce_bool(generate_video_text, settings.generate_video_transcripts)
        )
        if should_generate_video_text:
            transcript = transcribe_video(file_path)
            all_text = str(transcript or "").strip()
            confidences = [100] if all_text else []
        else:
            all_text = f"[Video transcript generation is disabled for uploaded videos: {display_name}]"
            confidences = [100]
    else:
        logger.info(f"Image file detected: {display_name}")
        img = Image.open(file_path).convert("RGB")
        result = paddle_ocr_extract(img)
        if result:
            all_text = result["raw_text"]
            confidences = [result["confidence"]]
        else:
            ocr_diag = get_last_ocr_diagnostics()
            if ocr_diag:
                logger.error(
                    f"OCR failed for image {display_name} | "
                    f"reason={ocr_diag.get('status')} detail={ocr_diag.get('detail')} "
                    f"raw_candidates={ocr_diag.get('raw_candidates')} "
                    f"kept_lines={ocr_diag.get('kept_lines')} "
                    f"min_confidence={ocr_diag.get('min_confidence')}"
                )

    return all_text, confidences, linked_targets, structured_content


def _normalize_link_identity(value: str) -> str:
    return " ".join(str(value or "").strip().lower().split())


def _build_linked_documents_text(linked_documents: List[dict]) -> str:
    sections = []
    for index, linked_doc in enumerate(linked_documents, start=1):
        linked_text = str(linked_doc.get("text") or "").strip()
        if not linked_text:
            continue
        source_url = linked_doc.get("source_url") or linked_doc.get("origin_link") or "unknown"
        sections.append(
            "\n".join(
                [
                    f"--- LINKED DOCUMENT {index}: {linked_doc.get('filename', 'linked_document')} ---",
                    f"Source: {source_url}",
                    f"Linked from page: {linked_doc.get('page_number', 'unknown')}",
                    f"Depth: {linked_doc.get('depth', 1)}",
                    "",
                    linked_text,
                ]
            ).strip()
        )
    return "\n\n".join(section for section in sections if section)


def _build_llm_input_text(main_text: str, linked_documents: List[dict], max_chars: int) -> str:
    if max_chars <= 0:
        combined_sections = [str(main_text or "").strip()]
        linked_text = _build_linked_documents_text(linked_documents)
        if linked_text:
            combined_sections.append(linked_text)
        return "\n\n".join(section for section in combined_sections if section).strip()

    main_text = str(main_text or "").strip()
    if not linked_documents:
        return main_text[:max_chars]

    sections = []
    reserved_main = max(30000, int(max_chars * 0.5))
    main_excerpt = main_text[: min(len(main_text), reserved_main)]
    if main_excerpt:
        sections.append(f"### MAIN DOCUMENT ###\n{main_excerpt}".strip())

    remaining_chars = max_chars - len("\n\n".join(sections))
    if remaining_chars <= 0:
        return "\n\n".join(sections)[:max_chars].strip()

    per_link_budget = max(1500, remaining_chars // max(len(linked_documents), 1))
    for index, linked_doc in enumerate(linked_documents, start=1):
        current_text = "\n\n".join(sections)
        remaining_chars = max_chars - len(current_text)
        if remaining_chars <= 300:
            break

        header = (
            f"### LINKED DOCUMENT {index}: {linked_doc.get('filename', 'linked_document')} ###\n"
            f"Source: {linked_doc.get('source_url') or linked_doc.get('origin_link') or 'unknown'}\n"
        )
        linked_body_budget = max(300, min(per_link_budget, remaining_chars) - len(header) - 2)
        linked_excerpt = str(linked_doc.get("text") or "").strip()[:linked_body_budget]
        if not linked_excerpt:
            continue
        sections.append(f"{header}{linked_excerpt}".strip())

    return "\n\n".join(section for section in sections if section)[:max_chars].strip()


EXCEL_FILL_SYSTEM_PROMPT = """
You fill an Excel template using structured data extracted from documents.

Return ONLY valid JSON in this schema:
{
  "cells": [
    {"row": 1, "col": 1, "value": "text"}
  ]
}

Rules:
1. Only fill cells that correspond to the provided targets (form_fields + tables).
2. If a value is missing, omit the cell from the output (leave it blank).
3. Do NOT invent data. Use only what is present in the documents.
4. Preserve original formatting (dates, numbers, IDs) as shown in the data.
5. Do NOT overwrite existing values; only fill empty targets.
"""


def _normalize_excel_header(value, index: int, seen: dict) -> str:
    raw = str(value or "").strip()
    name = raw if raw else f"column_{index + 1}"
    count = seen.get(name, 0) + 1
    seen[name] = count
    if count == 1:
        return name
    return f"{name}_{count}"


def _load_excel_template(
    excel_path: str,
    sheet_name: Optional[str] = None,
    keep_vba: bool = False,
):
    wb = load_workbook(excel_path, keep_vba=keep_vba)
    if sheet_name and sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
    else:
        sheet = wb.active

    header_row_idx = None
    headers = []
    for row in sheet.iter_rows(min_row=1, max_row=25):
        values = [cell.value for cell in row]
        if not any(v is not None and str(v).strip() for v in values):
            continue
        last_non_empty = max(
            idx for idx, v in enumerate(values) if v is not None and str(v).strip()
        )
        header_cells = values[: last_non_empty + 1]
        seen = {}
        headers = [
            _normalize_excel_header(value, index, seen)
            for index, value in enumerate(header_cells)
        ]
        header_row_idx = row[0].row
        break

    if not headers or header_row_idx is None:
        raise ValueError("No header row found in the Excel template.")

    return wb, sheet, headers, header_row_idx


def _coerce_excel_rows(headers: list[str], rows: list) -> list[dict]:
    normalized = []
    header_map = {str(header).strip().lower(): header for header in headers}
    for row in rows or []:
        payload = {}
        if isinstance(row, dict):
            if "cells" in row and isinstance(row.get("cells"), list):
                for cell in row.get("cells") or []:
                    if not isinstance(cell, dict):
                        continue
                    col = str(cell.get("column") or "").strip()
                    if not col:
                        continue
                    header_key = header_map.get(col.lower())
                    if not header_key:
                        continue
                    value = cell.get("value", "")
                    if value is None:
                        value = ""
                    if isinstance(value, (dict, list)):
                        value = json.dumps(value, ensure_ascii=False)
                    payload[header_key] = str(value)
            else:
                for header in headers:
                    value = row.get(header, "")
                    if value is None:
                        value = ""
                    if isinstance(value, (dict, list)):
                        value = json.dumps(value, ensure_ascii=False)
                    payload[header] = str(value)
        if not payload:
            continue
        for header in headers:
            payload.setdefault(header, "")
        normalized.append(payload)
    return normalized


def _cell_text(value) -> str:
    if value is None:
        return ""
    text = str(value).strip()
    return text


def _is_empty_cell(value) -> bool:
    return not _cell_text(value)


def _normalize_merged_target(sheet, row: int, col: int) -> tuple[int, int]:
    for merged_range in sheet.merged_cells.ranges:
        if merged_range.min_row <= row <= merged_range.max_row and merged_range.min_col <= col <= merged_range.max_col:
            return merged_range.min_row, merged_range.min_col
    return row, col


def _collect_form_fields(sheet, max_rows: int = 80, max_cols: int = 8) -> list[dict]:
    fields = []
    label_keywords = {
        "reviewer", "approver", "approved date", "reviewed date", "effective date",
        "tender", "bid", "customer", "contract", "amount", "date", "name",
    }
    for row in sheet.iter_rows(min_row=1, max_row=min(sheet.max_row, max_rows)):
        for col_idx, cell in enumerate(row[:max_cols], start=1):
            label = _cell_text(cell.value)
            if not label:
                continue
            if len(label) > 80:
                continue
            label_norm = re.sub(r"[^a-z0-9 ]+", "", label.lower()).strip()
            if not (label.endswith(":") or any(key in label_norm for key in label_keywords)):
                continue
            target_cell = sheet.cell(row=cell.row, column=col_idx + 1)
            if _is_empty_cell(target_cell.value):
                fields.append(
                    {
                        "row": cell.row,
                        "col": col_idx + 1,
                        "label": label,
                    }
                )
            if len(fields) >= 120:
                return fields
    return fields


def _looks_like_header(values: list[str]) -> bool:
    if len(values) < 3:
        return False
    normalized = " ".join(values).lower()
    header_tokens = [
        "s. no", "sr. no", "s no", "name", "date", "status", "remarks", "description",
        "qty", "quantity", "amount", "price", "total", "designation", "profile",
        "qualification", "unit rate", "approval", "department", "stage", "evidence",
        "responsibilities", "heading", "sub heading",
    ]
    return any(token in normalized for token in header_tokens)


def _detect_tables(sheet, max_scan_rows: int = 60, max_rows_per_table: int = 20) -> list[dict]:
    tables = []
    for row in sheet.iter_rows(min_row=1, max_row=min(sheet.max_row, max_scan_rows)):
        values = []
        col_map = []
        for cell in row:
            text = _cell_text(cell.value)
            if text:
                values.append(text)
                col_map.append((cell.column, text))
        if not _looks_like_header(values):
            continue
        columns = [{"col": col, "name": name} for col, name in col_map]
        header_row = row[0].row
        rows = []
        empty_slots_added = 0
        for r in range(header_row + 1, min(sheet.max_row, header_row + max_rows_per_table) + 1):
            existing = {}
            row_has_data = False
            for col in columns:
                cell = sheet.cell(row=r, column=col["col"])
                text = _cell_text(cell.value)
                if text:
                    row_has_data = True
                    existing[col["name"]] = text
            if row_has_data:
                rows.append({"row": r, "existing": existing})
            else:
                if empty_slots_added < 8:
                    rows.append({"row": r, "existing": {}})
                    empty_slots_added += 1
        if rows:
            tables.append(
                {
                    "header_row": header_row,
                    "columns": columns,
                    "rows": rows,
                }
            )
    return tables


def _chunk_list(items: list, size: int) -> list[list]:
    return [items[i:i + size] for i in range(0, len(items), size)]


def _build_target_batches(form_fields: list[dict], tables: list[dict]) -> list[dict]:
    batches = []
    for chunk in _chunk_list(form_fields or [], 30):
        batches.append({"form_fields": chunk, "tables": []})

    for table in tables or []:
        rows = table.get("rows") or []
        for chunk in _chunk_list(rows, 8):
            batches.append(
                {
                    "form_fields": [],
                    "tables": [
                        {
                            "header_row": table.get("header_row"),
                            "columns": table.get("columns") or [],
                            "rows": chunk,
                        }
                    ],
                }
            )
    return batches


def _fill_excel_sheet_with_gemini(sheet_name: str, targets: dict, documents: list[dict], prompt: Optional[str]):
    if not gemini_is_configured():
        raise RuntimeError("Gemini client is not configured. Set GEMINI_API_KEY.")

    response_schema = {
        "type": "object",
        "properties": {
            "cells": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "row": {"type": "integer"},
                        "col": {"type": "integer"},
                        "value": {"type": "string"},
                    },
                    "required": ["row", "col", "value"],
                },
            },
        },
        "required": ["cells"],
    }

    def _repair_and_parse(raw: str) -> dict:
        try:
            parsed = safe_json_parse(raw)
            if isinstance(parsed, dict) and isinstance(parsed.get("value"), list):
                return {"cells": parsed["value"]}
            return parsed
        except Exception:
            pass

        text = str(raw or "").strip()
        # Try to extract the largest JSON object
        obj_start = text.find("{")
        obj_end = text.rfind("}")
        if obj_start != -1 and obj_end != -1 and obj_end > obj_start:
            try:
                parsed = safe_json_parse(text[obj_start:obj_end + 1])
                if isinstance(parsed, dict) and isinstance(parsed.get("value"), list):
                    return {"cells": parsed["value"]}
                return parsed
            except Exception:
                pass

        # Try to extract a JSON array and wrap it as rows
        arr_start = text.find("[")
        arr_end = text.rfind("]")
        if arr_start != -1 and arr_end != -1 and arr_end > arr_start:
            try:
                parsed = safe_json_parse(text[arr_start:arr_end + 1])
                if isinstance(parsed, dict) and isinstance(parsed.get("value"), list):
                    return {"cells": parsed["value"]}
                if isinstance(parsed, list):
                    return {"cells": parsed}
                if isinstance(parsed, dict) and "cells" in parsed:
                    return parsed
            except Exception:
                pass

        raise ValueError("Cannot parse JSON response after repair attempts.")

    user_prompt = (
        f"Sheet: {sheet_name}\n\n"
        f"Targets (form_fields + tables): {json.dumps(targets, ensure_ascii=False)}\n\n"
        f"Documents (structured): {json.dumps(documents, ensure_ascii=False)}\n\n"
        f"User request: {str(prompt or '').strip() or 'Fill the Excel based on the documents.'}\n\n"
        "Return ONLY valid JSON."
    )
    max_retries = 3
    last_error = None

    for attempt in range(1, max_retries + 1):
        content = ""
        try:
            logger.info(
                f"[Excel Fill] Sending request to Gemini model {gemini_model_name()} (attempt {attempt}/{max_retries})"
            )
            retry_note = ""
            if attempt > 1:
                retry_note = (
                    "\n\nIMPORTANT RETRY NOTE:\n"
                    "The previous response was not valid JSON. "
                    "Return exactly ONE JSON object with a top-level 'rows' array. "
                    "No extra text or markdown."
                )
            content = gemini_generate_json_text(
                user_prompt=user_prompt + retry_note,
                system_instruction=EXCEL_FILL_SYSTEM_PROMPT,
                temperature=0.1,
                max_output_tokens=1800,
                response_schema=response_schema,
            )
            result = _repair_and_parse(content)
            cells = result.get("cells")
            if not isinstance(cells, list):
                raise ValueError("Gemini response did not include a 'cells' list.")
            return cells
        except Exception as exc:
            last_error = exc
            snippet = str(content or "")
            if snippet:
                logger.warning(
                    f"[Excel Fill] Raw Gemini response (attempt {attempt}) preview: {snippet[:2000]}"
                )
            logger.warning(f"[Excel Fill] Attempt {attempt} failed: {str(exc)[:200]}")
            time.sleep(0.6 * attempt)

    raise last_error


def _prepare_linked_document_metadata(linked_documents: List[dict]) -> List[dict]:
    prepared = []
    for linked_doc in linked_documents:
        prepared.append(
            {
                "filename": linked_doc.get("filename"),
                "source_url": linked_doc.get("source_url"),
                "origin_link": linked_doc.get("origin_link"),
                "page_number": linked_doc.get("page_number"),
                "depth": linked_doc.get("depth"),
                "char_count": linked_doc.get("char_count"),
                "confidence_percent": linked_doc.get("confidence_percent"),
            }
        )
    return prepared


def _first_non_empty_lines(text: str, max_lines: int = 6) -> List[str]:
    lines = []
    for raw_line in str(text or "").splitlines():
        clean_line = " ".join(raw_line.strip().split())
        if not clean_line:
            continue
        lines.append(clean_line)
        if len(lines) >= max_lines:
            break
    return lines


def _build_linked_document_preview_entry(linked_doc: dict) -> dict:
    lines = _first_non_empty_lines(linked_doc.get("text", ""), max_lines=8)
    title = " | ".join(lines[:2]) if lines else linked_doc.get("filename", "linked_document")
    preview = " ".join(lines[:6])[:_LINKED_PDF_PREVIEW_CHARS]
    return {
        "filename": linked_doc.get("filename"),
        "source_url": linked_doc.get("source_url"),
        "linked_from_page": linked_doc.get("page_number"),
        "depth": linked_doc.get("depth"),
        "char_count": linked_doc.get("char_count"),
        "confidence_percent": linked_doc.get("confidence_percent"),
        "document_type": "linked_pdf",
        "title": title[:240],
        "preview_excerpt": preview,
    }


def _build_structured_linked_document_entries(linked_documents: List[dict]):
    return [_build_linked_document_preview_entry(linked_doc) for linked_doc in linked_documents], []


def _resolve_link_target_files(link_target, graph_access_token: Optional[str] = None):
    warnings = []

    if getattr(link_target, "is_remote", False):
        remote_files, import_warnings = import_documents_from_link(
            link_target.target,
            _LINKED_PDF_ALLOWED_EXTENSIONS,
            graph_access_token=graph_access_token,
        )
        warnings.extend(import_warnings)
        truncated = len(remote_files) > _LINKED_PDF_MAX_DOCUMENTS_PER_LINK
        selected_files = remote_files[:_LINKED_PDF_MAX_DOCUMENTS_PER_LINK]
        if truncated:
            warnings.append(
                f"{link_target.target}: only the first {_LINKED_PDF_MAX_DOCUMENTS_PER_LINK} linked documents were processed."
            )
        return [
            {
                "filename": remote_file.filename,
                "path": remote_file.temp_path,
                "source_url": remote_file.source_url,
                "size_bytes": remote_file.size_bytes,
                "cleanup": True,
            }
            for remote_file in selected_files
        ], warnings

    local_path = Path(link_target.target)
    if not local_path.exists():
        return [], [f"{local_path.name or link_target.target}: linked file was not found."]

    if local_path.suffix.lower() not in _LINKED_PDF_ALLOWED_EXTENSIONS:
        return [], [f"{local_path.name}: linked file type '{local_path.suffix.lower()}' is not supported."]

    return [
        {
            "filename": local_path.name,
            "path": str(local_path),
            "source_url": str(local_path.resolve()),
            "size_bytes": local_path.stat().st_size,
            "cleanup": False,
        }
    ], warnings


def _collect_linked_document_texts(
    link_targets: List,
    graph_access_token: Optional[str] = None,
    depth: int = 1,
    visited_targets: Optional[set] = None,
    visited_documents: Optional[set] = None,
    remaining_slots: Optional[int] = None,
):
    visited_targets = visited_targets or set()
    visited_documents = visited_documents or set()
    remaining_slots = _LINKED_PDF_MAX_DOCUMENTS if remaining_slots is None else remaining_slots

    if depth > _LINKED_PDF_MAX_DEPTH or remaining_slots <= 0:
        return [], [], remaining_slots

    collected = []
    warnings = []
    limit_warning_added = False

    for link_target in link_targets:
        if remaining_slots <= 0:
            if not limit_warning_added:
                warnings.append(
                    f"Linked document limit reached ({_LINKED_PDF_MAX_DOCUMENTS}); remaining PDF links were skipped."
                )
                limit_warning_added = True
            break

        target_key = _normalize_link_identity(getattr(link_target, "target", ""))
        if not target_key or target_key in visited_targets:
            continue
        visited_targets.add(target_key)

        try:
            resolved_files, resolve_warnings = _resolve_link_target_files(
                link_target,
                graph_access_token=graph_access_token,
            )
            warnings.extend(resolve_warnings)
        except RemoteImportError as exc:
            warnings.append(f"{getattr(link_target, 'target', 'linked document')}: {exc}")
            continue
        except Exception as exc:
            logger.warning(f"Linked document resolution failed for {getattr(link_target, 'target', '')}: {exc}")
            warnings.append(f"{getattr(link_target, 'target', 'linked document')}: resolution failed")
            continue

        for resolved_file in resolved_files:
            if remaining_slots <= 0:
                break

            document_key = _normalize_link_identity(resolved_file.get("source_url") or resolved_file.get("path"))
            if not document_key or document_key in visited_documents:
                if resolved_file.get("cleanup") and os.path.exists(resolved_file["path"]):
                    os.remove(resolved_file["path"])
                continue
            visited_documents.add(document_key)

            try:
                linked_text, linked_confidences, nested_targets, _ = _extract_text_from_document(
                    resolved_file["path"],
                    resolved_file["filename"],
                )
                if not linked_text.strip():
                    warnings.append(f"{resolved_file['filename']}: no text extracted from linked document.")
                    continue

                linked_confidence = (
                    round(sum(linked_confidences) / len(linked_confidences), 2)
                    if linked_confidences else 0.0
                )
                collected.append(
                    {
                        "filename": resolved_file["filename"],
                        "source_url": resolved_file.get("source_url"),
                        "origin_link": getattr(link_target, "target", ""),
                        "page_number": getattr(link_target, "page_number", None),
                        "depth": depth,
                        "char_count": len(linked_text),
                        "confidence_percent": linked_confidence,
                        "text": linked_text,
                    }
                )
                remaining_slots -= 1

                if nested_targets and depth < _LINKED_PDF_MAX_DEPTH and remaining_slots > 0:
                    nested_docs, nested_warnings, remaining_slots = _collect_linked_document_texts(
                        nested_targets,
                        graph_access_token=graph_access_token,
                        depth=depth + 1,
                        visited_targets=visited_targets,
                        visited_documents=visited_documents,
                        remaining_slots=remaining_slots,
                    )
                    collected.extend(nested_docs)
                    warnings.extend(nested_warnings)
            except Exception as exc:
                logger.warning(f"Linked document extraction failed for {resolved_file['filename']}: {exc}")
                warnings.append(f"{resolved_file['filename']}: extraction failed ({exc})")
            finally:
                if resolved_file.get("cleanup") and os.path.exists(resolved_file["path"]):
                    os.remove(resolved_file["path"])

    return collected, warnings, remaining_slots

# ============= DOCUMENT RETRIEVAL ENDPOINTS =============

@router.get("/documents", response_model=DocumentsListResponse)
def get_documents(
    request: Request,  # NEW: RBAC
    limit: int = Query(5, ge=1, le=100),
    offset: int = Query(0, ge=0)
):
    """Get all documents with pagination - default 5 per page (Admin+ only)"""
    _require_permission(request, "view_documents", "Admins and Super Admins only can access documents.")

    result = db_service.get_all_documents(limit, offset)
    return {
        "total": result["total"],
        "documents": result["documents"],
        "offset": result.get("offset", offset),
        "limit": result.get("limit", limit),
        "pages": (result["total"] + limit - 1) // limit if result["total"] > 0 else 1
    }


@router.get("/documents/type/{doc_type}", response_model=DocumentsListResponse)
def get_documents_by_type(
    request: Request,
    doc_type: str,
    limit: int = Query(10, ge=1, le=100)
):
    """Get documents filtered by type"""
    _require_permission(request, "view_documents", "Admins and Super Admins only can access documents.")
    result = db_service.get_documents_by_type(doc_type, limit)
    return DocumentsListResponse(
        total=result["total"],
        documents=result["documents"]
    )


@router.get("/documents/duplicates")
@router.get("/documents/duplicate-analysis")
def detect_duplicate_documents(
    request: Request,
    limit: int = Query(40, ge=10, le=120),
    similarity_threshold: float = Query(0.86, ge=0.7, le=0.99),
    query: Optional[str] = Query(None),
    document_type: Optional[str] = Query(None),
):
    """Detect exact duplicates, similar versions, and fraud-risk copies."""
    _require_permission(request, "view_documents", "Admins and Super Admins only can access documents.")
    result = db_service.detect_duplicates(
        limit=limit,
        similarity_threshold=similarity_threshold,
        search_term=(query or "").strip() or None,
        document_type=(document_type or "").strip() or None,
    )
    result["search_query"] = (query or "").strip()
    result["document_type_filter"] = (document_type or "").strip()
    return result


@router.get("/documents/{doc_id}", response_model=DocumentDetailResponse)
def get_document(request: Request, doc_id: str):
    """Get single document by ID with full details"""
    return _get_document_for_read(request, doc_id)


@router.get("/documents/{doc_id}/original")
def open_original_document(request: Request, doc_id: str):
    """Open the original stored file or source URL for a document."""
    doc = _get_document_for_read(request, doc_id)

    target_type, target_value = _resolve_original_document_target(doc)
    if target_type == "file" and target_value:
        media_type, _ = mimetypes.guess_type(target_value)
        download_name = str(doc.get("file_name") or Path(target_value).name or "document")
        return FileResponse(
            path=target_value,
            media_type=media_type or "application/octet-stream",
            filename=download_name,
            content_disposition_type="inline",
        )

    if target_type == "url" and target_value:
        return RedirectResponse(url=target_value, status_code=307)

    raise HTTPException(
        status_code=404,
        detail="Original document is not available for this record",
    )


@router.get("/documents/{doc_id}/original-preview")
def preview_original_document(request: Request, doc_id: str):
    """Render an inline-friendly preview page for the original document."""
    doc = _get_document_for_read(request, doc_id)
    return HTMLResponse(
        content=_build_original_preview_html(
            doc_id,
            doc,
            base_path=str(request.scope.get("root_path") or "").rstrip("/"),
        )
    )


@router.get("/statistics", response_model=StatisticsResponse)
def get_statistics(request: Request):
    """Get database statistics: total documents and breakdown by type"""
    _require_permission(request, "view_statistics", "Admins and Super Admins only can access dashboard statistics.")
    stats = db_service.get_statistics()
    return StatisticsResponse(
        total_documents=stats.get("total_documents", 0),
        document_types=stats.get("document_types", []),
        average_confidence=stats.get("average_confidence", 0.0)
    )


@router.get("/dashboard-insights")
def get_dashboard_insights(request: Request):
    """Get dashboard manager insights for employee creation and upload usage charts."""
    _require_permission(request, "view_statistics", "Admins and Super Admins only can access dashboard insights.")
    return db_service.get_dashboard_insights()


@router.get("/employees", response_model=EmployeeProfilesListResponse)
def list_employee_profiles(
    request: Request,  # NEW: RBAC
    query: Optional[str] = Query(None),
    limit: int = Query(20, ge=1, le=100),
):
    """Get employee profiles visible to the current role."""
    user_role = _get_request_role(request)
    actor = _get_request_actor(request)
    if user_role == UserRole.MANAGER:
        result = db_service.search_employee_profiles(
            search_term=(query or "").strip() or None,
            limit=limit,
            created_by=actor,
        )
    elif user_role in {UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        result = db_service.search_employee_profiles(search_term=(query or "").strip() or None, limit=limit)
    else:
        raise HTTPException(status_code=403, detail="Employees cannot access employee profiles.")
    return EmployeeProfilesListResponse(
        total=result["total"],
        employees=result["employees"],
    )


@router.post("/employees", response_model=EmployeeProfileSaveResponse)
def create_employee_profile(request: Request, payload: EmployeeProfileRequest):
    user_role, existing_profile = _require_employee_profile_access(
        request,
        emp_id=(payload.empID or "").strip(),
        write=True,
    )
    if user_role not in {UserRole.MANAGER, UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        raise HTTPException(status_code=403, detail="You are not allowed to create employee profiles.")

    extra = {}
    if existing_profile and existing_profile.get("created_by"):
        extra["created_by"] = existing_profile.get("created_by")
    elif user_role == UserRole.MANAGER:
        extra["created_by"] = _get_request_actor(request)
    profile = db_service.create_or_update_employee_profile(
        emp_id=payload.empID,
        emp_name=payload.empName,
        extra=extra,
    )
    if not profile:
        raise HTTPException(status_code=400, detail="empID and empName are required")
    _invalidate_qa_caches()
    return {
        "status": "saved",
        "employee": {
            "empID": profile.get("empID"),
            "empName": profile.get("empName"),
            "uuid": profile.get("uuid"),
        },
    }


@router.delete("/employees/{emp_id}", response_model=EmployeeProfileDeleteResponse)
def delete_employee_profile(request: Request, emp_id: str):
    user_role, profile = _require_employee_profile_access(request, emp_id=emp_id, write=True)
    if user_role not in {UserRole.MANAGER, UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        raise HTTPException(status_code=403, detail="You are not allowed to delete employee profiles.")
    if not profile:
        raise HTTPException(status_code=404, detail="Employee profile not found")
    result = db_service.delete_employee_profile(emp_id=(emp_id or "").strip() or None)
    if not result:
        raise HTTPException(status_code=404, detail="Employee profile not found")
    _invalidate_qa_caches()
    return result


@router.post("/employees/bulk-delete", response_model=BulkDeleteEmployeeProfilesResponse)
def bulk_delete_employee_profiles(request: Request, payload: BulkDeleteEmployeeProfilesRequest):
    emp_ids = []
    seen = set()
    for raw_emp_id in payload.emp_ids:
        emp_id = str(raw_emp_id or "").strip()
        if not emp_id or emp_id in seen:
            continue
        seen.add(emp_id)
        emp_ids.append(emp_id)

    if not emp_ids:
        raise HTTPException(status_code=400, detail="At least one valid EMPID is required")

    user_role = _get_request_role(request)
    if user_role == UserRole.MANAGER:
        actor = _get_request_actor(request)
        for emp_id in emp_ids:
            _, profile = _require_employee_profile_access(request, emp_id=emp_id, write=True)
            if profile:
                owner = str(profile.get("created_by") or "").strip().lower()
                if owner and owner != actor.lower():
                    raise HTTPException(
                        status_code=403,
                        detail="Managers can only delete employee profiles they created.",
                    )
    elif user_role not in {UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        raise HTTPException(status_code=403, detail="You are not allowed to delete employee profiles.")

    result = db_service.delete_employee_profiles(emp_ids)
    if result.get("deleted_count", 0):
        _invalidate_qa_caches()
    return {
        "status": "completed",
        "requested_count": result["requested_count"],
        "deleted_count": result["deleted_count"],
        "deleted_ids": result["deleted_ids"],
        "not_found_ids": result["not_found_ids"],
        "deleted_documents": result["deleted_documents"],
        "deleted_chunks": result["deleted_chunks"],
    }


@router.get("/employee-documents", response_model=DocumentsListResponse)
def get_employee_documents(
    request: Request,
    empID: Optional[str] = Query(None),
    uuid: Optional[str] = Query(None),
    limit: int = Query(50, ge=1, le=200),
):
    if not (empID or uuid):
        raise HTTPException(status_code=400, detail="empID or uuid is required")
    user_role, profile = _require_employee_profile_access(request, emp_id=empID, profile_uuid=uuid, write=False)
    if user_role == UserRole.MANAGER and not profile:
        raise HTTPException(status_code=404, detail="Employee profile not found")
    result = db_service.get_documents_by_employee(
        emp_id=(empID or "").strip() or None,
        profile_uuid=(uuid or "").strip() or None,
        limit=limit,
        access_context=_build_request_kbac_context(request),
    )
    return DocumentsListResponse(total=result["total"], documents=result["documents"])


@router.get("/search", response_model=DocumentsListResponse)
def search_documents(
    request: Request,
    query: Optional[str] = Query(None),
    document_type: Optional[str] = Query(None),
    vendor_name: Optional[str] = Query(None),
    limit: int = Query(5, ge=1, le=100),
    offset: int = Query(0, ge=0)
):
    """Search documents by free text and optional type filter with pagination."""
    _require_permission(request, "view_documents", "Admins and Super Admins only can search documents.")
    # Backward compatibility: if old clients still pass vendor_name only,
    # treat it as a free-text search term.
    search_term = (query or vendor_name or "").strip() or None

    result = db_service.search_documents(
        search_term=search_term,
        document_type=document_type,
        limit=limit,
        offset=offset,
    )
    return {
        "total": result["total"],
        "documents": result["documents"],
        "offset": result.get("offset", offset),
        "limit": result.get("limit", limit),
        "pages": (result["total"] + limit - 1) // limit if result["total"] > 0 else 1
    }


@router.get("/search/semantic")
def semantic_search_documents(
    request: Request,
    query: str = Query(..., min_length=2),
    document_type: Optional[str] = Query(None),
    limit: int = Query(20, ge=1, le=100)
):
    """Semantic search with natural language query parsing and ranking."""
    _require_permission(request, "view_documents", "Admins and Super Admins only can search documents.")
    result = db_service.semantic_search(
        query=query,
        document_type=document_type,
        limit=limit,
    )
    return DocumentsListResponse(
        total=result["total"],
        documents=result["documents"]
    )


@router.get("/documents/{doc_id}/chunks")
def get_document_chunks(
    request: Request,
    doc_id: str,
    query: Optional[str] = Query(None, alias="q"),
    limit: int = Query(20, ge=1, le=200),
    offset: int = Query(0, ge=0),
    highlight: bool = Query(True),
    max_chars: int = Query(1200, ge=200, le=5000),
):
    """Retrieve stored chunks for a document with optional search + highlighting."""
    _require_permission(request, "view_document_chunks", "Admins and Super Admins only can view document chunks.")
    search_term = (query or "").strip()
    if search_term:
        result = db_service.search_document_chunks(
            query=search_term,
            doc_id=doc_id,
            limit=limit,
            offset=offset,
        )
    else:
        result = db_service.get_document_chunks(doc_id=doc_id, limit=limit, offset=offset)

    chunks = result.get("chunks", [])
    terms = _extract_highlight_terms(search_term) if highlight and search_term else []

    for chunk in chunks:
        text = str(chunk.get("chunk_text") or "")
        if max_chars and len(text) > max_chars:
            text = text[:max_chars].rstrip() + "..."
        chunk["chunk_text"] = text
        if terms:
            chunk["highlighted_text"] = _highlight_text(text, terms)

    total = int(result.get("total", 0))
    return {
        "total": total,
        "offset": offset,
        "limit": limit,
        "pages": (total + limit - 1) // limit if total > 0 else 1,
        "query": search_term or None,
        "chunks": chunks,
    }


@router.post("/employee-manifests/refresh")
def refresh_employee_manifests(request: Request, payload: RefreshManifestRequest):
    """Rebuild employee manifests (optionally limited for safety)."""
    _require_permission(request, "view_dashboard", "Admins and Super Admins only can refresh employee manifests.")
    limit = payload.limit
    result = db_service.refresh_employee_manifests(limit=limit)
    if result.get("refreshed", 0):
        _invalidate_qa_caches()
    return {
        "status": "success",
        "refreshed": result.get("refreshed", 0),
        "limit": limit,
    }


@router.get("/qa/history", response_model=QAHistoryListResponse)
def qa_history_list(request: Request):
    user_session = _get_session_data_safe(request)
    user_role = get_session_role(user_session) if user_session else None
    if not user_role or not check_permission(user_session, "view_ai_buzz"):
        raise HTTPException(status_code=403, detail="Not authorized to access Q&A history")

    user_key = _qa_history_user_key(request)
    conversations = _qa_history_list_conversations(user_key)
    return {
        "status": "success",
        "ttl_hours": max(1, int(_QA_HISTORY_TTL / 3600)),
        "conversations": conversations,
    }


@router.get("/qa/history/{conversation_id}", response_model=QAHistoryDetailResponse)
def qa_history_detail(request: Request, conversation_id: str):
    user_session = _get_session_data_safe(request)
    user_role = get_session_role(user_session) if user_session else None
    if not user_role or not check_permission(user_session, "view_ai_buzz"):
        raise HTTPException(status_code=403, detail="Not authorized to access Q&A history")

    user_key = _qa_history_user_key(request)
    conversation = _qa_history_get_conversation(user_key, conversation_id)
    if not isinstance(conversation, dict):
        raise HTTPException(status_code=404, detail="Q&A history not found")

    turns = conversation.get("turns") if isinstance(conversation.get("turns"), list) else []
    filtered_turns = [
        _qa_history_filter_turn_for_role(user_role, turn)
        for turn in turns
        if isinstance(turn, dict)
    ]
    filtered_turns = _qa_history_filter_live_graphs(_get_request_session_id(request), filtered_turns)
    return {
        "status": "success",
        "conversation_id": str(conversation.get("conversation_id") or conversation_id),
        "title": str(conversation.get("title") or "").strip(),
        "updated_at": conversation.get("updated_at"),
        "turns": filtered_turns,
    }


@router.post("/qa", response_model=QAResponse)
async def question_answer(request: Request, payload: QARequest):  # NEW: Add request for RBAC
    """Answer a question using stored documents with citations (Employee+ only)."""
    user_session = _get_session_data_safe(request)
    user_role = get_session_role(user_session) if user_session else None
    if not user_role or not check_permission(user_session, "view_ai_buzz"):
        raise HTTPException(status_code=403, detail="Not authorized to access Q&A")

    question = str(payload.question or "").strip()
    if not question:
        raise HTTPException(status_code=400, detail="Question is required")
    if len(question) == 1:
        history_user_key = _qa_history_user_key(request)
        requested_conversation_id = str(payload.conversation_id or "").strip() or str(uuid4())
        response_payload = _qa_incomplete_input_payload(requested_conversation_id)
        stored_conversation_id = _qa_history_store_turn(
            history_user_key,
            requested_conversation_id,
            question,
            response_payload,
            topic="",
        )
        response_payload["conversation_id"] = stored_conversation_id
        return response_payload

    try:
        limit = int(payload.limit)
    except Exception:
        limit = 6

    use_chunks = bool(payload.use_chunks)
    try:
        raw_chunk_limit = payload.chunk_limit
        chunk_limit = int(raw_chunk_limit if raw_chunk_limit is not None else max(limit * 2, 6))
    except Exception:
        chunk_limit = max(limit * 2, 6)

    simple_lookup = any(
        token in question.lower()
        for token in [
            "email", "mail", "phone", "mobile", "contact", "address",
            "gst", "gstin", "pan", "invoice", "bill number", "document number",
            "amount", "date", "number",
        ]
    )
    broad_document_listing = _qa_requests_broad_document_listing(question)
    if broad_document_listing:
        limit = 12
        chunk_limit = max(8, min(chunk_limit, 24))
    else:
        limit = max(1, min(limit, 3 if simple_lookup else 5))
        chunk_limit = max(1, min(chunk_limit, 12))

    logger.info(f"[QA API] Question: {question[:100]} (requesting {limit} sources) by role: {user_role}")

    access_context = _build_request_kbac_context(request)
    session_id = _get_request_session_id(request)
    _qa_graph_rollover_for_question(session_id)
    history_user_key = _qa_history_user_key(request)
    requested_conversation_id = str(payload.conversation_id or "").strip() or str(uuid4())

    qa_history = _qa_history_get_turns(history_user_key, requested_conversation_id)
    if not qa_history:
        qa_history = _qa_conversation_get(session_id)

    conversation_context = _qa_build_conversation_context(qa_history)
    retrieval_question = _qa_build_contextual_question(question, qa_history)
    contextual_sources = _qa_collect_context_sources(question, qa_history)
    context_signature = _qa_build_context_signature(qa_history)
    cache_scope_key = f"{access_context.get('scope_key', 'all')}|ctx:{context_signature}"
    cache_key = _qa_cache_key(question, limit, chunk_limit, use_chunks, cache_scope_key)

    cached_response = _qa_cache_get(cache_key)
    if cached_response is not None:
        logger.info("[QA API] Full response cache hit")
        cached_payload = dict(cached_response)
        cached_graph_data, cached_chart_options, _ = _qa_finalize_graph_payload(
            question,
            cached_payload.get("graph_data"),
            cached_payload.get("chart_options"),
            cached_payload.get("sources") if isinstance(cached_payload.get("sources"), list) else [],
        )
        cached_payload["graph_data"] = cached_graph_data
        cached_payload["chart_options"] = cached_chart_options
        cached_payload = _qa_graph_attach_to_response(session_id, cached_payload)
        cached_followups = cached_payload.get("follow_up_suggestions")
        if not isinstance(cached_followups, list) or not cached_followups:
            cached_payload["follow_up_suggestions"] = _qa_build_follow_up_suggestions(
                question=question,
                answer=str(cached_payload.get("answer") or ""),
                sources=cached_payload.get("sources") if isinstance(cached_payload.get("sources"), list) else [],
                has_graph_data=bool(cached_payload.get("graph_id")),
            )
        _qa_conversation_append(
            session_id,
            {
                "question": question,
                "answer": cached_payload.get("answer", ""),
                "topic": "",
                "sources": _build_answer_sources(
                    cached_payload.get("sources") if isinstance(cached_payload.get("sources"), list) else [],
                    cached_payload.get("chunk_sources") if isinstance(cached_payload.get("chunk_sources"), list) else [],
                )[:6],
            },
        )
        cached_conversation_id = _qa_history_store_turn(
            history_user_key,
            requested_conversation_id,
            question,
            cached_payload,
            topic="",
        )
        cached_payload["conversation_id"] = cached_conversation_id
        return filter_qa_response_by_role(user_role, cached_payload)

    doc_sources = db_service.qa_retrieve_documents(
        retrieval_question,
        limit=limit,
        access_context=access_context,
    )
    if contextual_sources:
        doc_sources = _qa_merge_sources(doc_sources, contextual_sources)
    logger.info(f"[QA API] Retrieved {len(doc_sources)} document sources from database")

    chunk_sources = []
    if use_chunks:
        allowed_chunk_doc_ids = [
            str(src.get("doc_id") or "").strip()
            for src in doc_sources
            if isinstance(src, dict) and str(src.get("doc_id") or "").strip()
        ]
        chunk_sources = db_service.qa_retrieve_chunks(
            retrieval_question,
            limit=chunk_limit,
            allowed_doc_ids=allowed_chunk_doc_ids or None,
            access_context=access_context,
        )
        logger.info(f"[QA API] Retrieved {len(chunk_sources)} chunk sources from database")

    sources_for_answer = _build_answer_sources(doc_sources, chunk_sources)
    if contextual_sources:
        sources_for_answer = _qa_merge_sources(sources_for_answer, contextual_sources)

    if not sources_for_answer:
        logger.warning(f"[QA API] No sources found for question: {question}")
        no_data_response = {
            "answer": _QA_SOFT_NOT_FOUND_ANSWER,
            "suggestion": "",
            "conversation_id": requested_conversation_id,
            "structured_answer": {
                "style": "paragraph",
                "title": "",
                "summary": _QA_SOFT_NOT_FOUND_ANSWER,
                "highlights": [],
                "sections": [],
                "closing": "",
            },
            "citations": [],
            "sources": [],
            "chunk_sources": [],
            "graph_id": "",
            "chart_options": [],
            "follow_up_suggestions": [
                "Try asking with an exact document number like PO-001 or INV-001.",
                "Ask me to list related documents first, then I can compare them.",
            ],
        }
        _qa_conversation_append(
            session_id,
            {
                "question": question,
                "answer": no_data_response["answer"],
                "topic": "",
                "sources": [],
            },
        )
        _qa_history_store_turn(
            history_user_key,
            requested_conversation_id,
            question,
            no_data_response,
            topic="",
        )
        return no_data_response

    result = answer_question(
        question,
        sources_for_answer,
        conversation_context=conversation_context,
    )
    logger.info(f"[QA API] Answer generated: {result.get('answer', '')[:100]}...")

    sources_for_ui = list(doc_sources or [])
    if chunk_sources:
        existing_ids = {src.get("doc_id") for src in sources_for_ui if isinstance(src, dict)}
        for chunk_src in chunk_sources:
            doc_id = chunk_src.get("doc_id") if isinstance(chunk_src, dict) else None
            if not doc_id or doc_id in existing_ids:
                continue
            sources_for_ui.append(
                {
                    "doc_id": doc_id,
                    "document_type": chunk_src.get("document_type", "unknown"),
                    "metadata": chunk_src.get("metadata", ""),
                    "snippet": chunk_src.get("snippet", ""),
                    "items": chunk_src.get("items", []),
                    "full_text": "",
                    "confidence": chunk_src.get("confidence", {}),
                    "score": chunk_src.get("score", 0),
                    "reasons": ["chunk_match"],
                }
            )
            existing_ids.add(doc_id)

    graph_data_for_response, chart_options_for_response, has_graph_data = _qa_finalize_graph_payload(
        question,
        result.get("graph_data"),
        result.get("chart_options"),
        sources_for_answer,
    )

    response_payload = {
        "answer": result.get("answer", _QA_SOFT_NOT_FOUND_ANSWER),
        "suggestion": result.get("suggestion", ""),
        "conversation_id": requested_conversation_id,
        "structured_answer": result.get("structured_answer", {}),
        "citations": result.get("citations", []),
        "sources": sources_for_ui,
        "chunk_sources": chunk_sources,
        "graph_data": graph_data_for_response,
        "chart_options": chart_options_for_response,
    }
    response_payload["follow_up_suggestions"] = _qa_build_follow_up_suggestions(
        question=question,
        answer=str(response_payload.get("answer") or ""),
        sources=sources_for_ui,
        has_graph_data=has_graph_data,
    )
    cache_payload = dict(response_payload)
    response_payload = _qa_graph_attach_to_response(session_id, response_payload)
    if response_payload.get("graph_id"):
        logger.info(
            f"[QA API] Graph attached: {response_payload.get('graph_id')} "
            f"options={response_payload.get('chart_options')}"
        )
    elif _qa_question_requests_graph(question):
        logger.info("[QA API] Graph requested but no graph data could be generated")
    stored_conversation_id = _qa_history_store_turn(
        history_user_key,
        requested_conversation_id,
        question,
        response_payload,
        topic=result.get("topic", ""),
    )
    response_payload["conversation_id"] = stored_conversation_id
    _qa_cache_set(cache_key, cache_payload)
    _qa_conversation_append(
        session_id,
        {
            "question": question,
            "answer": response_payload.get("answer"),
            "topic": result.get("topic", ""),
            "sources": sources_for_answer[:6],
        },
    )
    return filter_qa_response_by_role(user_role, response_payload)


@router.post("/qa/stream")
async def question_answer_streaming(request: Request, payload: QARequest):
    """Stream a question answer using stored documents with citations (Employee+ only)."""
    user_session = _get_session_data_safe(request)
    user_role = get_session_role(user_session) if user_session else None
    if not user_role or not check_permission(user_session, "view_ai_buzz"):
        raise HTTPException(status_code=403, detail="Not authorized to access Q&A")

    question = str(payload.question or "").strip()
    history_user_key = _qa_history_user_key(request)
    requested_conversation_id = str(payload.conversation_id or "").strip() or str(uuid4())
    replace_turn_id = str(payload.replace_turn_id or "").strip()
    replace_prefix_turns = []
    if not question:
        raise HTTPException(status_code=400, detail="Question is required")
    if replace_turn_id:
        resolved_conversation_id = requested_conversation_id or _qa_history_find_conversation_id_by_turn_id(
            history_user_key,
            replace_turn_id,
        )
        _, replace_turn_index, replace_turns = _qa_history_get_turn_slice(
            history_user_key,
            resolved_conversation_id,
            replace_turn_id,
        )
        if replace_turn_index < 0:
            fallback_conversation_id = _qa_history_find_conversation_id_by_turn_id(history_user_key, replace_turn_id)
            if fallback_conversation_id and fallback_conversation_id != resolved_conversation_id:
                requested_conversation_id = fallback_conversation_id
                _, replace_turn_index, replace_turns = _qa_history_get_turn_slice(
                    history_user_key,
                    requested_conversation_id,
                    replace_turn_id,
                )
        else:
            requested_conversation_id = resolved_conversation_id
        if replace_turn_index < 0:
            raise HTTPException(status_code=404, detail="The selected chat turn is no longer available.")
        target_turn = replace_turns[replace_turn_index]
        if not str(target_turn.get("answer") or "").strip():
            raise HTTPException(status_code=409, detail="Please wait for the current answer to finish before regenerating it.")
        replace_prefix_turns = replace_turns[:replace_turn_index]
    if len(question) == 1:
        response_payload = _qa_incomplete_input_payload(requested_conversation_id)
        if replace_turn_id:
            stored_conversation_id, stored_turn_id = _qa_history_replace_turn_result(
                history_user_key,
                requested_conversation_id,
                replace_turn_id,
                question,
                response_payload,
                topic="",
            )
        else:
            stored_conversation_id, stored_turn_id = _qa_history_store_turn_result(
                history_user_key,
                requested_conversation_id,
                question,
                response_payload,
                topic="",
            )
        response_payload["conversation_id"] = stored_conversation_id
        response_payload["turn_id"] = stored_turn_id

        def short_stream():
            yield _qa_stream_event({
                "type": "start",
                "conversation_id": stored_conversation_id,
            })
            yield _qa_stream_event({
                "type": "final",
                "data": response_payload,
            })

        return StreamingResponse(short_stream(), media_type="application/x-ndjson")

    try:
        limit = int(payload.limit)
    except Exception:
        limit = 6

    use_chunks = bool(payload.use_chunks)
    try:
        raw_chunk_limit = payload.chunk_limit
        chunk_limit = int(raw_chunk_limit if raw_chunk_limit is not None else max(limit * 2, 6))
    except Exception:
        chunk_limit = max(limit * 2, 6)

    simple_lookup = any(
        token in question.lower()
        for token in [
            "email", "mail", "phone", "mobile", "contact", "address",
            "gst", "gstin", "pan", "invoice", "bill number", "document number",
            "amount", "date", "number",
        ]
    )
    broad_document_listing = _qa_requests_broad_document_listing(question)
    if broad_document_listing:
        limit = 12
        chunk_limit = max(8, min(chunk_limit, 24))
    else:
        limit = max(1, min(limit, 3 if simple_lookup else 5))
        chunk_limit = max(1, min(chunk_limit, 12))

    logger.info(f"[QA API] Streaming question: {question[:100]} (requesting {limit} sources) by role: {user_role}")

    access_context = _build_request_kbac_context(request)
    session_id = _get_request_session_id(request)

    qa_history = list(replace_prefix_turns)
    if not qa_history and not replace_turn_id:
        qa_history = _qa_history_get_turns(history_user_key, requested_conversation_id)
    if not qa_history and not replace_turn_id:
        qa_history = _qa_conversation_get(session_id)

    conversation_context = _qa_build_conversation_context(qa_history)
    retrieval_question = _qa_build_contextual_question(question, qa_history)
    contextual_sources = _qa_collect_context_sources(question, qa_history)
    context_signature = _qa_build_context_signature(qa_history)
    cache_scope_key = f"{access_context.get('scope_key', 'all')}|ctx:{context_signature}"
    cache_key = _qa_cache_key(question, limit, chunk_limit, use_chunks, cache_scope_key)

    def stream():
        yield _qa_stream_event({
            "type": "start",
            "conversation_id": requested_conversation_id,
        })

        try:
            cached_response = _qa_cache_get(cache_key)
            if cached_response is not None:
                logger.info("[QA API] Full response cache hit (stream)")
                cached_payload = dict(cached_response)
                cached_graph_data, cached_chart_options, _ = _qa_finalize_graph_payload(
                    question,
                    cached_payload.get("graph_data"),
                    cached_payload.get("chart_options"),
                    cached_payload.get("sources") if isinstance(cached_payload.get("sources"), list) else [],
                )
                cached_payload["graph_data"] = cached_graph_data
                cached_payload["chart_options"] = cached_chart_options
                cached_payload = _qa_graph_attach_to_response(session_id, cached_payload)
                cached_followups = cached_payload.get("follow_up_suggestions")
                if not isinstance(cached_followups, list) or not cached_followups:
                    cached_payload["follow_up_suggestions"] = _qa_build_follow_up_suggestions(
                        question=question,
                        answer=str(cached_payload.get("answer") or ""),
                        sources=cached_payload.get("sources") if isinstance(cached_payload.get("sources"), list) else [],
                        has_graph_data=bool(cached_payload.get("graph_id")),
                    )
                cached_turn_summary = {
                    "question": question,
                    "answer": cached_payload.get("answer", ""),
                    "topic": "",
                    "sources": _build_answer_sources(
                        cached_payload.get("sources") if isinstance(cached_payload.get("sources"), list) else [],
                        cached_payload.get("chunk_sources") if isinstance(cached_payload.get("chunk_sources"), list) else [],
                    )[:6],
                }
                if replace_turn_id:
                    cached_conversation_id, cached_turn_id = _qa_history_replace_turn_result(
                        history_user_key,
                        requested_conversation_id,
                        replace_turn_id,
                        question,
                        cached_payload,
                        topic="",
                    )
                    _qa_conversation_set(session_id, replace_prefix_turns + [cached_turn_summary])
                else:
                    _qa_conversation_append(session_id, cached_turn_summary)
                    cached_conversation_id, cached_turn_id = _qa_history_store_turn_result(
                        history_user_key,
                        requested_conversation_id,
                        question,
                        cached_payload,
                        topic="",
                    )
                cached_payload["conversation_id"] = cached_conversation_id
                cached_payload["turn_id"] = cached_turn_id
                yield _qa_stream_event({
                    "type": "final",
                    "data": filter_qa_response_by_role(user_role, cached_payload),
                })
                return

            doc_sources = db_service.qa_retrieve_documents(
                retrieval_question,
                limit=limit,
                access_context=access_context,
            )
            if contextual_sources:
                doc_sources = _qa_merge_sources(doc_sources, contextual_sources)
            logger.info(f"[QA API] Retrieved {len(doc_sources)} document sources from database for stream")

            chunk_sources = []
            if use_chunks:
                allowed_chunk_doc_ids = [
                    str(src.get("doc_id") or "").strip()
                    for src in doc_sources
                    if isinstance(src, dict) and str(src.get("doc_id") or "").strip()
                ]
                chunk_sources = db_service.qa_retrieve_chunks(
                    retrieval_question,
                    limit=chunk_limit,
                    allowed_doc_ids=allowed_chunk_doc_ids or None,
                    access_context=access_context,
                )
                logger.info(f"[QA API] Retrieved {len(chunk_sources)} chunk sources from database for stream")

            sources_for_answer = _build_answer_sources(doc_sources, chunk_sources)
            if contextual_sources:
                sources_for_answer = _qa_merge_sources(sources_for_answer, contextual_sources)

            if not sources_for_answer:
                logger.warning(f"[QA API] No sources found for streaming question: {question}")
                no_data_response = {
                    "answer": _QA_SOFT_NOT_FOUND_ANSWER,
                    "suggestion": "",
                    "conversation_id": requested_conversation_id,
                    "structured_answer": {
                        "style": "paragraph",
                        "title": "",
                        "summary": _QA_SOFT_NOT_FOUND_ANSWER,
                        "highlights": [],
                        "sections": [],
                        "closing": "",
                    },
                    "citations": [],
                    "sources": [],
                    "chunk_sources": [],
                }
                if replace_turn_id:
                    stored_conversation_id, stored_turn_id = _qa_history_replace_turn_result(
                        history_user_key,
                        requested_conversation_id,
                        replace_turn_id,
                        question,
                        no_data_response,
                        topic="",
                    )
                    _qa_conversation_set(
                        session_id,
                        replace_prefix_turns + [{
                            "question": question,
                            "answer": no_data_response["answer"],
                            "topic": "",
                            "sources": [],
                        }],
                    )
                else:
                    stored_conversation_id, stored_turn_id = _qa_history_store_turn_result(
                        history_user_key,
                        requested_conversation_id,
                        question,
                        no_data_response,
                        topic="",
                    )
                    _qa_conversation_append(
                        session_id,
                        {
                            "question": question,
                            "answer": no_data_response["answer"],
                            "topic": "",
                            "sources": [],
                        },
                    )
                no_data_response["conversation_id"] = stored_conversation_id
                no_data_response["turn_id"] = stored_turn_id
                yield _qa_stream_event({"type": "final", "data": no_data_response})
                return

            final_result = None
            for event in answer_question_stream(
                question,
                sources_for_answer,
                conversation_context=conversation_context,
            ):
                event_type = str(event.get("type") or "").strip().lower()
                if event_type == "delta":
                    yield _qa_stream_event({"type": "delta", "delta": str(event.get("delta") or "")})
                    continue
                if event_type == "final":
                    final_result = dict(event.get("data") or {})

            if final_result is None:
                raise RuntimeError("Streaming QA completed without a final payload.")

            logger.info(f"[QA API] Streaming answer generated: {str(final_result.get('answer') or '')[:100]}...")
            sources_for_ui = _build_qa_ui_sources(doc_sources, chunk_sources)
            graph_data_for_response, chart_options_for_response, has_graph_data = _qa_finalize_graph_payload(
                question,
                final_result.get("graph_data"),
                final_result.get("chart_options"),
                sources_for_answer,
            )
            response_payload = {
                "answer": final_result.get("answer", _QA_SOFT_NOT_FOUND_ANSWER),
                "suggestion": final_result.get("suggestion", ""),
                "conversation_id": requested_conversation_id,
                "structured_answer": final_result.get("structured_answer", {}),
                "citations": final_result.get("citations", []),
                "sources": sources_for_ui,
                "chunk_sources": chunk_sources,
                "graph_data": graph_data_for_response,
                "chart_options": chart_options_for_response,
            }
            response_payload["follow_up_suggestions"] = _qa_build_follow_up_suggestions(
                question=question,
                answer=str(response_payload.get("answer") or ""),
                sources=sources_for_ui,
                has_graph_data=has_graph_data,
            )
            cache_payload = dict(response_payload)
            response_payload = _qa_graph_attach_to_response(session_id, response_payload)
            if response_payload.get("graph_id"):
                logger.info(
                    f"[QA API] Streaming graph attached: {response_payload.get('graph_id')} "
                    f"options={response_payload.get('chart_options')}"
                )
            elif _qa_question_requests_graph(question):
                logger.info("[QA API] Streaming graph requested but no graph data could be generated")
            turn_summary = {
                "question": question,
                "answer": response_payload.get("answer"),
                "topic": final_result.get("topic", ""),
                "sources": sources_for_answer[:6],
            }
            if replace_turn_id:
                stored_conversation_id, stored_turn_id = _qa_history_replace_turn_result(
                    history_user_key,
                    requested_conversation_id,
                    replace_turn_id,
                    question,
                    response_payload,
                    topic=final_result.get("topic", ""),
                )
                _qa_conversation_set(session_id, replace_prefix_turns + [turn_summary])
            else:
                stored_conversation_id, stored_turn_id = _qa_history_store_turn_result(
                    history_user_key,
                    requested_conversation_id,
                    question,
                    response_payload,
                    topic=final_result.get("topic", ""),
                )
                _qa_conversation_append(session_id, turn_summary)
            response_payload["conversation_id"] = stored_conversation_id
            response_payload["turn_id"] = stored_turn_id
            _qa_cache_set(cache_key, cache_payload)
            yield _qa_stream_event({
                "type": "final",
                "data": filter_qa_response_by_role(user_role, response_payload),
            })
        except Exception as exc:
            logger.exception("[QA API] Streaming error")
            yield _qa_stream_event({
                "type": "error",
                "message": str(exc) or "Failed to generate answer.",
            })

    return StreamingResponse(stream(), media_type="application/x-ndjson")
@router.get("/qa/chart")
def qa_chart_image(
    request: Request,
    graph_id: str = Query(..., min_length=5),
    chart_type: str = Query("bar"),
):
    user_session = _get_session_data_safe(request)
    user_role = get_session_role(user_session) if user_session else None
    if not user_role or not check_permission(user_session, "view_ai_buzz"):
        raise HTTPException(status_code=403, detail="Not authorized to access Q&A charts")

    if not _MATPLOTLIB_AVAILABLE:
        raise HTTPException(
            status_code=503,
            detail=f"Chart rendering is unavailable because matplotlib could not be loaded: {_MATPLOTLIB_IMPORT_ERROR}",
        )

    safe_chart_type = str(chart_type or "").strip().lower()
    session_id = _get_request_session_id(request)
    graph_payload = _qa_graph_get(session_id, graph_id)
    if not isinstance(graph_payload, dict):
        raise HTTPException(
            status_code=404,
            detail="Chart data expired for this question. Please ask again to regenerate chart data.",
        )

    chart_options = _normalize_qa_chart_options(graph_payload.get("chart_options"), has_graph_data=True)
    if safe_chart_type not in chart_options:
        raise HTTPException(status_code=400, detail="Chart type not available for this response")

    graph_data = _normalize_qa_graph_data(graph_payload.get("graph_data"))
    if not graph_data:
        raise HTTPException(status_code=400, detail="Invalid chart data")

    try:
        image_bytes = _qa_render_chart_png(graph_data, safe_chart_type)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc))
    except Exception as exc:
        logger.warning(f"[QA API] Chart rendering failed: {exc}")
        raise HTTPException(status_code=500, detail="Failed to render chart")

    return StreamingResponse(
        BytesIO(image_bytes),
        media_type="image/png",
        headers={"Cache-Control": "no-store, max-age=0"},
    )


@router.post("/qa/ppt")
def generate_qa_presentation(
    request: Request,
    payload: QAPresentationRequest,
    temporaryRetentionHours: int = Query(DEFAULT_TEMP_RETENTION_HOURS),
):
    _require_permission(request, "view_ai_buzz", "Only authorized users can generate AI Buzz presentations.")
    if not _ensure_pptx_available():
        raise HTTPException(
            status_code=500,
            detail=f"PowerPoint generation is unavailable because python-pptx could not be loaded: {_PPTX_IMPORT_ERROR}",
        )

    question = str(payload.question or "").strip()
    answer = str(payload.answer or "").strip()
    if not question:
        raise HTTPException(status_code=400, detail="Question is required to generate a presentation.")
    if not answer:
        raise HTTPException(status_code=400, detail="Answer is required to generate a presentation.")

    output_path = None
    retention_hours = resolve_retention_hours(temporaryRetentionHours)
    try:
        output_path = _qa_build_presentation_file(payload)
        download_name = _qa_ppt_build_download_name(payload)
        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=download_name,
        )
    except HTTPException:
        raise
    except Exception as exc:
        logger.error(f"[QA PPT] Failed to generate presentation: {exc}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Failed to generate AI Buzz presentation.")
    finally:
        if output_path:
            try:
                _schedule_temp_file_cleanup(output_path, retention_hours)
            except Exception:
                pass


@router.get("/health", response_model=HealthResponse)
def health_check():
    """Health check endpoint - confirms API is running"""
    return {"status": "ok", "message": "API is running"}


@router.get("/system-settings", response_model=SystemSettingsResponse)
def get_system_settings(request: Request):
    user_role = _get_request_role(request)
    if user_role not in {UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        raise HTTPException(status_code=403, detail="Only Admin or Super Admin can view system settings.")
    return {
        "generate_video_transcripts": settings.generate_video_transcripts,
        "save_original_documents": settings.save_original_documents,
        "guided_tour_enabled": True,
        "allowed_user_email_domain": settings.allowed_user_email_domain,
    }


@router.put("/system-settings", response_model=SystemSettingsResponse)
def update_system_settings(request: Request, payload: UpdateSystemSettingsRequest):
    user_role = _get_request_role(request)
    if user_role not in {UserRole.ADMIN, UserRole.SUPER_ADMIN}:
        raise HTTPException(status_code=403, detail="Only Admin or Super Admin can update system settings.")
    allowed_domain = str(payload.allowed_user_email_domain or "").strip().lower().lstrip("@")
    if not allowed_domain:
        raise HTTPException(status_code=400, detail="Allowed email domain is required.")
    settings.set_persistent_value("GENERATE_VIDEO_TRANSCRIPTS", payload.generate_video_transcripts)
    settings.set_persistent_value("SAVE_ORIGINAL_DOCUMENTS", payload.save_original_documents)
    settings.set_persistent_value("ALLOWED_USER_EMAIL_DOMAIN", allowed_domain)
    return {
        "generate_video_transcripts": settings.generate_video_transcripts,
        "save_original_documents": settings.save_original_documents,
        "guided_tour_enabled": True,
        "allowed_user_email_domain": settings.allowed_user_email_domain,
    }

# ============= PROCESSING STATUS ENDPOINT =============

@router.get("/status/{filename}", response_model=ProcessingStatusResponse)
def get_processing_status(request: Request, filename: str):
    """Get processing status for a specific file
    
    Checks:
    1. In-memory cache (fast)
    2. MongoDB persistent storage (fallback)
    Returns status with progress percentage
    """
    _require_permission(request, "upload_document", "Only Manager or higher roles can view processing status.")
    # Check memory first (most recent)
    if filename in processing_status:
        return _serialize_processing_status(processing_status[filename])
    
    # Check MongoDB for persisted status
    try:
        from db.mongo_db import MongoDBManager
        db = MongoDBManager()
        status_doc = db.db['processing_status'].find_one({"filename": filename})
        db.close()
        
        if status_doc:
            return _serialize_processing_status(status_doc)
    except Exception as e:
        logger.warning(f"Could not fetch status from DB: {e}")
    
    # Default: not found
    return {
        "status": "not_found",
        "progress": 0,
        "message": "File not found in processing queue"
    }


@router.post(
    "/employee-documents/{doc_id}/resolve-profile",
    response_model=EmployeeDocumentProfileResolutionResponse,
)
def resolve_employee_document_profile(request: Request, doc_id: str, payload: ResolveEmployeeProfileRequest):
    from bson import ObjectId

    action = _normalize_optional_text(payload.action).lower()
    if action not in {"save_profile", "keep_temporary"}:
        raise HTTPException(status_code=400, detail="action must be save_profile or keep_temporary")

    db = MongoDBManager()
    try:
        user_role = _require_permission(
            request,
            "resolve_employee_profile",
            "Only Manager or higher roles can resolve employee profiles.",
        )
        actor = _get_request_actor(request)
        doc = db.collection.find_one(db._active_documents_query({"_id": ObjectId(doc_id)}))
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")

        suggested = doc.get("suggested_profile") if isinstance(doc.get("suggested_profile"), dict) else {}
        retention_hours = resolve_retention_hours(
            payload.temporaryRetentionHours
            or doc.get("source", {}).get("temporaryRetentionHours")
            or DEFAULT_TEMP_RETENTION_HOURS
        )

        if action == "keep_temporary":
            current_storage_mode = normalize_choice(
                doc.get("storage_mode") or doc.get("source", {}).get("requested_storage_mode")
            )
            if current_storage_mode != "temporary":
                raise HTTPException(status_code=400, detail="This document was not uploaded in temporary mode")

            expiry_at = build_expiry_at(retention_hours)
            update_payload = {
                "storage_mode": "temporary",
                "isTemporary": True,
                "expiry_at": expiry_at,
                "assignment_status": "temporary_unassigned",
                "updated_at": datetime.utcnow(),
            }

            db.update_document(
                doc_id,
                update_payload,
            )
            _invalidate_qa_caches()
            return {
                "status": "saved_as_temporary",
                "document_id": doc_id,
                "expiry_at": _serialize_datetime(expiry_at),
            }

        emp_id = _normalize_optional_text(payload.empID or suggested.get("empID"))
        emp_name = _normalize_optional_text(payload.empName or suggested.get("empName"))
        if not emp_id or not emp_name:
            raise HTTPException(status_code=400, detail="empID and empName are required to save an employee profile")

        existing_profile = db.get_employee_profile(emp_id=emp_id)
        if user_role == UserRole.MANAGER and existing_profile:
            profile_owner = str(existing_profile.get("created_by") or "").strip().lower()
            if not profile_owner or profile_owner != actor.lower():
                raise HTTPException(
                    status_code=403,
                    detail="Managers can only link documents to employee profiles they created.",
                )

        profile = db.create_or_update_employee_profile(
            emp_id,
            emp_name,
            extra={"created_by": actor} if user_role == UserRole.MANAGER and not existing_profile else {},
        )
        if not profile:
            raise HTTPException(status_code=500, detail="Failed to save employee profile")

        current_storage_mode = normalize_choice(
            doc.get("storage_mode") or doc.get("source", {}).get("requested_storage_mode")
        )
        update_payload = {
            "empID": profile.get("empID"),
            "empName": profile.get("empName"),
            "employee_uuid": profile.get("uuid"),
            "storage_mode": current_storage_mode,
            "isTemporary": current_storage_mode == "temporary",
            "expiry_at": doc.get("expiry_at") if current_storage_mode == "temporary" else None,
            "assignment_status": "linked",
            "suggested_profile": build_profile_suggestion(profile.get("empID"), profile.get("empName")),
        }

        db.update_document(doc_id, update_payload)
        _invalidate_qa_caches()
        return {
            "status": "profile_linked",
            "document_id": doc_id,
            "employee": {
                "empID": profile.get("empID"),
                "empName": profile.get("empName"),
                "uuid": profile.get("uuid"),
            },
        }
    except HTTPException:
        raise
    except Exception as exc:
        logger.error(f"Failed to resolve employee profile for document {doc_id}: {exc}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Failed to resolve employee profile")
    finally:
        db.close()


@router.post("/cleanup-temporary-documents", response_model=CleanupTemporaryDocumentsResponse)
def cleanup_temporary_documents(request: Request):
    _require_permission(request, "delete_document", "Admins and Super Admins only can clean temporary documents.")
    result = db_service.cleanup_expired_temporary_documents()
    if result.get("deleted_documents", 0):
        _invalidate_qa_caches()
    return {
        "status": "completed",
        "deleted_documents": result.get("deleted_documents", 0),
        "deleted_files": result.get("deleted_files", 0),
    }

# ============= DOCUMENT DELETION ENDPOINT =============

@router.post("/documents/bulk-delete", response_model=BulkDeleteResponse)
def bulk_delete_documents(request: Request, payload: BulkDeleteDocumentsRequest):
    _require_permission(request, "delete_document", "Admins and Super Admins only can delete documents.")
    doc_ids = []
    seen = set()
    for raw_doc_id in payload.doc_ids:
        doc_id = str(raw_doc_id or "").strip()
        if not doc_id or doc_id in seen:
            continue
        seen.add(doc_id)
        doc_ids.append(doc_id)

    if not doc_ids:
        raise HTTPException(status_code=400, detail="At least one valid document ID is required")

    result = db_service.delete_documents(doc_ids)
    if result.get("deleted_count", 0):
        _invalidate_qa_caches()
    return {
        "status": "completed",
        "requested_count": result["requested_count"],
        "deleted_count": result["deleted_count"],
        "deleted_ids": result["deleted_ids"],
        "not_found_ids": result["not_found_ids"],
    }

@router.delete("/documents/{doc_id}", response_model=DeleteDocumentResponse)
def delete_document(request: Request, doc_id: str):
    """Delete document by ID"""
    _require_permission(request, "delete_document", "Admins and Super Admins only can delete documents.")
    result = db_service.delete_document(doc_id)
    if not result:
        raise HTTPException(status_code=404, detail="Document not found")
    _invalidate_qa_caches()
    return {"status": "deleted", "doc_id": doc_id}

# ============= FILE UPLOAD & PROCESSING ENDPOINTS =============

@router.post("/upload", response_model=UploadQueuedResponse)
async def upload_file(
    request: Request,  # NEW: Add request for RBAC
    file: UploadFile = File(...),
    empID: Optional[str] = Form(None),
    empName: Optional[str] = Form(None),
    choose: str = Form("permanent"),
    temporaryRetentionHours: int = Form(DEFAULT_TEMP_RETENTION_HOURS),
    background_tasks: BackgroundTasks = BackgroundTasks()
):
    """Upload and process single file asynchronously (Manager+ only)
    
    Supported formats:
    - Documents: PDF, DOCX, DOC, DOCM
    - Spreadsheets: XLSX, XLS, XLSM, XLSB
    - Images: JPG, JPEG, PNG, TIFF, TIF
    Processing: Extraction → LLM structuring → MongoDB storage
    Returns immediately while processing happens in background
    """
    # NEW: RBAC Check - Only Manager+ can upload
    user_role = _require_permission(request, "upload_document", "Only Manager or higher roles can upload documents")
    actor = _get_request_actor(request)

    if empID:
        _require_employee_profile_access(request, emp_id=empID, write=True)
    
    allowed_extensions = _ALLOWED_EXTENSIONS
    file_ext = Path(file.filename).suffix.lower()

    if file_ext not in allowed_extensions:
        raise HTTPException(
            status_code=400, 
            detail=f"File type '{file_ext}' not supported. Allowed: {', '.join(allowed_extensions)}"
        )

    normalized_choice = _resolve_storage_mode(choose)
    retention_hours = resolve_retention_hours(temporaryRetentionHours)
    display_filename = _build_tracking_filename(file.filename)
    tmp_path, size_bytes, file_hash = await _persist_uploaded_file(file)
    duplicate_match = db_service.find_duplicate_by_file_hash(
        file_hash,
        duplicate_owner=_normalize_kbac_owner(actor),
    )
    if duplicate_match:
        _safe_delete_file(tmp_path)
        message = _build_duplicate_skip_message(file.filename, duplicate_match)
        _register_duplicate_skip_status(display_filename, message, duplicate_match)
        return {
            "status": "duplicate_skipped",
            "filename": display_filename,
            "message": message,
            "duplicate_of": duplicate_match,
        }

    processing_path, persisted_path, preserve_file = _prepare_processing_file(tmp_path, file.filename)

    try:
        _queue_processing_task(
            background_tasks,
            processing_path,
            display_filename,
            "File uploaded, starting OCR...",
            {
                "type": "upload",
                "original_filename": file.filename,
                "file_name": file.filename,
                "file_path": persisted_path,
                "size_bytes": size_bytes,
                "file_hash": file_hash,
                "empID": _normalize_optional_text(empID),
                "empName": _normalize_optional_text(empName),
                "choose": normalized_choice,
                "temporaryRetentionHours": retention_hours,
                "uploaded_by": actor,
                "uploaded_by_role": user_role.value,
                "preserve_file": preserve_file,
                "save_original_documents": preserve_file,
                "generate_video_text": settings.generate_video_transcripts,
            },
        )
        return {
            "status": "processing",
            "filename": display_filename,
            "message": "File uploaded and queued for OCR processing"
        }
    except Exception as e:
        cleanup_path = persisted_path or processing_path
        if cleanup_path and os.path.exists(cleanup_path):
            os.remove(cleanup_path)
        processing_status[display_filename] = {
            "status": "error",
            "progress": 0,
            "message": str(e)
        }
        save_status_to_db(display_filename, processing_status[display_filename])
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/upload-multiple", response_model=MultiUploadQueuedResponse)
async def upload_multiple_files(
    request: Request,
    files: List[UploadFile] = File(...),
    empID: Optional[str] = Form(None),
    empName: Optional[str] = Form(None),
    choose: str = Form("permanent"),
    temporaryRetentionHours: int = Form(DEFAULT_TEMP_RETENTION_HOURS),
    background_tasks: BackgroundTasks = BackgroundTasks()
):
    """Upload and process multiple files asynchronously
    
    Supported formats:
    - Documents: PDF, DOCX, DOC, DOCM
    - Spreadsheets: XLSX, XLS, XLSM, XLSB
    - Images: JPG, JPEG, PNG, TIFF, TIF
    Processing: Extraction → LLM structuring → MongoDB storage
    Returns when files are queued (actual processing happens in background)
    """
    user_role = _require_permission(request, "upload_document", "Only Manager or higher roles can upload documents.")
    actor = _get_request_actor(request)
    if empID:
        _require_employee_profile_access(request, emp_id=empID, write=True)

    uploaded_files = []
    skipped_files = []
    seen_file_hashes: dict[str, str] = {}
    allowed_extensions = _ALLOWED_EXTENSIONS
    normalized_choice = _resolve_storage_mode(choose)
    retention_hours = resolve_retention_hours(temporaryRetentionHours)

    for file in files:
        file_ext = Path(file.filename).suffix.lower()
        if file_ext not in allowed_extensions:
            continue
        
        processing_path = None
        persisted_path = None
        preserve_file = False
        try:
            display_filename = _build_tracking_filename(file.filename)
            tmp_path, size_bytes, file_hash = await _persist_uploaded_file(file)

            duplicate_match = None
            if file_hash in seen_file_hashes:
                duplicate_match = {
                    "tracking_filename": seen_file_hashes[file_hash],
                    "file_name": file.filename,
                    "reason": "batch_duplicate",
                    "file_hash": file_hash,
                }
            else:
                duplicate_match = db_service.find_duplicate_by_file_hash(
                    file_hash,
                    duplicate_owner=_normalize_kbac_owner(actor),
                )

            if duplicate_match:
                _safe_delete_file(tmp_path)
                message = _build_duplicate_skip_message(file.filename, duplicate_match, duplicate_match.get("reason"))
                _register_duplicate_skip_status(display_filename, message, duplicate_match)
                uploaded_files.append(display_filename)
                skipped_files.append(
                    {
                        "filename": display_filename,
                        "original_filename": file.filename,
                        "message": message,
                        "duplicate_of": duplicate_match,
                    }
                )
                continue

            seen_file_hashes[file_hash] = display_filename
            processing_path, persisted_path, preserve_file = _prepare_processing_file(tmp_path, file.filename)

            _queue_processing_task(
                background_tasks,
                processing_path,
                display_filename,
                "File uploaded, starting OCR...",
                {
                    "type": "upload",
                    "original_filename": file.filename,
                    "file_name": file.filename,
                    "file_path": persisted_path,
                    "size_bytes": size_bytes,
                    "file_hash": file_hash,
                    "empID": _normalize_optional_text(empID),
                    "empName": _normalize_optional_text(empName),
                    "choose": normalized_choice,
                    "temporaryRetentionHours": retention_hours,
                    "uploaded_by": actor,
                    "uploaded_by_role": user_role.value,
                    "preserve_file": preserve_file,
                    "save_original_documents": preserve_file,
                    "generate_video_text": settings.generate_video_transcripts,
                },
            )
            uploaded_files.append(display_filename)
        except Exception as e:
            cleanup_path = persisted_path or processing_path
            if cleanup_path and os.path.exists(cleanup_path):
                os.remove(cleanup_path)
            logger.error(f"Error processing {file.filename}: {str(e)}")
            failed_name = _build_tracking_filename(file.filename)
            processing_status[failed_name] = {
                "status": "error",
                "progress": 0,
                "message": str(e)
            }
            save_status_to_db(failed_name, processing_status[failed_name])
            continue
    
    if not uploaded_files:
        raise HTTPException(status_code=400, detail="No valid files to process")
    
    return {
        "status": "processing",
        "uploaded_count": len(uploaded_files),
        "files": uploaded_files,
        "message": (
            f"Queued {len(uploaded_files) - len(skipped_files)} new file(s) for processing."
            + (f" Skipped {len(skipped_files)} duplicate file(s)." if skipped_files else "")
        ),
        "duplicate_skipped_count": len(skipped_files),
        "skipped_files": skipped_files,
    }


@router.post("/fill-excel")
async def fill_excel_from_documents(
    request: Request,
    files: List[UploadFile] = File(...),
    excel: UploadFile = File(...),
    prompt: Optional[str] = Form(None),
    sheet_name: Optional[str] = Form(None),
    temporaryRetentionHours: int = Form(DEFAULT_TEMP_RETENTION_HOURS),
    background_tasks: BackgroundTasks = BackgroundTasks(),
):
    """Fill an Excel template using uploaded documents (no DB write)."""
    _require_permission(request, "upload_document", "Only Manager or higher roles can fill Excel from documents.")
    if not gemini_is_configured():
        raise HTTPException(status_code=500, detail="Gemini is not configured (GEMINI_API_KEY missing)")

    excel_name = str(excel.filename or "").strip()
    excel_ext = Path(excel_name).suffix.lower()
    if excel_ext not in {".xlsx", ".xlsm"}:
        raise HTTPException(status_code=400, detail="Excel file must be .xlsx or .xlsm")

    if not files:
        raise HTTPException(status_code=400, detail="At least one document file is required")

    retention_hours = resolve_retention_hours(temporaryRetentionHours)
    temp_paths = []
    doc_payloads = []
    excel_path = None
    output_path = None

    try:
        excel_path, _, _ = await _persist_uploaded_file(excel)
        temp_paths.append(excel_path)

        for upload_file in files:
            filename = str(upload_file.filename or "").strip()
            if not filename:
                continue
            file_ext = Path(filename).suffix.lower()
            if file_ext in {".xlsx", ".xlsm"}:
                # Skip accidental Excel uploads in the documents list
                continue
            if file_ext not in _ALLOWED_EXTENSIONS:
                continue

            temp_path, _, _ = await _persist_uploaded_file(upload_file)
            temp_paths.append(temp_path)

            all_text, confidences, link_targets, _ = _extract_text_from_document(
                temp_path,
                filename,
                status_filename=None,
            )
            if not str(all_text or "").strip():
                logger.warning(f"[Excel Fill] No text extracted from {filename}")
                continue

            linked_documents = []
            if link_targets and _LINKED_PDF_MAX_DOCUMENTS > 0:
                logger.info(f"[Excel Fill] Linked PDF targets found: {len(link_targets)}")
                linked_documents, linked_warnings, _ = _collect_linked_document_texts(
                    link_targets,
                    graph_access_token=None,
                )
                if linked_warnings:
                    logger.warning(f"[Excel Fill] Linked PDF warnings: {len(linked_warnings)}")

            ocr_confidence = round(sum(confidences) / len(confidences), 2) if confidences else 0.0
            base_llm_chars = settings.get_int("LLM_MAX_CHARS", 120000)
            fill_llm_chars = settings.get_int("LLM_MAX_CHARS_FILL", 180000)
            max_llm_chars = max(base_llm_chars, fill_llm_chars)
            llm_text = _build_llm_input_text(all_text, linked_documents, max_llm_chars)
            try:
                metadata = extract_metadata_with_gemini(llm_text, ocr_confidence)
            except Exception as llm_err:
                logger.error(f"[Excel Fill] LLM extraction failed for {filename}: {llm_err}")
                metadata = {"document_type": "unknown", "high_level_metadata": {}, "items": []}

            doc_payloads.append(
                {
                    "filename": filename,
                    "document_type": metadata.get("document_type", "unknown"),
                    "high_level_metadata": metadata.get("high_level_metadata", {}),
                    "items": metadata.get("items", []),
                    "ocr_confidence": ocr_confidence,
                }
            )

        if not doc_payloads:
            raise HTTPException(status_code=400, detail="No usable documents were extracted")

        wb = load_workbook(excel_path, keep_vba=(excel_ext == ".xlsm"))
        if sheet_name and sheet_name in wb.sheetnames:
            sheets = [wb[sheet_name]]
        else:
            sheets = list(wb.worksheets)
        total_updates = 0
        for target_sheet in sheets:
            form_fields = _collect_form_fields(target_sheet)
            tables = _detect_tables(target_sheet)
            if not form_fields and not tables:
                continue
            target_batches = _build_target_batches(form_fields, tables)
            for batch in target_batches:
                cells = _fill_excel_sheet_with_gemini(
                    target_sheet.title,
                    batch,
                    doc_payloads,
                    prompt,
                )
                for cell in cells or []:
                    try:
                        row = int(cell.get("row"))
                        col = int(cell.get("col"))
                    except Exception:
                        continue
                    value = cell.get("value")
                    if value is None or str(value).strip() == "":
                        continue
                    row, col = _normalize_merged_target(target_sheet, row, col)
                    current = target_sheet.cell(row=row, column=col).value
                    if not _is_empty_cell(current):
                        continue
                    target_sheet.cell(row=row, column=col, value=str(value))
                    total_updates += 1

        if total_updates == 0:
            raise HTTPException(status_code=400, detail="No matching fields were filled for the template")

        with tempfile.NamedTemporaryFile(delete=False, suffix=excel_ext) as tmp:
            output_path = tmp.name
        wb.save(output_path)

        download_name = f"filled_{excel_name or 'template.xlsx'}"
        return FileResponse(
            output_path,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=download_name,
        )
    except HTTPException:
        raise
    except Exception as exc:
        logger.error(f"[Excel Fill] Failed: {exc}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Failed to fill Excel template")
    finally:
        if output_path:
            try:
                _schedule_temp_file_cleanup(output_path, retention_hours)
            except Exception:
                pass
        for path in temp_paths:
            try:
                if path:
                    _schedule_temp_file_cleanup(path, retention_hours)
            except Exception:
                pass


@router.post("/import-from-link", response_model=ImportFromLinkResponse)
async def import_from_link(
    request: Request,
    payload: ImportFromLinkRequest,
    background_tasks: BackgroundTasks = BackgroundTasks(),
):
    """Resolve a public file/folder link and queue supported files for OCR processing."""
    user_role = _require_permission(request, "upload_document", "Only Manager or higher roles can import documents from links.")
    actor = _get_request_actor(request)
    if payload.empID:
        _require_employee_profile_access(request, emp_id=payload.empID, write=True)
    session_id = request.cookies.get("session_id")
    graph_access_token = get_microsoft_access_token_for_session(session_id)
    try:
        remote_files, warnings = import_documents_from_link(
            payload.link,
            _ALLOWED_EXTENSIONS,
            graph_access_token=graph_access_token,
        )
    except RemoteImportError as exc:
        raise HTTPException(status_code=400, detail=str(exc))
    except Exception as exc:
        logger.error(f"Remote import failed for {payload.link}: {exc}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="Failed to import documents from link")

    normalized_choice = _resolve_storage_mode(payload.choose)
    retention_hours = resolve_retention_hours(payload.temporaryRetentionHours)
    queued_files = []
    skipped_files = []
    seen_file_hashes: dict[str, str] = {}
    for remote_file in remote_files:
        display_filename = _build_tracking_filename(remote_file.filename)
        processing_path = None
        persisted_path = None
        preserve_file = False
        try:
            file_hash = _compute_file_sha256(remote_file.temp_path)

            duplicate_match = None
            if file_hash in seen_file_hashes:
                duplicate_match = {
                    "tracking_filename": seen_file_hashes[file_hash],
                    "file_name": remote_file.filename,
                    "reason": "batch_duplicate",
                    "file_hash": file_hash,
                }
            else:
                duplicate_match = db_service.find_duplicate_by_file_hash(
                    file_hash,
                    duplicate_owner=_normalize_kbac_owner(actor),
                )

            if duplicate_match:
                _safe_delete_file(remote_file.temp_path)
                message = _build_duplicate_skip_message(remote_file.filename, duplicate_match, duplicate_match.get("reason"))
                _register_duplicate_skip_status(display_filename, message, duplicate_match)
                queued_files.append(
                    {
                        "filename": display_filename,
                        "original_filename": remote_file.filename,
                        "source_url": remote_file.source_url,
                        "size_bytes": remote_file.size_bytes,
                    }
                )
                skipped_files.append(
                    {
                        "filename": display_filename,
                        "original_filename": remote_file.filename,
                        "source_url": remote_file.source_url,
                        "size_bytes": remote_file.size_bytes,
                        "message": message,
                        "duplicate_of": duplicate_match,
                    }
                )
                continue

            seen_file_hashes[file_hash] = display_filename
            processing_path, persisted_path, preserve_file = _prepare_processing_file(
                remote_file.temp_path,
                remote_file.filename,
            )
            _queue_processing_task(
                background_tasks,
                processing_path,
                display_filename,
                "File imported from link, starting OCR...",
                {
                    "_graph_access_token": graph_access_token,
                    "type": "remote_link",
                    "import_link": payload.link,
                    "source_url": remote_file.source_url,
                    "original_filename": remote_file.filename,
                    "tracking_filename": display_filename,
                    "size_bytes": remote_file.size_bytes,
                    "file_name": remote_file.filename,
                    "file_path": persisted_path,
                    "file_hash": file_hash,
                    "empID": _normalize_optional_text(payload.empID),
                    "empName": _normalize_optional_text(payload.empName),
                    "choose": normalized_choice,
                    "temporaryRetentionHours": retention_hours,
                    "uploaded_by": actor,
                    "uploaded_by_role": user_role.value,
                    "preserve_file": preserve_file,
                    "save_original_documents": preserve_file,
                    "generate_video_text": settings.generate_video_transcripts,
                },
            )
            queued_files.append(
                {
                    "filename": display_filename,
                    "original_filename": remote_file.filename,
                    "source_url": remote_file.source_url,
                    "size_bytes": remote_file.size_bytes,
                }
            )
        except Exception as exc:
            cleanup_path = persisted_path or processing_path
            if cleanup_path and os.path.exists(cleanup_path):
                os.remove(cleanup_path)
            logger.error(f"Error queueing linked file {remote_file.filename}: {exc}")
            processing_status[display_filename] = {
                "status": "error",
                "progress": 0,
                "message": str(exc)
            }
            save_status_to_db(display_filename, processing_status[display_filename])

    if not queued_files:
        raise HTTPException(status_code=400, detail="No supported link files could be queued for processing")

    return {
        "status": "processing",
        "imported_count": len(queued_files),
        "files": queued_files,
        "warnings": warnings,
        "message": (
            f"Queued {len(queued_files) - len(skipped_files)} new file(s) from link for processing."
            + (f" Skipped {len(skipped_files)} duplicate file(s)." if skipped_files else "")
        ),
        "duplicate_skipped_count": len(skipped_files),
        "skipped_files": skipped_files,
    }

# ============= GEM PROCESSING ENDPOINT =============
@router.post("/gem-process")
async def process_gem_pdf(file: UploadFile = File(...)):
    tmp_path = None
    try: 
        if Path(file.filename).suffix.lower() != ".pdf":
            raise HTTPException(status_code=400, detail="Only PDF files are supported.")

        # Save uploaded file
        output_dir = Path.cwd() / "output"
        output_dir.mkdir(parents=True, exist_ok=True)
        tmp_path = output_dir / file.filename

        file_bytes = await file.read()
        with open(tmp_path, "wb") as f:
            f.write(file_bytes)

        # Process and get back the output Excel path
        result = gem_processing(str(tmp_path))

        if result.get("status") == "error":
            raise HTTPException(status_code=500, detail=result.get("message", "Processing failed."))

        output_path = Path(result["output_path"])
        if not output_path.exists():
            raise HTTPException(status_code=500, detail="Output file was not created.")

        logger.info(f"[GEM] completed processing {file.filename} → {output_path}")

        return FileResponse(
            path=str(output_path),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=output_path.name,
            headers={
                "Content-Disposition": f"attachment; filename={output_path.name}",
                "Cache-Control": "no-cache, no-store, must-revalidate",
                "Pragma": "no-cache",
                "Expires": "0",
            }
        )
    except HTTPException:
        logger.error(f"GEM processing HTTP error: {traceback.format_exc()}")
        raise
    finally:
        if tmp_path and tmp_path.exists():
            os.remove(tmp_path)
# ============= OCR PROCESSING PIPELINE =============

def process_ocr_file(file_path: str, filename: str, source_info: Optional[dict] = None):
    """Complete Processing Pipeline:
    STEP 1: OCR Text Extraction (PaddleOCR)
    STEP 2: Structure with LLM (LLAMA)
    STEP 3: Save to MongoDB
    STEP 4: Update Status to Completed
    """
    db = None
    source_info = source_info if isinstance(source_info, dict) else {}
    graph_access_token = source_info.get("_graph_access_token")
    sanitized_source = _sanitize_source_info(source_info)
    requested_emp_id = _normalize_optional_text(source_info.get("empID"))
    requested_emp_name = _normalize_optional_text(source_info.get("empName"))
    requested_choice = _resolve_storage_mode(source_info.get("choose"))
    retention_hours = resolve_retention_hours(source_info.get("temporaryRetentionHours"))
    preserve_file = bool(source_info.get("preserve_file"))
    generate_video_text = _coerce_bool(
        source_info.get("generate_video_text"),
        settings.generate_video_transcripts,
    )
    persist_file_path = str(source_info.get("file_path") or file_path) if preserve_file else None
    document_saved = False

    logger.info(f"{'='*70}")
    logger.info(f"[PIPELINE] PROCESSING STARTED | file={filename} | source={source_info.get('source_url') or 'upload'} | preserve_file={preserve_file}")
    logger.info(f"{'='*70}")
    try:
        # Initialize database connection first
        db = MongoDBManager()
        all_text = ""
        confidences = []
        link_targets = []
        structured_content = None
        linked_documents = []
        linked_document_warnings = []
        combined_full_text = ""
        excel_upload = is_excel_file(filename)
        
        # STEP 1: FILE CONTENT EXTRACTION
        extraction_label = "WORKBOOK EXTRACTION" if excel_upload else "OCR TEXT EXTRACTION (PaddleOCR)"
        extraction_message = "Extracting workbook data..." if excel_upload else "Extracting text with OCR..."
        logger.info(f"[PIPELINE] STEP 1: {extraction_label} | file={filename}")
        logger.info(f"{'─'*70}")
        processing_status[filename] = {
            "status": "processing",
            "progress": 20,
            "message": extraction_message
        }
        save_status_to_db(filename, processing_status[filename])
        
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            processing_status[filename] = {
                "status": "error",
                "progress": 0,
                "message": "File not found"
            }
            save_status_to_db(filename, processing_status[filename])
            return
        
        try:
            all_text, confidences, link_targets, structured_content = _extract_text_from_document(
                file_path,
                filename,
                status_filename=filename,
                generate_video_text=generate_video_text,
            )
        except Exception as ocr_err:
            logger.exception(f"[PIPELINE] STEP 1 ERROR: extraction failed for {filename}")
            processing_status[filename] = {
                "status": "error",
                "progress": 20,
                "message": f"Extraction failed: {str(ocr_err)}"
            }
            save_status_to_db(filename, processing_status[filename])
            return
        
        if not all_text.strip():
            ocr_diag = get_last_ocr_diagnostics()
            if ocr_diag:
                logger.error(
                    f"[PIPELINE] STEP 1 ERROR: No text extracted from {filename} | "
                    f"reason={ocr_diag.get('status')} detail={ocr_diag.get('detail')} "
                    f"raw_candidates={ocr_diag.get('raw_candidates')} "
                    f"kept_lines={ocr_diag.get('kept_lines')} "
                    f"min_confidence={ocr_diag.get('min_confidence')}"
                )
            else:
                logger.error(f"[PIPELINE] STEP 1 ERROR: No text extracted from {filename}")
            if db:
                db.close()
            if os.path.exists(file_path):
                os.remove(file_path)
            processing_status[filename] = {
                "status": "error",
                "progress": 20,
                "message": "No text extracted from file"
            }
            save_status_to_db(filename, processing_status[filename])
            return
        
        ocr_confidence = round(sum(confidences) / len(confidences), 2) if confidences else 0.0
        logger.info(f"[PIPELINE] STEP 1 SUCCESS: extracted {len(all_text)} chars, confidence={ocr_confidence}% for {filename}")

        if link_targets and _LINKED_PDF_MAX_DOCUMENTS > 0:
            logger.info(f"Linked PDF targets discovered: {len(link_targets)}")
            _update_processing_progress(
                filename,
                48,
                f"Extracting linked PDFs... ({len(link_targets)} link(s) found)",
            )
            linked_documents, linked_document_warnings, _ = _collect_linked_document_texts(
                link_targets,
                graph_access_token=graph_access_token,
            )
            logger.info(
                f"Linked PDF extraction complete: {len(linked_documents)} document(s), "
                f"{len(linked_document_warnings)} warning(s)"
            )
        elif link_targets:
            linked_document_warnings.append(
                "Linked PDF extraction is disabled because PDF_LINK_MAX_DOCUMENTS is set to 0."
            )

        combined_full_text = _build_llm_input_text(all_text, linked_documents, max_chars=0)
        
        # STEP 2: STRUCTURING / METADATA
        is_structured_excel = bool(structured_content and structured_content.get("source_type") == "excel")
        logger.info(
            f"[PIPELINE] STEP 2: {'DETERMINISTIC EXCEL STRUCTURING' if is_structured_excel else 'STRUCTURE WITH LLM'} | file={filename}"
        )
        logger.info(f"{'─'*70}")
        processing_status[filename] = {
            "status": "processing",
            "progress": 50,
            "message": "Structuring workbook data..." if is_structured_excel else "Structuring data with AI..."
        }
        save_status_to_db(filename, processing_status[filename])
        if is_structured_excel:
            logger.info("[PIPELINE] STEP 2 ACTION: using deterministic Excel parser output")
            metadata = {
                "document_type": str(structured_content.get("document_type") or "excel_workbook"),
                "high_level_metadata": dict(structured_content.get("high_level_metadata") or {}),
                "items": [],
            }
        else:
            logger.info(f"[PIPELINE] STEP 2 ACTION: running Gemini for metadata extraction")
            try:
                max_llm_chars = settings.get_int("LLM_MAX_CHARS", 120000)
                llm_text = _build_llm_input_text(all_text, linked_documents, max_llm_chars)
                if max_llm_chars > 0 and len(combined_full_text) > len(llm_text):
                    logger.info(
                        f"Large document detected. Using budgeted LLM input ({len(llm_text)}/{len(combined_full_text)} chars)."
                    )
                metadata = extract_metadata_with_gemini(llm_text, ocr_confidence)
                logger.info("LLM complete!")
                logger.info(f"     • Type: {metadata.get('document_type', 'unknown').upper()}")
                logger.info(f"     • Fields: {len(metadata.get('high_level_metadata', {}))}")
            except Exception as llm_err:
                logger.exception(f"[PIPELINE] STEP 2 ERROR: Gemini LLM structuring failed for {filename}")
                metadata = {
                    "document_type": "unknown",
                    "high_level_metadata": {}
                }

        structured_metadata = dict(metadata.get("high_level_metadata") or {})
        structured_items = metadata.get("items") if isinstance(metadata.get("items"), list) else []
        if structured_items:
            structured_metadata["items"] = structured_items

        linked_document_entries, linked_metadata_warnings = _build_structured_linked_document_entries(linked_documents)
        if linked_document_entries:
            structured_metadata["referenced_documents"] = linked_document_entries
            structured_metadata["linked_documents_count"] = len(linked_document_entries)
            structured_metadata["linked_document_titles"] = [
                entry.get("title")
                for entry in linked_document_entries
                if entry.get("title")
            ]
        if linked_document_warnings or linked_metadata_warnings:
            structured_metadata["linked_document_warnings"] = (
                linked_document_warnings + linked_metadata_warnings
            )

        # Optional: Build embeddings for RAG retrieval
        embedding = None
        embedding_model = None
        if settings.enable_embeddings:
            try:
                from services.rag_embedder import RagEmbedder, build_embedding_text
                if RagEmbedder.available():
                    embedding_text = build_embedding_text(
                        metadata.get("document_type"),
                        structured_metadata,
                        combined_full_text,
                        max_chars=settings.embedding_max_chars,
                    )
                    embedding = RagEmbedder.embed_text(embedding_text)
                    embedding_model = settings.get_string(
                        "EMBEDDING_MODEL", "sentence-transformers/all-MiniLM-L6-v2"
                    )
                    logger.info("[RAG] Embedding generated")
            except Exception as emb_err:
                logger.warning(f"[RAG] Embedding generation failed: {emb_err}")
        
        # STEP 3: DATABASE SAVE
        logger.info(f"STEP 3: SAVE TO MONGODB")
        logger.info(f"{'─'*70}")
        processing_status[filename] = {
            "status": "processing",
            "progress": 80,
            "message": "Saving to database..."
        }
        save_status_to_db(filename, processing_status[filename])
        try:
            logger.info(f"[PIPELINE] STEP 3: saving document to MongoDB | file={filename}")
            original_filename = (
                sanitized_source.get("original_filename")
                if sanitized_source.get("original_filename")
                else filename
            )
            assignment = _resolve_employee_assignment(
                db=db,
                requested_emp_id=requested_emp_id,
                requested_emp_name=requested_emp_name,
                requested_choice=requested_choice,
                retention_hours=retention_hours,
                structured_metadata=structured_metadata,
                owner_email=str(source_info.get("uploaded_by") or "").strip() or None,
                owner_role=validate_role_string(source_info.get("uploaded_by_role") or "employee")
                if str(source_info.get("uploaded_by_role") or "").strip()
                else None,
            )
            source_payload = dict(sanitized_source or {"type": "upload", "original_filename": filename})
            source_payload["uploaded_by_normalized"] = _normalize_kbac_owner(source_payload.get("uploaded_by"))
            source_payload["linked_documents_count"] = len(linked_documents)
            source_payload["selected_empID"] = requested_emp_id or None
            source_payload["selected_empName"] = requested_emp_name or None
            source_payload["requested_storage_mode"] = requested_choice
            source_payload["temporaryRetentionHours"] = retention_hours
            source_payload["generate_video_text"] = generate_video_text
            source_payload["save_original_documents"] = preserve_file
            kbac_scope, kbac_owner = _resolve_document_kbac_scope(assignment, source_payload)
            source_payload["kbac_scope"] = kbac_scope
            source_payload["kbac_owner"] = kbac_owner
            if linked_documents:
                source_payload["linked_documents"] = _prepare_linked_document_metadata(linked_documents)
            if linked_document_warnings:
                source_payload["linked_document_warnings"] = linked_document_warnings
            content_payload = {
                "full_text": combined_full_text,
                "main_document_text": (
                    str(structured_content.get("preview_text") or all_text)
                    if is_structured_excel
                    else all_text
                ),
            }
            if is_structured_excel:
                content_payload["structured_workbook"] = structured_content.get("workbook")
                content_payload["retrieval_text"] = structured_content.get("retrieval_text")
                content_payload["preview_text"] = structured_content.get("preview_text")
            final_output = {
                "document_type": metadata.get("document_type", "unknown"),
                "high_level_metadata": structured_metadata,
                "file_name": original_filename,
                "file_path": persist_file_path,
                "tracking_filename": filename,
                "source": source_payload,
                "empID": assignment.get("empID"),
                "empName": assignment.get("empName"),
                "employee_uuid": assignment.get("employee_uuid"),
                "document_uuid": str(uuid4()),
                "storage_mode": assignment.get("storage_mode"),
                "isTemporary": assignment.get("isTemporary", False),
                "expiry_at": assignment.get("expiry_at"),
                "assignment_status": assignment.get("assignment_status"),
                "kbac_scope": kbac_scope,
                "kbac_owner": kbac_owner,
                "suggested_profile": assignment.get("suggested_profile"),
                "confidence": {
                    "ocr_engine": "StructuredExcelParser" if is_structured_excel else "PaddleOCR",
                    "confidence_percent": ocr_confidence
                },
                "content": content_payload,
                "embedding": embedding,
                "embedding_model": embedding_model,
                "saved_at": datetime.utcnow()
            }
            if linked_documents:
                final_output["content"]["linked_documents_text"] = _build_linked_documents_text(linked_documents)
            duplicate_match = db.find_duplicate_for_document(
                final_output,
                duplicate_owner=_normalize_kbac_owner(source_payload.get("uploaded_by")),
            )
            if duplicate_match:
                message = _build_duplicate_skip_message(original_filename, duplicate_match, duplicate_match.get("reason"))
                logger.info(f"Duplicate detected for {filename}; skipping save. Matched existing document {duplicate_match.get('document_id')}")
                processing_status[filename] = {
                    "status": "duplicate_skipped",
                    "progress": 100,
                    "message": message,
                    "duplicate_of": duplicate_match,
                    "duplicate_reason": duplicate_match.get("reason"),
                    "document_id": duplicate_match.get("document_id"),
                    "document_uuid": duplicate_match.get("document_uuid"),
                    "storage_mode": final_output.get("storage_mode"),
                    "empID": final_output.get("empID"),
                    "empName": final_output.get("empName"),
                    "employee_uuid": final_output.get("employee_uuid"),
                    "isTemporary": final_output.get("isTemporary"),
                    "expiry_at": final_output.get("expiry_at"),
                    "assignment_status": final_output.get("assignment_status"),
                }
                save_status_to_db(filename, processing_status[filename])
                return
            logger.info(f"Saving to MongoDB...")
            doc_id = db.save_document(final_output)
            document_saved = True
            _invalidate_qa_caches()
            logger.info(f"[PIPELINE] STEP 3 SUCCESS: saved document id={doc_id} for {filename}")
            # STEP 4: Mark as completed
            logger.info(f"[PIPELINE] STEP 4: COMPLETED | file={filename}")
            logger.info(f"{'─'*70}")
            processing_status[filename] = {
                "status": "completed",
                "progress": 100,
                "message": (
                    "Document parsed successfully. Employee profile can be added later."
                    if assignment.get("employee_action_required")
                    else f"Successfully processed! Document ID: {doc_id}"
                ),
                "document_id": doc_id,
                "document_uuid": final_output.get("document_uuid"),
                "storage_mode": final_output.get("storage_mode"),
                "empID": final_output.get("empID"),
                "empName": final_output.get("empName"),
                "employee_uuid": final_output.get("employee_uuid"),
                "isTemporary": final_output.get("isTemporary"),
                "expiry_at": final_output.get("expiry_at"),
                "assignment_status": final_output.get("assignment_status"),
                "employee_action_required": assignment.get("employee_action_required", False),
                "suggested_profile": assignment.get("suggested_profile"),
            }
            save_status_to_db(filename, processing_status[filename])
            logger.info(f"PROCESSING COMPLETED!")
            logger.info(f"{'='*70}")
        except Exception as db_err:
            logger.exception(f"[PIPELINE] STEP 3 ERROR: Database save failed for {filename}")
            processing_status[filename] = {
                "status": "error",
                "progress": 80,
                "message": f"Database error: {str(db_err)}"
            }
            save_status_to_db(filename, processing_status[filename])
            return
        
    except Exception as e:
        logger.exception(f"[PIPELINE] CRITICAL ERROR: file={filename}")
        logger.error(f"{'='*70}")
        processing_status[filename] = {
            "status": "error",
            "progress": 0,
            "message": f"Critical error: {str(e)}"
        }
        save_status_to_db(filename, processing_status[filename])
    
    finally:
        if db:
            db.close()
        should_delete_source = not preserve_file
        if preserve_file and not document_saved and os.path.exists(file_path):
            os.remove(file_path)
        if should_delete_source and os.path.exists(file_path):
            os.remove(file_path)
# ===========================
# PUBLIC TRANSCRIPT ROUTES
# Put this in same routes file
# BELOW: router = APIRouter(prefix="/api"...)
# ===========================
templates = Jinja2Templates(directory="templates")


# -------------------------------------------------
# 1. PUBLIC PAGE ROUTE (NO LOGIN REQUIRED)
# This bypasses auth because no _require_permission()
# -------------------------------------------------


@router.get("/public-transcript", response_class=HTMLResponse, include_in_schema=False)
async def public_transcript_page(request: Request):
    return templates.TemplateResponse(
        "public_transcript.html",
        {
            "request": request,
            "base_path": str(request.scope.get("root_path") or "").rstrip("/"),
        }
    )


# -------------------------------------------------
# 2. PUBLIC UPLOAD + TRANSCRIPT ROUTE
# NO LOGIN REQUIRED
# -------------------------------------------------
@router.post("/generate-transcript", include_in_schema=False)
async def generate_transcript(
    request: Request,
    file: UploadFile = File(...)
):
    try:
        # create temp folder
        upload_dir = Path("output/public_transcripts")
        upload_dir.mkdir(parents=True, exist_ok=True)

        file_path = upload_dir / file.filename

        # save uploaded video
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # -----------------------------------------
        # PLACE YOUR WHISPER / AI MODEL HERE
        file_ext = file.filename.rsplit('.', 1)[-1].lower() if '.' in file.filename else ''
        ALLOWED_EXTENSIONS = {'mp4', 'avi', 'mov', 'mkv', 'flv', 'wmv', "pdf", "img", "jpg", "jpeg", "png"}
        if file_ext not in ALLOWED_EXTENSIONS:
            logger.error(f"File type not allowed: {file_ext}")
            raise HTTPException(
            status_code=400,
            detail=f"File type not allowed. Allowed: {', '.join(ALLOWED_EXTENSIONS)}"
            )
        if file_ext in {'mp4', 'avi', 'mov', 'mkv', 'flv', 'wmv'}:    
            try:
                import whisper
                model = whisper.load_model("small")
                if model:
                    logger.info(f"Whisper model loaded successfully for file {file.filename}")
                result = model.transcribe(str(file_path))
            except Exception as e:
                logger.error(f"Error transcribing video {file_path}: {str(e)}")
                raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")
            if result:
                transcript_text = result["text"]
            else:
                transcript_text = "No transcript available."   
            return JSONResponse({
                "status": "success",
                "filename": file.filename,
                "transcript": transcript_text
            })
        elif file_ext in {"pdf"}:
            import fitz
            text = ""
            doc = fitz.open(file_path)
            for page in doc:
                try:
                    page_text = page.get_text().strip()

                    if page_text:
                        text += page_text + "\n"

                except Exception as e:
                    logger.error(
                        f"Error processing page {page.number + 1} of PDF {file_path}: {str(e)}"
                    )

            doc.close()

            return JSONResponse({
                "status": "success",
                "filename": file.filename,
                "transcript": text.strip() if text.strip() else "No transcript available."
            })
        elif file_ext in {"img", "jpg", "jpeg", "png"}:
            try:
                    from PIL import Image
                    img = Image.open(file_path).convert("RGB")
                    img = np.array(img, dtype=np.uint8)
                    result = paddle_ocr_extract(img)
                    text = result.get("raw_text", "") if result else ""
                    logger.info(f"OCR extracted text from image {file.filename}: {text[:100]}...") 

                    if not text:
                        text = "No transcript available."

                    return JSONResponse({
                        "status": "success",
                        "filename": file.filename,
                        "transcript": text
                    })
            except Exception as e:
                logger.error(str(e))
                raise HTTPException(status_code=400, detail="Failed to process image file.")

    except Exception as e:
        logger.error(f"Failed to process image file {file_path}: {str(e)}")
        raise HTTPException(
        status_code=400,
        detail="Failed to process image file."
        )
                    


    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

